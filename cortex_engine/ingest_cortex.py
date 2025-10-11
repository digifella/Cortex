# ## File: ingest_cortex.py
# Version: 14.0.0 (Docling Integration with Migration Strategy)
# Date: 2025-08-22
# Purpose: Core ingestion script for Project Cortex with integrated knowledge graph extraction.
#          - FEATURE (v13.0.0): Integrated entity extraction and knowledge graph building
#            during the ingestion process. The system now extracts people, organizations,
#            projects, and their relationships while maintaining backward compatibility.

import os
import sys
os.environ['ANONYMIZED_TELEMETRY'] = 'False'

import argparse
import json
import logging
import warnings

# Suppress common warnings that don't affect functionality
warnings.filterwarnings("ignore", message=".*attention_mask.*")
warnings.filterwarnings("ignore", message=".*pad_token_id.*")
warnings.filterwarnings("ignore", category=FutureWarning, module="transformers")

# Suppress Docling warnings about API compatibility (non-blocking)
warnings.filterwarnings("ignore", message=".*Docling.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*docling.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*docling_reader.*", category=UserWarning)

# Suppress Pydantic field name warnings (non-functional)
warnings.filterwarnings("ignore", message=".*field names.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*pydantic.*field.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*has conflict with protected namespace.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*model_.*", category=UserWarning)

# Suppress torchvision and model loading warnings (informational only)
warnings.filterwarnings("ignore", message=".*torchvision.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*CUDA.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*GPU.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*torch.load.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*vulnerability.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*CVE-.*", category=UserWarning)

# Suppress sentence-transformers verbosity
warnings.filterwarnings("ignore", category=FutureWarning, module="sentence_transformers")
warnings.filterwarnings("ignore", category=UserWarning, module="sentence_transformers")

# Suppress huggingface model warnings
warnings.filterwarnings("ignore", category=FutureWarning, module="huggingface_hub")

# Reduce verbosity of specific loggers that tend to be noisy
logging.getLogger("unstructured").setLevel(logging.ERROR)
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("sentence_transformers").setLevel(logging.WARNING)
logging.getLogger("huggingface_hub").setLevel(logging.WARNING)
logging.getLogger("torch").setLevel(logging.WARNING)
logging.getLogger("torchvision").setLevel(logging.ERROR)
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Literal, Tuple
import re

from pydantic import BaseModel, Field, ValidationError
from llama_index.core import Document
from llama_index.core.settings import Settings
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.core import VectorStoreIndex, StorageContext
import chromadb
from chromadb.config import Settings as ChromaSettings
import hashlib

from llama_index.readers.file import (
    DocxReader,
    PptxReader,
    PyMuPDFReader,
    FlatReader,
    UnstructuredReader
)

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from cortex_engine.config import INGESTION_LOG_PATH, STAGING_INGESTION_FILE, INGESTED_FILES_LOG, COLLECTION_NAME, EMBED_MODEL
from cortex_engine.utils.file_utils import get_file_hash
from cortex_engine.utils.logging_utils import get_logger
from cortex_engine.utils.smart_ollama_llm import create_smart_ollama_llm
try:
    from cortex_engine.migration_to_docling import create_migration_manager
except Exception:
    def create_migration_manager():
        class _Noop:
            def migrate(self, *args, **kwargs):
                return None
        return _Noop()

logger = get_logger(__name__)
from cortex_engine.query_cortex import describe_image_with_vlm_for_ingestion, describe_image_with_vlm_async
from cortex_engine.embedding_adapters import EmbeddingServiceAdapter
from cortex_engine.entity_extractor import EntityExtractor, ExtractedEntity, ExtractedRelationship
from cortex_engine.graph_manager import EnhancedGraphManager
from cortex_engine.batch_manager import BatchState

LOG_FILE = INGESTION_LOG_PATH
# STAGING_FILE will be set dynamically based on runtime db_path
COLLECTIONS_FILE = str(project_root / "working_collections.json")

def get_staging_file_path(db_path: str) -> str:
    """Get the staging file path for the given database path"""
    return os.path.join(db_path, "staging_ingestion.json")

# Define supported image extensions - Enhanced Visual Processing Support
IMAGE_EXTENSIONS = {
    # Standard formats
    ".png", ".jpg", ".jpeg",
    # Additional formats for comprehensive visual processing
    ".gif", ".bmp", ".webp", ".tiff", ".tif",
    # Vector and specialized formats
    ".svg",  # Will be converted to raster for VLM processing
    ".ico"   # Icon files
}

os.makedirs(os.path.dirname(LOG_FILE), exist_ok=True)

def initialize_script():
    """Sets up robust logging and configures LlamaIndex models."""
    log_formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    root_logger = logging.getLogger()

    if root_logger.hasHandlers():
        root_logger.handlers.clear()

    root_logger.setLevel(logging.INFO)

    file_handler = logging.FileHandler(LOG_FILE, mode='a', encoding='utf-8')
    file_handler.setFormatter(log_formatter)
    root_logger.addHandler(file_handler)

    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setFormatter(log_formatter)
    root_logger.addHandler(console_handler)

    logging.info("--- Script Initialized: Configuring models... ---")
    
    # Check if Ollama is available before configuring
    from cortex_engine.utils.ollama_utils import check_ollama_service, format_ollama_error_for_user
    
    chk = check_ollama_service()
    # Backward compatible unpack
    is_running = chk[0]
    error_msg = chk[1] if len(chk) > 1 else None
    resolved_url = chk[2] if len(chk) > 2 else None
    if not is_running:
        logging.warning(f"Ollama service not available: {error_msg}")
        logging.warning("AI-enhanced metadata extraction will be disabled. Documents will be processed with basic metadata only.")
        Settings.llm = None  # Will be handled in analysis function
    else:
        # Set reasonable timeout for metadata extraction (3 minutes)
        # Very long documents may timeout and use fallback metadata
        Settings.llm = create_smart_ollama_llm(model="mistral:latest", request_timeout=180.0)
        logging.info(f"Ollama connected successfully at {resolved_url}")
    
    # Unify embeddings with search/async via adapter around embedding_service (Docker)
    try:
        Settings.embed_model = EmbeddingServiceAdapter(model_name=EMBED_MODEL)
        logging.info(f"Models configured via EmbeddingServiceAdapter (Embed: {EMBED_MODEL}).")
    except Exception as e:
        logging.warning(f"EmbeddingServiceAdapter failed ({e}); falling back to HuggingFaceEmbedding")
        Settings.embed_model = HuggingFaceEmbedding(model_name=EMBED_MODEL)
        logging.info(f"Models configured (Embed: {EMBED_MODEL}).")

class RichMetadata(BaseModel):
    document_type: Literal[
        "Project Plan", "Technical Documentation", "Proposal/Quote", "Case Study / Trophy",
        "Final Report", "Draft Report", "Presentation", "Contract/SOW",
        "Meeting Minutes", "Financial Report", "Research Paper", "Email Correspondence", "Image/Diagram", "Other"
    ] = Field(..., description="The primary category of the document.")
    proposal_outcome: Literal["Won", "Lost", "Pending", "N/A"] = Field(..., description="The outcome of the proposal, if applicable.")
    summary: str = Field(..., description="A concise, 1-3 sentence summary of the document's content and purpose.")
    thematic_tags: List[str] = Field(default_factory=list, description="A list of 3-5 key themes, topics, or technologies discussed.")

class DocumentMetadata(BaseModel):
    doc_id: str
    doc_posix_path: str
    file_name: str
    last_modified_date: str
    rich_metadata: Optional[RichMetadata] = None
    exclude_from_final: bool = False
    extracted_entities: List[Dict] = Field(default_factory=list)
    extracted_relationships: List[Dict] = Field(default_factory=list)
    target_collection: Optional[str] = None  # Collection assignment for finalization

def load_processed_files_log(log_path: str) -> Dict[str, str]:
    if not os.path.exists(log_path): return {}
    with open(log_path, 'r') as f:
        try:
            log_data = json.load(f)
            return {k: (v[0] if isinstance(v, list) else v) for k, v in log_data.items()}
        except json.JSONDecodeError: return {}

def write_to_processed_log(log_path: str, file_path: str, doc_id: str):
    """Write a processed file entry to the log."""
    processed_files = {}
    if os.path.exists(log_path):
        try:
            with open(log_path, 'r') as f:
                processed_files = json.load(f)
        except (json.JSONDecodeError, IOError):
            processed_files = {}
    
    processed_files[file_path] = doc_id
    
    try:
        with open(log_path, 'w') as f:
            json.dump(processed_files, f, indent=2)
    except IOError as e:
        logging.error(f"Failed to write to processed log: {e}")

def get_document_content(file_path: str, skip_image_processing: bool = False) -> Dict:
    """Extract content from a single document file
    
    Returns:
        Dict with keys: 'content' (str), 'metadata' (dict), 'source_type' (str)
    """
    from cortex_engine.query_cortex import describe_image_with_vlm_for_ingestion
    
    path = Path(file_path)
    extension = path.suffix.lower()
    
    # Set up readers
    reader_map = {".pdf": PyMuPDFReader(), ".docx": DocxReader(), ".pptx": PptxReader(), 
                  ".doc": UnstructuredReader(), ".ppt": UnstructuredReader()}
    default_reader = FlatReader()
    
    try:
        # Handle image files using the VLM (if enabled)
        if extension in IMAGE_EXTENSIONS:
            if skip_image_processing:
                return {
                    'content': f"Image file: {path.name} (processing skipped)",
                    'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                    'source_type': 'image_skipped'
                }
            else:
                try:
                    description = describe_image_with_vlm_for_ingestion(file_path)
                    return {
                        'content': description,
                        'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                        'source_type': 'image'
                    }
                except Exception as e:
                    return {
                        'content': f"Image file: {path.name} (VLM processing failed: {str(e)})",
                        'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                        'source_type': 'image_error'
                    }

        # Text-based document handling
        reader = reader_map.get(extension, default_reader)
        
        # Handle different reader types with error handling
        if isinstance(reader, PyMuPDFReader):
            try:
                docs_from_file = reader.load_data(path)
            except Exception as pdf_error:
                if "cmsOpenProfileFromMem" not in str(pdf_error):
                    logger.warning(f"PDF processing warning for {path.name}: {pdf_error}")
                docs_from_file = reader.load_data(path)
        elif isinstance(reader, UnstructuredReader):
            try:
                docs_from_file = reader.load_data(path)
            except Exception as unstructured_error:
                error_str = str(unstructured_error).lower()
                if not any(x in error_str for x in ['warning', 'deprecation', 'future']):
                    logger.warning(f"UnstructuredReader processing issue for {path.name}: {unstructured_error}")
                docs_from_file = [Document(
                    text=f"Error processing old Office document: {path.name}. Reason: {unstructured_error}",
                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                )]
        else:
            try:
                docs_from_file = reader.load_data(file=path)
            except UnicodeDecodeError as encoding_error:
                logger.warning(f"Encoding error for {path.name}: {encoding_error}")
                docs_from_file = [Document(
                    text=f"File could not be processed due to encoding issues: {path.name}",
                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                )]
            except OSError as wmf_error:
                if "cannot find loader for this WMF file" in str(wmf_error):
                    logger.warning(f"WMF image error in {path.name}: {wmf_error}")
                    docs_from_file = [Document(
                        text=f"Document processing failed due to WMF image issues: {path.name}",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
                else:
                    raise wmf_error
        
        # Combine text from multiple documents if needed
        if docs_from_file:
            combined_text = "\n\n".join([doc.text for doc in docs_from_file if doc.text])
            metadata = docs_from_file[0].metadata if docs_from_file else {}
            metadata.update({'file_path': str(path.as_posix()), 'file_name': path.name})
            
            return {
                'content': combined_text,
                'metadata': metadata,
                'source_type': metadata.get('source_type', 'document')
            }
        else:
            return {
                'content': f"No content extracted from {path.name}",
                'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                'source_type': 'document_error'
            }
            
    except Exception as e:
        logger.error(f"Failed to process {file_path}: {e}")
        return {
            'content': f"Error processing {path.name}: {str(e)}",
            'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
            'source_type': 'document_error'
        }


def create_document_with_entities(content_result: Dict, rich_metadata: Optional[RichMetadata] = None) -> Document:
    """Create a Document object with entity extraction"""
    content = content_result.get('content', '')
    metadata = content_result.get('metadata', {})
    source_type = content_result.get('source_type', 'document')
    
    # Create the document
    doc = Document(text=content)
    doc.metadata.update(metadata)
    doc.metadata['source_type'] = source_type
    
    # Add rich metadata if provided
    if rich_metadata:
        doc.metadata['rich_metadata'] = rich_metadata
    
    return doc


def manual_load_documents(file_paths: List[str], args=None) -> List[Document]:
    """
    Enhanced document loading with Docling integration and intelligent migration.
    
    This function now uses a migration strategy to gradually transition from legacy
    LlamaIndex readers to Docling-enhanced processing while maintaining stability.
    """
    
    # Check for migration mode preference in args
    # In Docker environments, default to legacy mode to avoid Docling dependency conflicts
    import os
    default_mode = 'legacy' if os.path.exists('/.dockerenv') else 'gradual'
    migration_mode = getattr(args, 'migration_mode', default_mode) if args else default_mode
    skip_image_processing = getattr(args, 'skip_image_processing', False) if args else False
    
    # Use migration manager for intelligent processing
    try:
        migration_manager = create_migration_manager(mode=migration_mode)
        documents = migration_manager.process_documents(file_paths, skip_image_processing)
        
        # Log migration insights
        if hasattr(migration_manager, 'comparison_results') and migration_manager.comparison_results:
            logger.info("📊 Migration manager provided processing insights")
        
        return documents
        
    except Exception as e:
        logger.warning(f"⚠️ Migration processing failed, falling back to legacy: {e}")
        # Fallback to legacy processing
        return _legacy_manual_load_documents(file_paths, args)


def _process_images_batch(
    image_files: List[str],
    skip_image_processing: bool = False
) -> List[Document]:
    """
    Process multiple images in parallel with VLM.

    This function significantly improves image processing performance by:
    1. Processing up to 3 images concurrently
    2. Using 30s timeout per image (vs 120s previously)
    3. Graceful fallback on timeout/error

    Args:
        image_files: List of image file paths to process
        skip_image_processing: Skip VLM processing if True (fast mode)

    Returns:
        List of Document objects with image descriptions or placeholders
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from pathlib import Path

    documents = []

    if skip_image_processing:
        logging.info(f"⚡ Skipping VLM processing for {len(image_files)} images (fast mode)")
        for file_path in image_files:
            path = Path(file_path)
            doc = Document(text=f"Image file: {path.name} (processing skipped)")
            doc.metadata['file_path'] = str(path.as_posix())
            doc.metadata['file_name'] = path.name
            doc.metadata['source_type'] = 'image_skipped'
            documents.append(doc)
        return documents

    logging.info(f"🖼️ Processing {len(image_files)} images with VLM (parallel, 30s timeout each)")

    with ThreadPoolExecutor(max_workers=3) as executor:
        # Submit all image processing tasks
        future_to_path = {
            executor.submit(describe_image_with_vlm_async, img_path, 30): img_path
            for img_path in image_files
        }

        # Process results as they complete
        for idx, future in enumerate(as_completed(future_to_path), 1):
            img_path = future_to_path[future]
            path = Path(img_path)

            logging.info(f"📄 Processing image {idx}/{len(image_files)}: {path.name}")

            try:
                description = future.result()

                if description:
                    doc = Document(text=description)
                    doc.metadata['file_path'] = str(path.as_posix())
                    doc.metadata['file_name'] = path.name
                    doc.metadata['source_type'] = 'image'
                    documents.append(doc)
                    logging.info(f"✅ Successfully processed image: {path.name}")
                else:
                    # Timeout or error - create placeholder
                    doc = Document(text=f"Image file: {path.name} (VLM timeout/error)")
                    doc.metadata['file_path'] = str(path.as_posix())
                    doc.metadata['file_name'] = path.name
                    doc.metadata['source_type'] = 'image_fallback'
                    documents.append(doc)
                    logging.warning(f"⚠️ Image processing failed, using fallback: {path.name}")

            except Exception as e:
                logging.error(f"❌ Unexpected error processing {path.name}: {e}")
                # Create error placeholder
                doc = Document(text=f"Image file: {path.name} (processing error)")
                doc.metadata['file_path'] = str(path.as_posix())
                doc.metadata['file_name'] = path.name
                doc.metadata['source_type'] = 'image_error'
                documents.append(doc)

    successful = len([d for d in documents if d.metadata.get('source_type') == 'image'])
    logging.info(f"🎯 Image processing complete: {successful}/{len(image_files)} successful")

    return documents


def _legacy_manual_load_documents(file_paths: List[str], args=None) -> List[Document]:
    """
    Legacy document loading function (original implementation).
    Kept as fallback for compatibility and robustness.
    """
    documents = []
    
    # Initialize readers with graceful fallback for torch version issues
    reader_map = {".pdf": PyMuPDFReader(), ".docx": DocxReader()}
    
    # Try to initialize PptxReader, fallback to UnstructuredReader if torch issues
    try:
        reader_map[".pptx"] = PptxReader()
        reader_map[".ppt"] = PptxReader()
        logging.info("✅ PowerPoint reader initialized successfully")
    except Exception as e:
        logging.warning(f"⚠️ PowerPoint reader failed to initialize (torch version issue): {e}")
        logging.info("📋 Using UnstructuredReader fallback for PowerPoint files")
        reader_map[".pptx"] = UnstructuredReader()
        reader_map[".ppt"] = UnstructuredReader()
    
    # Add other readers
    reader_map[".doc"] = UnstructuredReader()
    default_reader = FlatReader()
    
    # Suppress PyMuPDF warnings
    import fitz
    fitz.TOOLS.mupdf_display_errors(False)
    
    # Expand directories to get actual files
    expanded_file_paths = []
    for file_path in file_paths:
        path = Path(file_path)
        if path.is_dir():
            # Recursively find all files in directory
            for file in path.rglob('*'):
                if file.is_file():
                    expanded_file_paths.append(str(file))
        elif path.is_file():
            expanded_file_paths.append(str(path))
        else:
            logging.warning(f"Path does not exist or is not a file/directory: {file_path}")

    logging.info(f"Expanded {len(file_paths)} paths to {len(expanded_file_paths)} files")

    # Separate images from other documents for batch processing
    skip_images = getattr(args, 'skip_image_processing', False) if hasattr(args, 'skip_image_processing') else False
    image_files = []
    other_files = []

    for file_path in expanded_file_paths:
        path = Path(file_path)
        extension = path.suffix.lower()
        if extension in IMAGE_EXTENSIONS:
            image_files.append(file_path)
        else:
            other_files.append(file_path)

    logging.info(f"📊 File breakdown: {len(image_files)} images, {len(other_files)} documents")

    # Process images in batch (parallel with timeout)
    if image_files:
        image_docs = _process_images_batch(image_files, skip_images)
        documents.extend(image_docs)

    # Process other documents normally
    total_files = len(other_files)
    for idx, file_path in enumerate(other_files, 1):
        try:
            path = Path(file_path)
            extension = path.suffix.lower()
            logging.info(f"📄 Processing file {idx}/{total_files}: {path.name}")

            # Images are already processed in batch above, this loop only handles documents
            # Original text-based document handling
            reader = reader_map.get(extension, default_reader)
            logging.info(f"Loading '{file_path}' with reader: {reader.__class__.__name__}")
            
            # For PDFs, add extra error handling
            if isinstance(reader, PyMuPDFReader):
                try:
                    docs_from_file = reader.load_data(path)
                except Exception as pdf_error:
                    # If PDF fails completely, log it but continue
                    if "cmsOpenProfileFromMem" not in str(pdf_error):
                        logging.warning(f"PDF processing warning for {path.name}: {pdf_error}")
                    docs_from_file = reader.load_data(path)
            elif isinstance(reader, UnstructuredReader):
                # UnstructuredReader also uses positional path parameter
                try:
                    docs_from_file = reader.load_data(path)
                except Exception as unstructured_error:
                    # Many UnstructuredReader warnings are harmless parsing issues - only log real errors
                    error_str = str(unstructured_error).lower()
                    if not any(x in error_str for x in ['warning', 'deprecation', 'future']):
                        logging.warning(f"UnstructuredReader processing issue for {path.name}: {unstructured_error}")
                    # Create error document if UnstructuredReader fails
                    docs_from_file = [Document(
                        text=f"Error processing old Office document: {path.name}. Reason: {unstructured_error}",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
            else:
                try:
                    docs_from_file = reader.load_data(file=path)
                except UnicodeDecodeError as encoding_error:
                    logging.warning(f"Encoding error for {path.name} - file appears to be binary or uses unsupported encoding: {encoding_error}")
                    # Create an informational document noting the file couldn't be processed
                    docs_from_file = [Document(
                        text=f"File could not be processed due to encoding issues: {path.name}. This file may be binary or use unsupported text encoding.",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
                except OSError as wmf_error:
                    if "cannot find loader for this WMF file" in str(wmf_error):
                        logging.warning(f"WMF image error in {path.name}, attempting text-only extraction: {wmf_error}")
                        # For PowerPoint files with WMF image issues, try alternative extraction
                        if isinstance(reader, PptxReader):
                            try:
                                # Try to extract just text content without images
                                from pptx import Presentation
                                prs = Presentation(path)
                                text_content = []
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text") and shape.text.strip():
                                            text_content.append(shape.text.strip())
                                
                                if text_content:
                                    combined_text = "\n\n".join(text_content)
                                    docs_from_file = [Document(
                                        text=f"PowerPoint content (text-only extraction due to WMF image issues):\n\n{combined_text}",
                                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_partial'}
                                    )]
                                else:
                                    docs_from_file = [Document(
                                        text=f"PowerPoint file processed but no extractable text found: {path.name}",
                                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_partial'}
                                    )]
                            except Exception as pptx_error:
                                logging.error(f"Alternative PowerPoint extraction also failed for {path.name}: {pptx_error}")
                                docs_from_file = [Document(
                                    text=f"PowerPoint file could not be processed: {path.name}. WMF image and text extraction both failed.",
                                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                                )]
                        else:
                            # For non-PowerPoint files with WMF issues, create error document
                            docs_from_file = [Document(
                                text=f"Document processing failed due to WMF image issues: {path.name}",
                                metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                            )]
                    else:
                        raise wmf_error
                
            for doc in docs_from_file:
                doc.metadata['file_path'] = str(path.as_posix())
                doc.metadata['file_name'] = path.name
                doc.metadata['source_type'] = 'document'
            documents.extend(docs_from_file)
            
        except Exception as e:
            logging.error(f"Failed to load file {file_path}: {e}", exc_info=True)
            documents.append(Document(
                text=f"Error reading this document. Could not load content from {Path(file_path).name}. Reason: {e}",
                metadata={'file_path': str(Path(file_path).as_posix()), 'file_name': Path(file_path).name}
            ))
    return documents

def extract_entities_and_relationships(doc_text: str, metadata: Dict, entity_extractor: EntityExtractor) -> Tuple[List[Dict], List[Dict]]:
    """Extract entities and relationships from document and convert to serializable format."""
    try:
        entities, relationships = entity_extractor.extract_entities_and_relationships(doc_text, metadata)
        
        # Convert to dictionaries for JSON serialization
        entities_dict = [entity.model_dump() for entity in entities]
        relationships_dict = [rel.model_dump() for rel in relationships]
        
        return entities_dict, relationships_dict
    except Exception as e:
        logging.error(f"Failed to extract entities: {e}")
        return [], []

def analyze_documents(include_paths: List[str], fresh_start: bool, args=None, target_collection: str = None):
    logging.info(f"--- Starting Stage 2: Document Analysis with Graph Extraction (Cortex v13.0.0) ---")
    print("CORTEX_STAGE::ANALYSIS_START", flush=True)
    
    # Initialize batch manager if db_path is available
    batch_manager = None
    if hasattr(args, 'db_path') and args.db_path:
        batch_manager = BatchState(args.db_path)
        
        # Handle resume logic
        if not fresh_start:
            batch_id, files_to_process, completed_count = batch_manager.resume_or_create_batch(include_paths)
            if not files_to_process:
                logging.info("No files to process - all files already completed")
                return
            
            logging.info(f"Batch {batch_id}: Processing {len(files_to_process)} files ({completed_count} already completed)")
            include_paths = files_to_process
        else:
            # Fresh start - clear any existing batch
            batch_manager.clear_batch()
            batch_manager.create_batch(include_paths)
    
    # Get staging file path based on batch manager's db_path
    staging_file = get_staging_file_path(batch_manager.db_path) if batch_manager else STAGING_INGESTION_FILE
    if fresh_start and os.path.exists(staging_file): os.remove(staging_file)
    
    if target_collection:
        logging.info(f"Collection assignment configured: documents will be assigned to '{target_collection}'")
    
    # Initialize entity extractor
    entity_extractor = EntityExtractor()
    
    # Check memory and implement chunked processing for large batches
    if len(include_paths) > 1000:
        logging.warning(f"Large batch detected ({len(include_paths)} files). Consider processing in smaller chunks to avoid memory issues.")
    
    docs = manual_load_documents(include_paths, args)
    logging.info(f"Loaded {len(docs)} document objects from {len(include_paths)} files.")
    unique_docs = list({doc.metadata['file_path']: doc for doc in docs}.values())
    logging.info(f"--- Found {len(unique_docs)} unique files to process. ---")
    
    # Memory check
    try:
        import psutil
        memory_percent = psutil.virtual_memory().percent
        if memory_percent > 80:
            logging.warning(f"High memory usage detected: {memory_percent}%. Consider reducing batch size.")
    except ImportError:
        pass  # psutil not available

    staged_docs = []
    schema_json_str = json.dumps(RichMetadata.model_json_schema(), indent=2)
    prompt_lines = [
        "Analyze the document content and its file path to return a single, valid JSON object that strictly conforms to this Pydantic schema:", "```json", "{schema}", "```", "---",
        "**SPECIAL INSTRUCTIONS:**",
        '- If the source is an image (`source_type` is "image"), you **MUST** set `document_type` to "Image/Diagram".',
        '- If the file is a `.pptx`, you **MUST** set `document_type` to "Presentation".',
        '- If the `file_path` contains "trophy" or the filename mentions "Case Study", you **MUST** set `document_type` to "Case Study / Trophy".',
        '- If the filename contains the word "Final", you should strongly prefer the "Final Report" type.', '- If the filename contains the word "Draft", you should strongly prefer the "Draft Report" type.',
        '- If no other category seems appropriate, you **MUST** use "Other" as a fallback.', "---", "File Path: {file_path}", "Source Type: {source_type}", "Document Content (first 8000 characters):",
        "-----------------", "{text}", "-----------------", "IMPORTANT: Your response must be ONLY the JSON object itself, with no extra text, explanations, or wrapper keys."
    ]
    metadata_prompt_template = "\n".join(prompt_lines)

    for i, doc in enumerate(unique_docs):
        # Check for pause request
        if batch_manager and batch_manager.is_paused():
            logging.info("Batch processing paused by user request")
            break
            
        rich_metadata = None
        file_path_str, file_name = doc.metadata.get('file_path', ''), doc.metadata.get('file_name', 'Unknown File')
        source_type = doc.metadata.get('source_type', 'document')
        logging.info(f"--- ({i+1}/{len(unique_docs)}) Analyzing: {file_name} ({source_type.upper()}) ---")
        
        # Print machine-readable progress for the UI
        print(f"CORTEX_PROGRESS::{i+1}/{len(unique_docs)}::{file_name}", flush=True)
        
        try:
            if not doc.text.strip(): 
                raise ValueError("Document is empty or could not be read.")
            
            # Check if Ollama LLM is available
            if Settings.llm is None:
                logging.info(f"Ollama not available - using basic metadata for {file_name}")
                # Create basic metadata when Ollama is not available
                rich_metadata = RichMetadata(
                    document_type="Other",
                    proposal_outcome="N/A", 
                    summary=f"Document processed without AI analysis (Ollama unavailable). File: {file_name}",
                    thematic_tags=["basic-processing", "no-ai-analysis"]
                )
            else:
                try:
                    prompt = metadata_prompt_template.format(schema=schema_json_str, text=doc.get_content()[:8000], file_path=file_path_str, source_type=source_type)
                    logging.info("Sending prompt to LLM for metadata extraction...")

                    # Add timeout handling with progress feedback
                    import time
                    start_time = time.time()
                    try:
                        response_str = str(Settings.llm.complete(prompt))
                        elapsed = time.time() - start_time
                        if elapsed > 60:
                            logging.info(f"LLM call took {elapsed:.1f}s (longer document)")
                    except TimeoutError as te:
                        logging.warning(f"LLM timeout after {time.time() - start_time:.1f}s for {file_name}, using fallback metadata")
                        raise ValueError(f"LLM timeout: {te}")
                    except Exception as llm_ex:
                        logging.warning(f"LLM error after {time.time() - start_time:.1f}s for {file_name}: {llm_ex}")
                        raise

                    logging.info("Received raw response from LLM. Cleaning and parsing JSON...")
                    
                    # Check if response is empty (indicates LLM error)
                    if not response_str or response_str.strip() == "":
                        raise ValueError("LLM returned empty response")
                    
                    json_match = re.search(r'\{.*\}', response_str, re.DOTALL)
                    if not json_match:
                        error_snippet = response_str.strip().replace('\n', ' ')[:200]
                        logging.error(f"LLM did not return valid JSON for {file_name}. Response: {error_snippet}...")
                        raise ValueError("LLM did not return a valid JSON object.")
                    
                    clean_json_str = json_match.group(0)
                    metadata_json = json.loads(clean_json_str)
                    
                    # Only process if we got valid JSON data
                    if 'metadata_json' in locals():
                        # Normalize common validation fields before Pydantic
                        if 'proposal_outcome' in metadata_json and isinstance(metadata_json['proposal_outcome'], str):
                            metadata_json['proposal_outcome'] = metadata_json['proposal_outcome'].title()

                        try: 
                            rich_metadata = RichMetadata.model_validate(metadata_json)
                        except ValidationError as e:
                            logging.warning(f"Initial validation failed for {file_name}. Retrying with nested key check...")
                            if isinstance(metadata_json, dict) and len(metadata_json) == 1:
                                nested_key = list(metadata_json.keys())[0]
                                if isinstance(metadata_json[nested_key], dict): 
                                    rich_metadata = RichMetadata.model_validate(metadata_json[nested_key])
                                else: 
                                    raise e
                            else:
                                raise e
                                
                except (ValueError, json.JSONDecodeError, RuntimeError, Exception) as llm_error:
                    logging.warning(f"LLM processing failed for {file_name}: {llm_error}")
                    logging.info(f"Falling back to basic metadata for {file_name}")
                    # Create fallback metadata when LLM fails
                    rich_metadata = RichMetadata(
                        document_type="Other",
                        proposal_outcome="N/A", 
                        summary=f"Document processed with fallback metadata (LLM error). File: {file_name}",
                        thematic_tags=["llm-error", "fallback-metadata"]
                    )
            logging.info(f"Successfully parsed and validated metadata for {file_name}.")
            
            # Extract entities and relationships
            logging.info(f"Extracting entities and relationships from {file_name}...")
            if rich_metadata is not None:
                entities_dict, relationships_dict = extract_entities_and_relationships(
                    doc.get_content()[:8000],
                    {
                        'document_type': rich_metadata.document_type,
                        'file_name': file_name,
                        'summary': rich_metadata.summary,
                        'thematic_tags': rich_metadata.thematic_tags
                    },
                    entity_extractor
                )
            else:
                logging.warning(f"Skipping entity extraction for {file_name} - metadata is None")
                entities_dict, relationships_dict = {}, {}
            logging.info(f"Extracted {len(entities_dict)} entities and {len(relationships_dict)} relationships.")
            
        except Exception as e:
            logging.error(f"CRITICAL ERROR analyzing {file_name}: {e}", exc_info=True)
            # Record error in batch manager
            if batch_manager:
                batch_manager.record_error(file_path_str, str(e))
            
            default_doc_type = "Image/Diagram" if source_type == 'image' else 'Other'
            rich_metadata = RichMetadata(
                document_type=default_doc_type, 
                proposal_outcome="N/A", 
                summary=f"ERROR: Could not analyze document. Reason: {e}", 
                thematic_tags=["error", "analysis-failed"]
            )
            entities_dict = []
            relationships_dict = []

        doc_meta = DocumentMetadata(
            doc_id=get_file_hash(file_path_str), 
            doc_posix_path=Path(file_path_str).as_posix(),
            file_name=doc.metadata.get('file_name'), 
            last_modified_date=str(datetime.fromtimestamp(os.path.getmtime(file_path_str))),
            rich_metadata=rich_metadata,
            extracted_entities=entities_dict,
            extracted_relationships=relationships_dict
        )
        staged_docs.append(doc_meta.model_dump())
        
        # Update batch progress for successfully processed file
        if batch_manager:
            batch_manager.update_progress(file_path_str)

    # Create staging data structure with collection assignment
    staging_data = {
        "documents": staged_docs,  # staged_docs already contains dicts from doc_meta.model_dump()
        "target_collection": target_collection,
        "created_at": str(datetime.now()),
        "version": "2.0"  # New format with collection assignment
    }
    
    with open(staging_file, 'w') as f: 
        json.dump(staging_data, f, indent=2)
    
    collection_info = f" -> '{target_collection}'" if target_collection else " -> 'default'"
    logging.info(f"--- Analysis complete. {len(staged_docs)} documents staged{collection_info} at {staging_file} ---")
    print(f"CORTEX_STAGE::ANALYSIS_DONE::{len(staged_docs)}::{staging_file}", flush=True)

def finalize_ingestion(db_path: str, args=None):
    logging.info(f"--- Starting Stage 3: Finalize from Staging with Graph Building (Cortex v13.0.0) ---")
    print("CORTEX_STAGE::FINALIZE_START", flush=True)
    staging_file = get_staging_file_path(db_path)
    if not os.path.exists(staging_file): 
        logging.error(f"Staging file not found at: {staging_file}")
        return
    
    chroma_db_path = os.path.join(db_path, "knowledge_hub_db")
    os.makedirs(chroma_db_path, exist_ok=True)
    
    # Initialize graph manager
    graph_file_path = os.path.join(db_path, "knowledge_cortex.gpickle")
    graph_manager = EnhancedGraphManager(graph_file_path)
    
    with open(staging_file, 'r') as f: 
        staging_data = json.load(f)
    
    # Handle both old format (list) and new format (dict with metadata)
    if isinstance(staging_data, list):
        docs_to_process = [DocumentMetadata(**data) for data in staging_data]
        target_collection = None  # No collection assignment for old format
    else:
        docs_to_process = [DocumentMetadata(**data) for data in staging_data.get('documents', [])]
        target_collection = staging_data.get('target_collection', None)
    
    logging.info(f"Target collection for finalization: {target_collection or 'default'}")

    docs_to_index_paths, metadata_map, doc_ids_to_add_to_default = [], {}, []
    processed_log_path = os.path.join(chroma_db_path, INGESTED_FILES_LOG)
    
    for doc_meta in docs_to_process:
        if doc_meta.exclude_from_final:
            logging.info(f"User excluded {doc_meta.file_name}. Skipping.")
            write_to_processed_log(processed_log_path, doc_meta.doc_posix_path, doc_meta.doc_id)
            continue
        if doc_meta.rich_metadata and "ERROR:" in doc_meta.rich_metadata.summary:
            logging.warning(f"Skipping finalization for {doc_meta.file_name} due to prior analysis error.")
            continue
        
        docs_to_index_paths.append(doc_meta.doc_posix_path)
        metadata_map[doc_meta.doc_posix_path] = doc_meta
        doc_ids_to_add_to_default.append(doc_meta.doc_id)
        
        # Add document to graph
        logging.info(f"Adding {doc_meta.file_name} to knowledge graph...")
        graph_manager.add_entity(
            doc_meta.doc_id,
            'Document',
            file_name=doc_meta.file_name,
            document_type=doc_meta.rich_metadata.document_type if doc_meta.rich_metadata else 'Unknown',
            summary=doc_meta.rich_metadata.summary if doc_meta.rich_metadata else '',
            last_modified=doc_meta.last_modified_date,
            posix_path=doc_meta.doc_posix_path
        )
        
        # Add entities and relationships to graph
        entity_map = {}  # Track entity names to node IDs
        
        for entity_dict in doc_meta.extracted_entities:
            entity_id = f"{entity_dict['entity_type']}:{entity_dict['name']}"
            entity_map[entity_dict['name']] = entity_id
            
            # Add entity if it doesn't exist
            if entity_id not in graph_manager.graph:
                graph_manager.add_entity(
                    entity_id,
                    entity_dict['entity_type'],
                    name=entity_dict['name'],
                    aliases=entity_dict.get('aliases', [])
                )
            
            # Link entity to document
            if entity_dict['entity_type'] == 'person':
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'authored'
                )
            elif entity_dict['entity_type'] == 'organization':
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'client_of'
                )
            else:
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'mentioned_in'
                )
        
        # Add extracted relationships
        for rel_dict in doc_meta.extracted_relationships:
            source_id = entity_map.get(rel_dict['source'], rel_dict['source'])
            target_id = entity_map.get(rel_dict['target'], rel_dict['target'])
            
            # Handle special case where target might be the document
            if rel_dict['target'] == doc_meta.file_name:
                target_id = doc_meta.doc_id
            
            graph_manager.add_relationship(
                source_id,
                target_id,
                rel_dict['relationship_type'],
                context=rel_dict.get('context', '')
            )

    if not docs_to_index_paths:
        logging.warning("No new, valid documents to ingest. Finalization complete.")
        if os.path.exists(staging_file): 
            os.remove(staging_file)
        return

    # Save the graph
    graph_manager.save_graph()
    logging.info(f"Knowledge graph saved to {graph_file_path}")

    # Clear query cache since new data was added
    try:
        from cortex_engine.graph_query import clear_query_cache
        clear_query_cache()
        logging.info("🔄 Query cache cleared due to new ingestion")
    except Exception as e:
        logging.warning(f"Failed to clear query cache: {e}")

    # Continue with regular vector indexing
    documents_for_indexing = manual_load_documents(docs_to_index_paths, args)
    for doc in documents_for_indexing:
        path_key = doc.metadata['file_path']
        if path_key in metadata_map:
            doc_meta = metadata_map[path_key]
            doc.doc_id = doc_meta.doc_id
            flat_metadata = {
                "doc_id": doc_meta.doc_id, 
                "file_name": doc_meta.file_name, 
                "doc_posix_path": doc_meta.doc_posix_path, 
                "last_modified_date": doc_meta.last_modified_date
            }
            if doc_meta.rich_metadata:
                flat_metadata.update(doc_meta.rich_metadata.model_dump())
                flat_metadata['thematic_tags'] = ', '.join(flat_metadata.get('thematic_tags', []))
            doc.metadata = flat_metadata

    db_settings = ChromaSettings(anonymized_telemetry=False)
    chroma_client = chromadb.PersistentClient(path=chroma_db_path, settings=db_settings)
    chroma_collection = chroma_client.get_or_create_collection(COLLECTION_NAME)
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    storage_context = StorageContext.from_defaults(vector_store=vector_store)
    index = VectorStoreIndex.from_documents(documents_for_indexing, storage_context=storage_context, show_progress=True)
    logging.info(f"Persisting index to disk at {chroma_db_path}...")
    index.storage_context.persist(persist_dir=chroma_db_path)
    
    for doc in documents_for_indexing: 
        write_to_processed_log(processed_log_path, doc.metadata['doc_posix_path'], doc.metadata['doc_id'])
    os.remove(staging_file)

    if doc_ids_to_add_to_default:
        # Use target collection if specified, otherwise default
        collection_name = target_collection or "default"
        logging.info(f"Adding {len(doc_ids_to_add_to_default)} new documents to the '{collection_name}' collection.")
        try:
            from cortex_engine.collection_manager import WorkingCollectionManager
            collection_mgr = WorkingCollectionManager()
            collection_mgr.add_docs_by_id_to_collection(collection_name, doc_ids_to_add_to_default)
            logging.info("Collections updated via WorkingCollectionManager (Docker)")
        except Exception as e:
            logging.error(f"Could not automatically add documents to '{collection_name}' collection: {e}")
    
    logging.info("--- Finalization complete. Knowledge base and graph are up to date. ---")
    print("CORTEX_STAGE::FINALIZE_DONE", flush=True)

def main():
    parser = argparse.ArgumentParser(description="Cortex Ingestion Engine with GraphRAG")
    parser.add_argument("--db-path", type=str, required=True)
    parser.add_argument("--include", type=str, nargs='*')
    parser.add_argument("--analyze-only", action="store_true")
    parser.add_argument("--finalize-from-staging", action="store_true")
    parser.add_argument("--fresh", action="store_true")
    parser.add_argument("--skip-image-processing", action="store_true", 
                       help="Skip VLM image processing for faster ingestion")
    parser.add_argument("--resume", action="store_true",
                       help="Resume from existing batch state (opposite of --fresh)")
    parser.add_argument("--target-collection", type=str, default=None,
                       help="Target collection for document assignment")
    args = parser.parse_args()

    initialize_script()

    try:
        if args.analyze_only:
            if not args.include: 
                logging.error("--include paths are required.")
                sys.exit(1)
            # Resume takes precedence over fresh
            fresh_start = args.fresh and not args.resume
            analyze_documents(args.include, fresh_start, args, getattr(args, 'target_collection', None))
        elif args.finalize_from_staging:
            finalize_ingestion(args.db_path, args)
        else:
            logging.error("Specify either --analyze-only or --finalize-from-staging.")
            sys.exit(1)
    except KeyboardInterrupt:
        logging.info("Process interrupted by user (Ctrl+C)")
        sys.exit(0)
    except Exception as e:
        logging.error(f"CRITICAL ERROR: Batch processing failed with exception: {e}", exc_info=True)
        # Try to save any progress made so far
        try:
            if hasattr(args, 'db_path') and args.db_path:
                batch_manager = BatchState(args.db_path)
                batch_manager.record_error("SYSTEM_CRASH", f"Critical system error: {e}")
        except Exception as save_error:
            logging.error(f"Failed to save error state: {save_error}")  # Log but don't re-raise
        sys.exit(1)

if __name__ == "__main__":
    main()
