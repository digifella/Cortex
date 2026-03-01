"""
Document Textifier Module
Converts PDF, DOCX, and PPTX documents to rich Markdown with optional
vision-model descriptions of images and tables.
"""

import io
import os
import re
import base64
import subprocess
import tempfile
import statistics
import time
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple, Any

from cortex_engine.utils.logging_utils import get_logger

logger = get_logger(__name__)


class DocumentTextifier:
    """Converts documents to Markdown with optional VLM image descriptions."""

    # Preferred vision models in order of priority
    VISION_MODELS = ["qwen3-vl:8b", "qwen3-vl:4b", "qwen3-vl"]

    def __init__(
        self,
        use_vision: bool = True,
        on_progress: Optional[Callable[[float, str], None]] = None,
        pdf_strategy: str = "hybrid",
        cleanup_provider: Optional[str] = None,
        cleanup_model: Optional[str] = None,
        docling_timeout_seconds: Optional[float] = None,
        image_description_timeout_seconds: Optional[float] = None,
        image_enrich_max_seconds: Optional[float] = None,
    ):
        self.use_vision = use_vision
        self.on_progress = on_progress
        self.pdf_strategy = (pdf_strategy or "docling").strip().lower()
        self.cleanup_provider = (cleanup_provider or "").strip().lower()
        self.cleanup_model = (cleanup_model or "").strip()
        self.docling_timeout_seconds = float(
            docling_timeout_seconds
            if docling_timeout_seconds is not None
            else os.getenv("CORTEX_TEXTIFIER_DOCLING_TIMEOUT_SECONDS", "240")
        )
        self.image_description_timeout_seconds = float(
            image_description_timeout_seconds
            if image_description_timeout_seconds is not None
            else os.getenv("CORTEX_TEXTIFIER_IMAGE_TIMEOUT_SECONDS", "20")
        )
        self.image_enrich_max_seconds = float(
            image_enrich_max_seconds
            if image_enrich_max_seconds is not None
            else os.getenv("CORTEX_TEXTIFIER_IMAGE_ENRICH_MAX_SECONDS", "120")
        )
        self._vlm_client = None
        self._vlm_model = None
        self._vlm_skip_models: set[str] = set()
        self._last_cleanup_applied = False

    @classmethod
    def from_options(
        cls,
        options: Optional[Dict[str, Any]] = None,
        on_progress: Optional[Callable[[float, str], None]] = None,
    ) -> "DocumentTextifier":
        """Create a DocumentTextifier from a shared options payload."""
        opts = dict(options or {})
        return cls(
            use_vision=bool(opts.get("use_vision", True)),
            on_progress=on_progress,
            pdf_strategy=str(opts.get("pdf_strategy", "hybrid") or "hybrid").strip().lower(),
            cleanup_provider=str(opts.get("cleanup_provider", "") or "").strip().lower() or None,
            cleanup_model=str(opts.get("cleanup_model", "") or "").strip() or None,
            docling_timeout_seconds=float(opts.get("docling_timeout_seconds", 240)),
            image_description_timeout_seconds=float(opts.get("image_description_timeout_seconds", 20)),
            image_enrich_max_seconds=float(opts.get("image_enrich_max_seconds", 120)),
        )

    def _report(self, fraction: float, message: str):
        """Send a progress update if a callback is registered."""
        if self.on_progress:
            self.on_progress(fraction, message)

    @staticmethod
    def _is_structural_markdown_line(line: str) -> bool:
        stripped = (line or "").strip()
        if not stripped:
            return True
        if stripped.startswith(("#", ">", "|", "```", "---")):
            return True
        if re.match(r"^[-*+]\s+", stripped):
            return True
        if re.match(r"^\d+\.\s+", stripped):
            return True
        return False

    def _normalize_markdown_output(self, markdown_text: str) -> str:
        """
        Normalize markdown text by:
        - removing literal CR/CRLF markers
        - normalizing newlines
        - reflowing hard-wrapped plain lines into paragraphs
        while preserving markdown structure.
        """
        text = (markdown_text or "")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"\s*<CRLF>\s*", "\n", text, flags=re.IGNORECASE)
        text = re.sub(r"\s*<CR>\s*", "\n", text, flags=re.IGNORECASE)
        text = re.sub(r"\n{3,}", "\n\n", text)

        lines = [ln.rstrip() for ln in text.split("\n")]
        out_lines: List[str] = []
        para_buffer: List[str] = []
        in_code_fence = False

        def flush_paragraph() -> None:
            if not para_buffer:
                return
            paragraph = " ".join(part.strip() for part in para_buffer if part.strip())
            if paragraph:
                out_lines.append(paragraph)
            para_buffer.clear()

        for raw in lines:
            line = raw.strip()
            if line.startswith("```"):
                flush_paragraph()
                in_code_fence = not in_code_fence
                out_lines.append(raw)
                continue

            if in_code_fence:
                out_lines.append(raw)
                continue

            if not line:
                flush_paragraph()
                if out_lines and out_lines[-1] != "":
                    out_lines.append("")
                continue

            if self._is_structural_markdown_line(line):
                flush_paragraph()
                out_lines.append(raw)
                continue

            para_buffer.append(line)

        flush_paragraph()
        normalized = "\n".join(out_lines).strip()
        return normalized + "\n" if normalized else ""

    def _init_vlm(self):
        """Lazy-init the Ollama VLM client, preferring Qwen3-VL models."""
        if self._vlm_client is not None:
            return
        try:
            import ollama
            client = ollama.Client(timeout=120)

            # Try preferred Qwen3-VL models first
            for model in self.VISION_MODELS:
                try:
                    client.show(model)
                    self._vlm_client = client
                    self._vlm_model = model
                    logger.info(f"Textifier using vision model: {model}")
                    return
                except Exception:
                    continue

            # Fall back to configured VLM_MODEL (e.g. llava:7b)
            from cortex_engine.config import VLM_MODEL
            client.show(VLM_MODEL)
            self._vlm_client = client
            self._vlm_model = VLM_MODEL
            logger.info(f"Textifier falling back to VLM_MODEL: {VLM_MODEL}")
        except Exception as e:
            logger.warning(f"No vision model available: {e}")
            self._vlm_client = None
            self._vlm_model = None

    def _vlm_model_candidates(self, include_current: bool = True) -> List[str]:
        """Return a de-duplicated ordered list of VLM model names to try."""
        candidates: List[str] = []
        if include_current and self._vlm_model:
            candidates.append(self._vlm_model)
        candidates.extend(self.VISION_MODELS)
        try:
            from cortex_engine.config import VLM_MODEL
            if VLM_MODEL:
                candidates.append(str(VLM_MODEL).strip())
        except Exception:
            pass
        candidates.extend([
            "llava:7b",
            "llava:latest",
            "llava",
            "minicpm-v:latest",
            "minicpm-v",
            "moondream:latest",
            "moondream",
        ])
        ordered: List[str] = []
        seen = set()
        for model in candidates:
            name = str(model or "").strip()
            if not name or name in seen or name in self._vlm_skip_models:
                continue
            seen.add(name)
            ordered.append(name)
        return ordered

    def _reset_vlm_session(self, reason: str = "") -> None:
        """Drop cached Ollama client/model so the next image starts with a fresh VLM session."""
        if reason:
            logger.info(f"Resetting cached VLM session: {reason}")
        self._vlm_client = None
        self._vlm_model = None

    def _describe_with_model(self, model: str, encoded_image: str, simple_prompt: bool = False) -> str:
        """Call a specific VLM model and normalize the returned text."""
        prompt = (
            "Describe this photograph very briefly in 1-2 short plain sentences (max 35 words total). "
            "Identify only the main subject and setting. "
            "If the image is primarily a logo/icon/watermark or tiny decorative graphic, "
            "return exactly: [Image: logo/icon omitted]. "
            "Do not use markdown, headings, or bullet points."
        )
        if not simple_prompt:
            prompt += " /no_think"
        response = self._vlm_client.chat(
            model=model,
            messages=[{
                "role": "user",
                "content": prompt,
                "images": [encoded_image],
            }],
            options={"temperature": 0.1, "num_predict": 140},
        )
        result = ((response or {}).get("message", {}) or {}).get("content", "")
        result = str(result or "").strip()
        return re.sub(r"<think>.*?</think>", "", result, flags=re.DOTALL).strip()

    def _prepare_vlm_image_bytes(self, image_bytes: bytes) -> bytes:
        """
        Downscale large images before sending them to the VLM.

        This reduces Ollama memory pressure without changing the source file on disk.
        """
        max_edge = int(float(os.environ.get("CORTEX_VLM_MAX_EDGE_PX", "1600") or 1600))
        max_payload_bytes = int(float(os.environ.get("CORTEX_VLM_MAX_INPUT_BYTES", str(1600 * 1024)) or (1600 * 1024)))
        if max_edge <= 0 or not image_bytes:
            return image_bytes
        if len(image_bytes) <= max_payload_bytes:
            # Still inspect dimensions, but small files often don't need work.
            pass
        try:
            from PIL import Image
        except Exception:
            return image_bytes

        try:
            with Image.open(io.BytesIO(image_bytes)) as img:
                width, height = img.size
                long_edge = max(width, height)
                if long_edge <= max_edge and len(image_bytes) <= max_payload_bytes:
                    return image_bytes

                # Ensure a broadly compatible RGB payload for Ollama VLMs.
                if img.mode not in ("RGB",):
                    img = img.convert("RGB")

                working = img.copy()
                working.thumbnail((max_edge, max_edge), Image.Resampling.LANCZOS)
                out = io.BytesIO()
                working.save(out, format="JPEG", quality=88, optimize=True)
                reduced = out.getvalue()
                if reduced and len(reduced) < len(image_bytes):
                    logger.info(
                        "Downscaled VLM input from %sx%s (%s KB) to %sx%s (%s KB)",
                        width,
                        height,
                        max(1, len(image_bytes) // 1024),
                        working.size[0],
                        working.size[1],
                        max(1, len(reduced) // 1024),
                    )
                    return reduced
        except Exception as e:
            logger.debug(f"Could not downscale VLM input image: {e}")
        return image_bytes

    @staticmethod
    def _looks_like_logo_icon_description(text: str) -> bool:
        """Heuristic: detect short logo/icon-only descriptions that add metadata noise."""
        t = (text or "").strip().lower()
        if not t:
            return False
        logo_markers = [
            "logo", "icon", "icons", "emblem", "symbol", "badge", "watermark",
            "seal", "silhouette", "initials", "letters",
        ]
        visual_markers = [
            "left icon", "right icon", "circular", "circle", "black and white",
            "background", "strip", "main colors", "main colours",
        ]
        marker_hits = sum(1 for m in logo_markers if m in t)
        visual_hits = sum(1 for m in visual_markers if m in t)
        return (marker_hits >= 1 and visual_hits >= 1) or marker_hits >= 2

    def describe_image(self, image_bytes: bytes) -> str:
        """Describe an image using the VLM. Returns placeholder on failure."""
        if not self.use_vision:
            return "[Image: vision model disabled]"
        self._init_vlm()
        if self._vlm_client is None:
            return "[Image: could not be described — vision model unavailable]"
        try:
            prepared_bytes = self._prepare_vlm_image_bytes(image_bytes)
            encoded = base64.b64encode(prepared_bytes).decode("utf-8")
            models_tried: List[str] = []
            for idx, model in enumerate(self._vlm_model_candidates(include_current=True)):
                try:
                    if idx > 0:
                        self._vlm_client.show(model)
                    result = self._describe_with_model(model, encoded, simple_prompt=False)
                except Exception as model_error:
                    err_text = str(model_error)
                    if "status code: 404" in err_text:
                        self._vlm_skip_models.add(model)
                    if "unexpectedly stopped" in err_text and model == self._vlm_model:
                        self._reset_vlm_session(f"{model} runner stopped unexpectedly")
                    logger.warning(f"VLM describe failed for model {model}: {model_error}")
                    models_tried.append(f"{model} (error)")
                    continue
                models_tried.append(model)
                if not result:
                    logger.warning(f"VLM returned empty description (model: {model})")
                    if model.startswith("qwen3-vl"):
                        try:
                            retry_result = self._describe_with_model(model, encoded, simple_prompt=True)
                        except Exception as retry_error:
                            logger.warning(f"VLM simple-prompt retry failed for model {model}: {retry_error}")
                            retry_result = ""
                        if retry_result:
                            logger.info(f"Recovered image description using simple-prompt retry on {model}")
                            if self._looks_like_logo_icon_description(retry_result):
                                return "[Image: logo/icon omitted]"
                            return retry_result
                    continue
                if model != self._vlm_model:
                    logger.info(f"Recovered image description using fallback vision model: {model}")
                    self._vlm_model = model
                if self._looks_like_logo_icon_description(result):
                    return "[Image: logo/icon omitted]"
                return result
            self._reset_vlm_session("no non-empty VLM description returned")
            logger.warning(f"No non-empty VLM description returned after trying: {', '.join(models_tried)}")
            return "[Image: vision model returned empty description]"
        except Exception as e:
            logger.warning(f"VLM describe failed: {e}")
            return "[Image: could not be described — vision model error]"

    def _describe_image_with_timeout(self, image_bytes: bytes) -> str:
        """Run image description with a hard timeout so long VLM calls don't stall conversion."""
        timeout_s = self.image_description_timeout_seconds
        if timeout_s is None or timeout_s <= 0:
            return self.describe_image(image_bytes)
        executor = ThreadPoolExecutor(max_workers=1)
        future = executor.submit(self.describe_image, image_bytes)
        try:
            return future.result(timeout=timeout_s)
        except FuturesTimeoutError:
            future.cancel()
            return "[Image: description timed out]"
        except Exception:
            return "[Image: could not be described — vision model error]"
        finally:
            executor.shutdown(wait=False, cancel_futures=True)

    @staticmethod
    def _run_with_timeout(fn, timeout_seconds: Optional[float], *args, **kwargs):
        """Run a callable with optional timeout. Returns None when timed out."""
        if timeout_seconds is None or timeout_seconds <= 0:
            return fn(*args, **kwargs)
        executor = ThreadPoolExecutor(max_workers=1)
        future = executor.submit(fn, *args, **kwargs)
        try:
            return future.result(timeout=timeout_seconds)
        except FuturesTimeoutError:
            future.cancel()
            return None
        finally:
            executor.shutdown(wait=False, cancel_futures=True)

    @staticmethod
    def table_to_markdown(rows: List[List[str]]) -> str:
        """Convert a list-of-lists table to Markdown."""
        if not rows:
            return ""
        # Ensure all rows have same column count
        max_cols = max(len(r) for r in rows)
        normalised = [r + [""] * (max_cols - len(r)) for r in rows]
        header = normalised[0]
        lines = ["| " + " | ".join(header) + " |"]
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in normalised[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)

    @staticmethod
    def _clean_table_rows(rows: List[List[str]]) -> List[List[str]]:
        cleaned_rows: List[List[str]] = []
        for row in rows or []:
            cleaned_row = []
            for cell in row or []:
                text = re.sub(r"\s+", " ", str(cell or "")).strip()
                cleaned_row.append(text)
            cleaned_rows.append(cleaned_row)
        return cleaned_rows

    @staticmethod
    def _is_chart_noise_line(line: str) -> bool:
        s = (line or "").strip()
        if not s:
            return False
        low = s.lower()
        if re.fullmatch(r"\d+(?:\.\d+)?%?", s):
            return True
        if re.fullmatch(r"(?:19|20)\d{2}", s):
            return True
        if re.search(r"\bvs\b", low) and len(s) <= 10:
            return True
        if re.fullmatch(r"[0-9%\s\.\-–—]+", s):
            return True
        if re.search(r"\b(life expectancy|key facts|key figures|infographic|chart)\b", low):
            return True
        tokens = s.split()
        if 2 <= len(tokens) <= 8:
            numeric_like = sum(1 for t in tokens if re.fullmatch(r"\d+(?:\.\d+)?%?", t))
            if numeric_like >= max(2, len(tokens) // 2):
                return True
        return False

    def _strip_infographic_noise(self, lines: List[str]) -> List[str]:
        if not lines:
            return []
        marker_idx = -1
        for i, line in enumerate(lines):
            if re.search(r"\b(infographic|figure|chart)\b", line, flags=re.IGNORECASE):
                marker_idx = i
                break

        if marker_idx >= 0:
            # For infographic/chart pages, keep preface text only and drop the visual payload lines.
            kept_prefix = [
                ln for ln in lines[:marker_idx]
                if ln and not self._is_chart_noise_line(ln)
            ]
            kept_prefix.append("[Infographic/Figure content omitted in strict text-only mode.]")
            return kept_prefix

        return [ln for ln in lines if not self._is_chart_noise_line(ln)]

    @staticmethod
    def _is_simple_table_quality(rows: List[List[str]]) -> bool:
        if not rows or len(rows) < 2:
            return False
        col_count = max((len(r) for r in rows), default=0)
        if col_count < 2:
            return False
        if col_count > 12:
            return False

        total_cells = len(rows) * col_count
        filled_cells = 0
        long_cell_count = 0
        for row in rows:
            normalized = row + [""] * (col_count - len(row))
            for cell in normalized:
                if cell:
                    filled_cells += 1
                if len(cell) > 140:
                    long_cell_count += 1

        fill_ratio = (filled_cells / total_cells) if total_cells else 0.0
        if fill_ratio < 0.35:
            return False
        if long_cell_count > max(2, int(total_cells * 0.15)):
            return False

        header = rows[0] + [""] * (col_count - len(rows[0]))
        header_nonempty = sum(1 for c in header if c)
        if header_nonempty < 1:
            return False
        return True

    def _extract_pdf_tables(self, page) -> List[Dict[str, object]]:
        """Extract simple tables from a PDF page, returning parsed/failed statuses."""
        results: List[Dict[str, object]] = []
        if not hasattr(page, "find_tables"):
            return results
        try:
            found = page.find_tables()
            table_objs = list(getattr(found, "tables", []) or [])
        except Exception as e:
            logger.debug(f"PDF table detection failed on page {page.number + 1}: {e}")
            return results

        for idx, table_obj in enumerate(table_objs, 1):
            try:
                raw_rows = table_obj.extract() or []
                rows = self._clean_table_rows(raw_rows)
                if self._is_simple_table_quality(rows):
                    results.append({"status": "parsed", "index": idx, "rows": rows})
                else:
                    results.append({"status": "unreliable", "index": idx, "rows": []})
            except Exception as e:
                logger.debug(f"Table extraction failed for table {idx} on page {page.number + 1}: {e}")
                results.append({"status": "unreliable", "index": idx, "rows": []})
        return results

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------

    def textify_file(self, file_path: str) -> str:
        """Dispatch to the correct converter based on file extension."""
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self.textify_pdf(file_path)
        elif ext == ".docx":
            return self.textify_docx(file_path)
        elif ext == ".pptx":
            return self.textify_pptx(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"):
            return self.textify_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def textify_pdf(self, file_path: str) -> str:
        """Convert a PDF to Markdown using PyMuPDF."""
        strategy = self.pdf_strategy
        if strategy in {"docling", "hybrid"}:
            self._report(0.0, "Checking Docling availability...")
            docling_md = self._try_docling(file_path, timeout_seconds=self.docling_timeout_seconds)
            if docling_md:
                self._report(1.0, "Converted via Docling")
                return docling_md

        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        total_pages = len(doc)
        md_parts: List[str] = []
        page_lines: List[List[str]] = []
        header_footer_counts: Dict[str, int] = {}

        # Pass 1: collect lines and detect repeated header/footer boilerplate.
        for page_num in range(total_pages):
            page = doc[page_num]
            raw_text = page.get_text("text") or ""
            lines = [re.sub(r"\s+", " ", ln).strip() for ln in raw_text.splitlines() if ln.strip()]
            page_lines.append(lines)

            candidates = lines[:4] + lines[-4:]
            for line in candidates:
                if len(line) < 4:
                    continue
                if re.fullmatch(r"[\d\W]+", line):
                    continue
                header_footer_counts[line] = header_footer_counts.get(line, 0) + 1

        repeated_boilerplate = {line for line, count in header_footer_counts.items() if count >= 2}

        for page_num in range(total_pages):
            self._report(page_num / total_pages, f"Page {page_num + 1}/{total_pages}")
            page = doc[page_num]
            md_parts.append(f"## Page {page_num + 1}\n")

            lines = self._extract_pdf_page_structured_lines(page)
            if not lines:
                lines = page_lines[page_num] if page_num < len(page_lines) else []
            filtered_lines: List[str] = []
            for idx, line in enumerate(lines):
                in_header_footer_zone = idx < 5 or idx >= max(len(lines) - 5, 0)
                if in_header_footer_zone and line in repeated_boilerplate:
                    continue
                filtered_lines.append(line)
            filtered_lines = self._strip_infographic_noise(filtered_lines)

            text = "\n".join(filtered_lines).strip()
            if text:
                md_parts.append(text)
                md_parts.append("")
            elif self.use_vision:
                # If no extractable text, include a vision summary so the page isn't lost.
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                    desc = self.describe_image(pix.tobytes("png"))
                    md_parts.append(f"> **[Page {page_num + 1} image summary]**: {desc}")
                    md_parts.append("")
                except Exception as e:
                    logger.debug(f"Page image summary failed on page {page_num + 1}: {e}")

            # Strict text-only PDF mode:
            # intentionally ignore embedded tables/figures/images to avoid noisy pseudo-structured output.

            # Explicit page break marker for downstream parsers.
            if page_num < total_pages - 1:
                md_parts.append("\n---\n")

        doc.close()
        self._report(1.0, "PDF conversion complete")
        normalized = self._normalize_markdown_output("\n".join(md_parts))
        cleanup_enabled = strategy in {"qwen30b", "hybrid"} or (
            os.getenv("CORTEX_TEXTIFIER_SEMANTIC_CLEANUP", "1").strip().lower() not in {"0", "false", "no"}
        )
        if cleanup_enabled:
            cleaned = self._semantic_markdown_cleanup(normalized)
            if cleaned and cleaned.strip():
                out = self._normalize_markdown_output(cleaned)
                # If explicit Qwen mode, return cleanup output even if only small change.
                if strategy == "qwen30b":
                    return out
                # Hybrid accepts cleanup only when it actually improved structure.
                if strategy != "hybrid" or (out != normalized and self._last_cleanup_applied):
                    return out

        return normalized

    def extract_pdf_images(
        self,
        file_path: str,
        min_width_px: int = 500,
        min_height_px: int = 500,
        min_page_coverage_pct: float = 2.0,
        ignore_edge_decorations: bool = True,
        edge_margin_pct: float = 8.0,
        render_scale: float = 2.0,
        split_full_page_scans: bool = False,
        split_min_crop_coverage_pct: float = 6.0,
    ) -> Dict[str, Any]:
        """
        Extract likely photographic image regions from a PDF.

        Uses PyMuPDF image blocks so the extracted JPEG matches what is visually
        placed on the page, then filters out likely icons / header-footer art.
        """
        import fitz  # PyMuPDF

        pdf_path = Path(file_path)
        images: List[Dict[str, Any]] = []
        detected_blocks = 0
        skipped_small = 0
        skipped_edge = 0
        skipped_invalid = 0
        split_generated = 0

        doc = fitz.open(file_path)
        try:
            total_pages = len(doc)
            for page_idx, page in enumerate(doc):
                if total_pages > 0:
                    self._report(page_idx / total_pages, f"Scanning page {page_idx + 1}/{total_pages}")

                try:
                    text_dict = page.get_text("dict") or {}
                    blocks = list(text_dict.get("blocks", []) or [])
                except Exception as e:
                    logger.warning(f"Could not inspect PDF image blocks on page {page_idx + 1}: {e}")
                    continue

                page_rect = page.rect
                page_area = max(page_rect.get_area(), 1.0)
                margin_x = page_rect.width * max(edge_margin_pct, 0.0) / 100.0
                margin_y = page_rect.height * max(edge_margin_pct, 0.0) / 100.0
                page_extract_count = 0

                for block in blocks:
                    if block.get("type") != 1:
                        continue
                    detected_blocks += 1

                    bbox_raw = block.get("bbox")
                    if not bbox_raw or len(bbox_raw) != 4:
                        skipped_invalid += 1
                        continue

                    try:
                        bbox = fitz.Rect(
                            float(bbox_raw[0]),
                            float(bbox_raw[1]),
                            float(bbox_raw[2]),
                            float(bbox_raw[3]),
                        )
                    except Exception:
                        skipped_invalid += 1
                        continue

                    bbox = bbox & page_rect
                    if bbox.is_empty or bbox.width < 8 or bbox.height < 8:
                        skipped_invalid += 1
                        continue

                    intrinsic_width = int(block.get("width", 0) or 0)
                    intrinsic_height = int(block.get("height", 0) or 0)
                    coverage_pct = (bbox.get_area() / page_area) * 100.0

                    if intrinsic_width < int(min_width_px) or intrinsic_height < int(min_height_px):
                        skipped_small += 1
                        continue
                    if coverage_pct < float(min_page_coverage_pct):
                        skipped_small += 1
                        continue

                    if ignore_edge_decorations:
                        near_top = bbox.y1 <= margin_y
                        near_bottom = bbox.y0 >= (page_rect.height - margin_y)
                        near_left = bbox.x1 <= margin_x
                        near_right = bbox.x0 >= (page_rect.width - margin_x)
                        banner_like = bbox.height <= (page_rect.height * 0.18)
                        icon_like = bbox.width <= (page_rect.width * 0.18)
                        if (near_top or near_bottom or near_left or near_right) and (banner_like or icon_like):
                            skipped_edge += 1
                            continue

                    try:
                        pix = page.get_pixmap(
                            clip=bbox,
                            matrix=fitz.Matrix(float(render_scale), float(render_scale)),
                            alpha=False,
                        )
                        try:
                            image_bytes = pix.tobytes("jpg")
                        except Exception:
                            image_bytes = pix.tobytes("jpeg")
                    except Exception as e:
                        logger.debug(f"Could not render PDF image block on page {page_idx + 1}: {e}")
                        skipped_invalid += 1
                        continue

                    split_crops: List[Dict[str, Any]] = []
                    if split_full_page_scans:
                        split_crops = self._split_scanned_page_image(
                            image_bytes,
                            parent_name=pdf_path.stem,
                            page_number=page_idx + 1,
                            image_index=page_extract_count + 1,
                            min_width_px=int(min_width_px),
                            min_height_px=int(min_height_px),
                            min_crop_coverage_pct=float(split_min_crop_coverage_pct),
                        )

                    if len(split_crops) >= 2:
                        split_generated += len(split_crops)
                        for crop in split_crops:
                            crop["bbox"] = [round(float(v), 2) for v in (bbox.x0, bbox.y0, bbox.x1, bbox.y1)]
                            crop["source_pdf"] = pdf_path.name
                            crop["intrinsic_width"] = intrinsic_width
                            crop["intrinsic_height"] = intrinsic_height
                            crop["coverage_pct"] = round(coverage_pct, 2)
                            images.append(crop)
                        continue

                    page_extract_count += 1
                    out_name = f"{pdf_path.stem}_page{page_idx + 1:03d}_image{page_extract_count:02d}.jpg"
                    images.append(
                        {
                            "file_name": out_name,
                            "bytes": image_bytes,
                            "page": page_idx + 1,
                            "bbox": [round(float(v), 2) for v in (bbox.x0, bbox.y0, bbox.x1, bbox.y1)],
                            "source_pdf": pdf_path.name,
                            "intrinsic_width": intrinsic_width,
                            "intrinsic_height": intrinsic_height,
                            "coverage_pct": round(coverage_pct, 2),
                        }
                    )

            self._report(1.0, "PDF image scan complete")
        finally:
            doc.close()

        message = (
            f"Extracted {len(images)} image(s) from {pdf_path.name}. "
            f"Detected {detected_blocks} image block(s); "
            f"skipped {skipped_small} small block(s), {skipped_edge} edge decoration(s), "
            f"and {skipped_invalid} invalid block(s). "
            f"Generated {split_generated} split crop(s)."
        )

        return {
            "success": bool(images),
            "source_pdf": str(pdf_path),
            "images": images,
            "pages_scanned": total_pages if 'total_pages' in locals() else 0,
            "detected_blocks": detected_blocks,
            "skipped_small": skipped_small,
            "skipped_edge": skipped_edge,
            "skipped_invalid": skipped_invalid,
            "split_generated": split_generated,
            "message": message,
        }

    @staticmethod
    def _merge_crop_boxes(boxes: List[Tuple[int, int, int, int]]) -> List[Tuple[int, int, int, int]]:
        """Merge overlapping / nearly-touching crop boxes."""
        if not boxes:
            return []
        merged: List[List[int]] = [list(boxes[0])]
        for x, y, w, h in boxes[1:]:
            merged_any = False
            for box in merged:
                bx, by, bw, bh = box
                pad_x = max(12, min(w, bw) // 12)
                pad_y = max(12, min(h, bh) // 12)
                overlaps = not (
                    x > bx + bw + pad_x
                    or bx > x + w + pad_x
                    or y > by + bh + pad_y
                    or by > y + h + pad_y
                )
                if overlaps:
                    nx0 = min(bx, x)
                    ny0 = min(by, y)
                    nx1 = max(bx + bw, x + w)
                    ny1 = max(by + bh, y + h)
                    box[0] = nx0
                    box[1] = ny0
                    box[2] = nx1 - nx0
                    box[3] = ny1 - ny0
                    merged_any = True
                    break
            if not merged_any:
                merged.append([x, y, w, h])
        return [tuple(box) for box in merged]

    def _split_scanned_page_image(
        self,
        image_bytes: bytes,
        parent_name: str,
        page_number: int,
        image_index: int,
        min_width_px: int,
        min_height_px: int,
        min_crop_coverage_pct: float,
    ) -> List[Dict[str, Any]]:
        """
        Try to split a large scanned-page image into separate photo crops.

        This is a heuristic: it looks for large non-white rectangular regions on
        the page image and returns crops only when it finds multiple candidates.
        """
        try:
            import cv2
            import numpy as np
        except Exception as e:
            logger.debug(f"OpenCV unavailable for PDF scan splitting: {e}")
            return []

        image_np = np.frombuffer(image_bytes, dtype=np.uint8)
        decoded = cv2.imdecode(image_np, cv2.IMREAD_COLOR)
        if decoded is None:
            return []

        img_h, img_w = decoded.shape[:2]
        if img_w < max(400, min_width_px) or img_h < max(400, min_height_px):
            return []

        gray = cv2.cvtColor(decoded, cv2.COLOR_BGR2GRAY)
        blurred = cv2.GaussianBlur(gray, (5, 5), 0)
        _, thresh = cv2.threshold(blurred, 245, 255, cv2.THRESH_BINARY_INV)

        kernel_w = max(9, (img_w // 60) | 1)
        kernel_h = max(9, (img_h // 60) | 1)
        kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (kernel_w, kernel_h))
        closed = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel, iterations=2)
        opened = cv2.morphologyEx(closed, cv2.MORPH_OPEN, kernel, iterations=1)

        contours, _ = cv2.findContours(opened, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        if not contours:
            return []

        image_area = float(img_w * img_h)
        min_crop_area = image_area * max(min_crop_coverage_pct, 0.5) / 100.0
        candidate_boxes: List[Tuple[int, int, int, int]] = []

        for contour in contours:
            x, y, w, h = cv2.boundingRect(contour)
            area = float(w * h)
            if w < max(160, int(min_width_px * 0.4)):
                continue
            if h < max(160, int(min_height_px * 0.4)):
                continue
            if area < min_crop_area:
                continue
            if area > image_area * 0.92:
                continue
            aspect = w / float(max(h, 1))
            if aspect < 0.3 or aspect > 4.0:
                continue
            candidate_boxes.append((x, y, w, h))

        if len(candidate_boxes) < 2:
            return []

        candidate_boxes.sort(key=lambda box: (box[1], box[0]))
        candidate_boxes = self._merge_crop_boxes(candidate_boxes)
        if len(candidate_boxes) < 2:
            return []

        results: List[Dict[str, Any]] = []
        for crop_idx, (x, y, w, h) in enumerate(candidate_boxes, 1):
            pad_x = max(8, w // 40)
            pad_y = max(8, h // 40)
            x0 = max(0, x - pad_x)
            y0 = max(0, y - pad_y)
            x1 = min(img_w, x + w + pad_x)
            y1 = min(img_h, y + h + pad_y)
            crop = decoded[y0:y1, x0:x1]
            if crop.size == 0:
                continue
            ok, encoded = cv2.imencode(".jpg", crop, [int(cv2.IMWRITE_JPEG_QUALITY), 92])
            if not ok:
                continue
            results.append(
                {
                    "file_name": (
                        f"{parent_name}_page{page_number:03d}_image{image_index:02d}_crop{crop_idx:02d}.jpg"
                    ),
                    "bytes": encoded.tobytes(),
                    "page": page_number,
                    "split_from_scan": True,
                    "split_bbox_px": [int(x0), int(y0), int(x1), int(y1)],
                }
            )

        return results if len(results) >= 2 else []

    def _extract_pdf_page_structured_lines(self, page) -> List[str]:
        """Extract page text while preserving visual structure better than plain text dump."""
        lines_out: List[str] = []
        try:
            text_dict = page.get_text("dict")
        except Exception as e:
            logger.debug(f"Structured PDF extraction failed on page {page.number + 1}: {e}")
            return lines_out

        # Gather page-level font stats so we can detect heading-like lines.
        font_sizes: List[float] = []
        blocks = text_dict.get("blocks", []) if isinstance(text_dict, dict) else []
        for block in blocks:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    try:
                        size = float(span.get("size", 0))
                    except Exception:
                        size = 0
                    if size > 0:
                        font_sizes.append(size)

        base_font = statistics.median(font_sizes) if font_sizes else 10.0
        heading_threshold = max(base_font * 1.25, base_font + 1.5)

        for block in blocks:
            if block.get("type") != 0:
                continue
            block_lines_added = 0
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue
                text = "".join((span.get("text", "") or "") for span in spans).strip()
                if not text:
                    continue

                max_font = max(float(span.get("size", 0) or 0) for span in spans)
                is_bold = any("bold" in str(span.get("font", "")).lower() for span in spans)
                looks_heading = (
                    (max_font >= heading_threshold and len(text) <= 120)
                    or (is_bold and len(text) <= 80 and text[0].isupper())
                    or re.match(r"^[A-Z][A-Z0-9 &/\-]{4,}$", text) is not None
                )

                # Keep numbering / bullets as markdown list items.
                if re.match(r"^[\u2022\u25CF\u25AA]\s*", text):
                    text = "- " + re.sub(r"^[\u2022\u25CF\u25AA]\s*", "", text)
                elif re.match(r"^\d+\.\s+", text):
                    pass
                elif looks_heading:
                    text = f"### {text}"

                lines_out.append(text)
                block_lines_added += 1

            if block_lines_added > 0:
                lines_out.append("")

        return lines_out

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def textify_docx(self, file_path: str) -> str:
        """Convert a DOCX to Markdown."""
        from docx import Document as DocxDocument
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        import io

        doc = DocxDocument(file_path)
        md_parts: List[str] = []

        total_paras = len(doc.paragraphs)
        self._report(0.0, "Processing paragraphs...")

        for p_idx, para in enumerate(doc.paragraphs):
            if total_paras > 0 and p_idx % 20 == 0:
                self._report(p_idx / total_paras * 0.5, f"Paragraph {p_idx}/{total_paras}")
            style_name = (para.style.name or "").lower()
            text = para.text.strip()
            if not text:
                md_parts.append("")
                continue
            if "heading 1" in style_name:
                md_parts.append(f"# {text}")
            elif "heading 2" in style_name:
                md_parts.append(f"## {text}")
            elif "heading 3" in style_name:
                md_parts.append(f"### {text}")
            elif "heading" in style_name:
                md_parts.append(f"#### {text}")
            elif "list" in style_name:
                md_parts.append(f"- {text}")
            else:
                md_parts.append(text)
            md_parts.append("")

        self._report(0.5, "Processing tables...")
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                md_parts.append(f"**Table {table_idx + 1}:**")
                md_parts.append(self.table_to_markdown(rows))
                md_parts.append("")

        # Count images first for progress
        image_rels = [r for r in doc.part.rels.values() if "image" in r.reltype]
        total_images = len(image_rels)
        self._report(0.7, f"Processing {total_images} embedded images...")

        img_idx = 0
        for rel_i, rel in enumerate(image_rels):
            try:
                image_bytes = rel.target_part.blob
                self._report(0.7 + 0.3 * (rel_i / max(total_images, 1)),
                             f"Describing image {rel_i + 1}/{total_images}")
                desc = self.describe_image(image_bytes)
                img_idx += 1
                md_parts.append(f"> **[Image {img_idx}]**: {desc}")
                md_parts.append("")
            except Exception as e:
                logger.debug(f"Could not extract DOCX image: {e}")

        return self._normalize_markdown_output("\n".join(md_parts))

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def textify_pptx(self, file_path: str) -> str:
        """Convert a PPTX to Markdown."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        md_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            self._report((slide_num - 1) / total_slides, f"Slide {slide_num}/{total_slides}")
            md_parts.append(f"## Slide {slide_num}\n")

            img_count = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            md_parts.append(text)
                    md_parts.append("")

                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        md_parts.append(self.table_to_markdown(rows))
                        md_parts.append("")

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_count += 1
                        self._report(
                            (slide_num - 1) / total_slides,
                            f"Slide {slide_num}/{total_slides} — describing image {img_count}",
                        )
                        image_bytes = shape.image.blob
                        desc = self.describe_image(image_bytes)
                        md_parts.append(f"> **[Image]**: {desc}")
                        md_parts.append("")
                    except Exception as e:
                        logger.debug(f"Could not extract PPTX image: {e}")

            md_parts.append("---\n")

        self._report(1.0, "PPTX conversion complete")

        return self._normalize_markdown_output("\n".join(md_parts))

    # ------------------------------------------------------------------
    # Image files (PNG, JPG, etc.)
    # ------------------------------------------------------------------

    def textify_image(self, file_path: str) -> str:
        """Convert an image file to Markdown with a VLM description."""
        self._report(0.0, "Reading image...")
        file_name = Path(file_path).name

        with open(file_path, "rb") as f:
            image_bytes = f.read()

        md_parts = [f"# {file_name}\n"]

        self._report(0.3, "Describing image with vision model...")
        description = self.describe_image(image_bytes)
        md_parts.append(f"> **[Image]**: {description}")
        md_parts.append("")

        # Try OCR extraction via pytesseract if available
        try:
            from PIL import Image
            import pytesseract
            self._report(0.7, "Extracting text via OCR...")
            img = Image.open(file_path)
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text:
                md_parts.append("## Extracted Text (OCR)\n")
                md_parts.append(ocr_text)
                md_parts.append("")
        except ImportError:
            logger.debug("pytesseract not available — skipping OCR")
        except Exception as e:
            logger.debug(f"OCR extraction failed: {e}")

        self._report(1.0, "Image conversion complete")
        return "\n".join(md_parts)

    # ------------------------------------------------------------------
    # Photo Keywords — describe, extract keywords, write EXIF
    # ------------------------------------------------------------------

    # Text models for keyword extraction (VLMs need images, so use a text LLM)
    TEXT_MODELS = ["mistral:latest", "mistral-small3.2", "llama3:latest", "gemma:latest"]
    CLEANUP_MODELS = ["qwen2.5:32b", "qwen2.5:14b", "mistral:latest", "llama3:latest"]

    def extract_keywords(self, description: str) -> List[str]:
        """Extract flat keywords from an image description using a text LLM.

        Uses a text model (Mistral etc.) rather than the VLM, because VLMs
        like Qwen3-VL return empty responses for text-only prompts.
        """
        # Don't hallucinate keywords from empty/failed descriptions
        if not description or description.startswith("[Image:"):
            logger.warning("No valid description available — skipping keyword extraction")
            return []
        try:
            import ollama
            client = ollama.Client(timeout=30)

            # Find an available text model
            text_model = None
            for model in self.TEXT_MODELS:
                try:
                    client.show(model)
                    text_model = model
                    break
                except Exception:
                    continue

            if not text_model:
                logger.warning("No text model available for keyword extraction")
                return self._extract_keywords_simple(description)

            response = client.chat(
                model=text_model,
                messages=[{
                    "role": "user",
                    "content": (
                        "Extract 10-15 photo tags from this image description. "
                        "Return ONLY a comma-separated list. Each tag should be 1-2 "
                        "simple lowercase words (no underscores, no hyphens). "
                        "Good tags: specific subjects, species, colours, season, "
                        "location type, weather. "
                        "Do NOT include photography jargon (bokeh, depth of field, "
                        "backlit, composition, close up) or vague words (atmosphere, "
                        "mood, scene, tones). Maximum 15 tags.\n\n"
                        f"Description: {description}"
                    ),
                }],
                options={"temperature": 0.1, "num_predict": 200},
            )
            raw = response["message"]["content"].strip()
            # Parse comma-separated keywords, clean up
            keywords = [k.strip().lower().strip('"\'') for k in raw.split(",")]
            # Replace underscores with spaces
            keywords = [k.replace("_", " ") for k in keywords]
            # Filter out photography jargon and empty/too-long entries
            _jargon = {
                "bokeh", "bokeh effect", "depth of field", "shallow depth of field",
                "backlit", "backlighting", "side lit", "composition", "close up",
                "blurred background", "out of focus", "warm tones", "cool tones",
                "atmosphere", "mood", "ambiance", "scene", "tones", "texture",
                "glossy texture", "gradient",
            }
            keywords = [
                k for k in keywords
                if k and len(k) > 1 and len(k) < 50 and k not in _jargon
            ]
            # Cap at 15
            keywords = keywords[:15]
            if keywords:
                return keywords
            # If LLM returned something unparseable, fall back
            return self._extract_keywords_simple(description)
        except Exception as e:
            logger.warning(f"LLM keyword extraction failed: {e}")
            return self._extract_keywords_simple(description)

    def _semantic_markdown_cleanup(self, markdown_text: str) -> str:
        """Use a local text model to improve fallback markdown structure."""
        content = (markdown_text or "").strip()
        if not content:
            return markdown_text
        self._last_cleanup_applied = False

        preferred = os.getenv("CORTEX_TEXTIFIER_CLEANUP_MODEL", "").strip()
        if self.cleanup_model:
            preferred = self.cleanup_model
        candidates: List[str] = [preferred] if preferred else []
        candidates.extend(self.CLEANUP_MODELS)
        dedup_candidates: List[str] = []
        seen_models = set()
        for candidate in candidates:
            if not candidate or candidate in seen_models:
                continue
            dedup_candidates.append(candidate)
            seen_models.add(candidate)
        candidates = dedup_candidates

        max_chars = int(os.getenv("CORTEX_TEXTIFIER_CLEANUP_MAX_CHARS", "18000"))
        payload = content[:max_chars]
        prompt = (
            "Rewrite this extracted markdown to improve readability while preserving facts.\n"
            "Rules:\n"
            "- Keep existing page headings like '## Page N'.\n"
            "- Preserve wording and meaning; do not invent content.\n"
            "- Split run-on lines into proper paragraphs.\n"
            "- Convert obvious enumerations into markdown lists.\n"
            "- Return markdown only.\n\n"
            f"{payload}"
        )

        # Path 1: shared LLM interface (supports LM Studio and Ollama via env config).
        for model_name in candidates:
            if not model_name:
                continue
            try:
                from cortex_engine.llm_interface import LLMInterface
                provider_override = self.cleanup_provider or None
                llm_kwargs: Dict[str, Any] = {
                    "model": model_name,
                    "temperature": 0.0,
                    "request_timeout": 90.0,
                }
                if provider_override:
                    llm_kwargs["provider"] = provider_override
                llm = LLMInterface(**llm_kwargs)
                result = llm.generate(prompt=prompt, system_prompt=None, max_tokens=4096).strip()
                result = re.sub(r"<think>.*?</think>", "", result, flags=re.DOTALL).strip()
                if result:
                    logger.info(f"Applied semantic markdown cleanup via LLMInterface model={model_name}")
                    self._last_cleanup_applied = True
                    return result
            except Exception:
                if self._is_connection_error_exception():
                    return markdown_text
                continue

        # Path 2: direct Ollama fallback with model availability probing.
        if self.cleanup_provider and self.cleanup_provider != "ollama":
            logger.info(f"Semantic cleanup skipped (provider '{self.cleanup_provider}' unavailable)")
            return markdown_text
        try:
            import ollama
            client = ollama.Client(timeout=90)
        except Exception as e:
            logger.info(f"Semantic cleanup skipped (ollama unavailable): {e}")
            return markdown_text

        cleanup_model = None
        for model in candidates:
            if not model:
                continue
            try:
                client.show(model)
                cleanup_model = model
                break
            except Exception:
                continue

        if not cleanup_model:
            logger.info("Semantic cleanup skipped (no cleanup model available)")
            return markdown_text

        try:
            response = client.chat(
                model=cleanup_model,
                messages=[{"role": "user", "content": prompt}],
                options={"temperature": 0.0, "num_predict": 4096},
            )
            result = (response.get("message", {}) or {}).get("content", "").strip()
            result = re.sub(r"<think>.*?</think>", "", result, flags=re.DOTALL).strip()
            if result:
                logger.info(f"Applied semantic markdown cleanup with {cleanup_model}")
                self._last_cleanup_applied = True
                return result
            return markdown_text
        except Exception as e:
            logger.info(f"Semantic cleanup failed ({cleanup_model}): {e}")
            return markdown_text

    @staticmethod
    def _is_connection_error_exception() -> bool:
        """Best-effort check for endpoint/network availability errors from provider clients."""
        import traceback
        exc_text = traceback.format_exc().lower()
        markers = [
            "connection error",
            "failed to connect",
            "connection refused",
            "timed out",
            "max retries exceeded",
            "name or service not known",
            "couldn't connect to server",
        ]
        return any(marker in exc_text for marker in markers)

    def _try_docling_subprocess(
        self,
        file_path: str,
        timeout_seconds: float,
    ) -> Optional[Tuple[str, List[Dict[str, Any]]]]:
        """
        Run Docling conversion in a subprocess so timeout can terminate work cleanly.
        Returns (markdown_text, figures) on success, else None.
        """
        import json
        import sys

        root_dir = str(Path(__file__).resolve().parents[1])
        with tempfile.NamedTemporaryFile(prefix="cortex_docling_", suffix=".json", delete=False) as tf:
            output_json_path = tf.name
        script = "\n".join(
            [
                "import json, sys",
                "from cortex_engine.docling_reader import DoclingDocumentReader",
                "fp = sys.argv[1]",
                "outp = sys.argv[2]",
                "reader = DoclingDocumentReader()",
                "payload = {'ok': False, 'text': '', 'figures': []}",
                "if reader.is_available and reader.can_process_file(fp):",
                "    docs = reader.load_data(fp)",
                "    if docs and docs[0].text and docs[0].text.strip():",
                "        md = docs[0].metadata or {}",
                "        payload = {",
                "            'ok': True,",
                "            'text': docs[0].text,",
                "            'figures': (md.get('docling_figures', []) or []),",
                "        }",
                "with open(outp, 'w', encoding='utf-8') as f:",
                "    json.dump(payload, f)",
            ]
        )
        try:
            proc = subprocess.run(
                [sys.executable, "-c", script, str(file_path), str(output_json_path)],
                cwd=root_dir,
                capture_output=True,
                text=True,
                timeout=timeout_seconds,
            )
            if proc.returncode != 0:
                logger.info(f"Docling subprocess unavailable for {Path(file_path).name}: {proc.stderr.strip()}")
                return None
            payload = json.loads(Path(output_json_path).read_text(encoding="utf-8"))
            if not payload.get("ok"):
                return None
            text = str(payload.get("text") or "")
            figures = payload.get("figures") or []
            if not text.strip():
                return None
            return text, figures
        except subprocess.TimeoutExpired:
            logger.warning(
                f"Docling timed out after {timeout_seconds}s for {Path(file_path).name}; "
                "falling back to legacy reader"
            )
            return None
        except Exception as e:
            logger.info(f"Docling subprocess failed for {Path(file_path).name}: {e}")
            return None
        finally:
            try:
                Path(output_json_path).unlink(missing_ok=True)
            except Exception:
                pass

    def _enrich_docling_image_markers(
        self,
        markdown_text: str,
        file_path: str,
        figures: Optional[List[Dict[str, Any]]] = None,
    ) -> str:
        """
        Replace Docling '<!-- image -->' markers with VLM descriptions.
        Uses figure bbox/page metadata to crop images directly from the PDF when possible.
        """
        text = markdown_text or ""
        marker = "<!-- image -->"
        marker_count = text.count(marker)
        if marker_count == 0:
            return text

        if not self.use_vision:
            return text.replace(marker, "> **[Image]**: [Image: vision model disabled]")

        figures = figures or []
        descriptions: List[str] = []
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            try:
                enrich_started = time.monotonic()
                page_desc_cache: Dict[int, str] = {}

                def describe_page(page_idx: int) -> str:
                    if page_idx in page_desc_cache:
                        return page_desc_cache[page_idx]
                    page_idx = min(max(page_idx, 0), len(doc) - 1)
                    page = doc[page_idx]
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                    desc = self._describe_image_with_timeout(pix.tobytes("png"))
                    page_desc_cache[page_idx] = desc
                    return desc

                for i in range(marker_count):
                    elapsed = time.monotonic() - enrich_started
                    if self.image_enrich_max_seconds > 0 and elapsed > self.image_enrich_max_seconds:
                        descriptions.extend(
                            ["[Image: skipped to keep processing responsive]"] * (marker_count - i)
                        )
                        break
                    figure = figures[i] if i < len(figures) else {}
                    desc = figure.get("vlm_description") or ""
                    if desc:
                        descriptions.append(desc)
                        continue

                    page_raw = figure.get("page")
                    bbox = figure.get("bbox")
                    # Fallback path when Docling doesn't expose figure-level metadata:
                    # estimate image context from page index aligned to marker order.
                    if page_raw is None:
                        if marker_count > 1:
                            ratio = i / float(marker_count - 1)
                            page_idx = int(round(ratio * (len(doc) - 1)))
                        else:
                            page_idx = 0
                        descriptions.append(describe_page(page_idx))
                        continue
                    if not bbox or len(bbox) != 4:
                        page_idx = int(page_raw) - 1 if int(page_raw) >= 1 else int(page_raw)
                        descriptions.append(describe_page(page_idx))
                        continue

                    try:
                        page_idx = int(page_raw)
                        # Docling page indices are typically 1-based.
                        if page_idx >= 1:
                            page_idx -= 1
                    except Exception:
                        descriptions.append("[Image: invalid figure page metadata]")
                        continue

                    if page_idx < 0 or page_idx >= len(doc):
                        descriptions.append("[Image: figure page out of range]")
                        continue

                    page = doc[page_idx]
                    try:
                        rect = fitz.Rect(float(bbox[0]), float(bbox[1]), float(bbox[2]), float(bbox[3]))
                    except Exception:
                        descriptions.append("[Image: invalid figure bounds]")
                        continue

                    # Defensive clamp/cleanup of rect.
                    rect = rect & page.rect
                    if rect.is_empty or rect.width < 4 or rect.height < 4:
                        descriptions.append(describe_page(page_idx))
                        continue

                    pix = page.get_pixmap(clip=rect, matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                    desc = self._describe_image_with_timeout(pix.tobytes("png"))
                    descriptions.append(desc or describe_page(page_idx))
            finally:
                doc.close()
        except Exception as e:
            logger.warning(f"Docling image marker enrichment failed: {e}")
            descriptions = ["[Image: could not be described]"] * marker_count

        parts = text.split(marker)
        out: List[str] = [parts[0]]
        for i in range(marker_count):
            desc = descriptions[i] if i < len(descriptions) else "[Image: description unavailable]"
            out.append(f"> **[Image {i + 1}]**: {desc}\n")
            out.append(parts[i + 1] if i + 1 < len(parts) else "")
        return "".join(out)

    @staticmethod
    def _extract_keywords_simple(description: str) -> List[str]:
        """Fallback keyword extraction using simple NLP."""
        import re
        # Remove bracketed placeholders
        text = re.sub(r'\[.*?\]', '', description).lower()
        # Common stop words
        stop = {
            'a', 'an', 'the', 'is', 'are', 'was', 'were', 'be', 'been', 'being',
            'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could',
            'should', 'may', 'might', 'can', 'shall', 'of', 'in', 'to', 'for',
            'with', 'on', 'at', 'from', 'by', 'about', 'as', 'into', 'through',
            'and', 'but', 'or', 'so', 'if', 'that', 'this', 'it', 'its', 'not',
            'what', 'which', 'who', 'how', 'there', 'their', 'they', 'them',
            'image', 'shows', 'appears', 'visible', 'seen', 'also', 'very',
        }
        words = re.findall(r'[a-z]+', text)
        keywords = []
        seen = set()
        for w in words:
            if w not in stop and len(w) > 2 and w not in seen:
                seen.add(w)
                keywords.append(w)
        return keywords[:20]

    @staticmethod
    def read_exif_keywords(file_path: str) -> List[str]:
        """Read existing XMP Subject / IPTC Keywords from image."""
        import shutil
        import json
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return []
        try:
            result = subprocess.run(
                [exiftool_path, "-json", "-XMP-dc:Subject", "-IPTC:Keywords",
                 file_path],
                capture_output=True, text=True, timeout=10,
            )
            if result.returncode != 0:
                return []
            data = json.loads(result.stdout)
            if not data:
                return []
            existing = set()
            for field in ("Subject", "Keywords"):
                val = data[0].get(field, [])
                if isinstance(val, str):
                    val = [val]
                for v in val:
                    existing.add(v.strip().lower())
            return sorted(existing)
        except Exception as e:
            logger.warning(f"Read existing keywords failed: {e}")
            return []

    @staticmethod
    def write_exif_keywords(file_path: str, keywords: List[str],
                            description: str = "") -> Dict[str, any]:
        """Write keywords and description to image EXIF/XMP using exiftool.

        Merges with existing keywords (uses += to append, not replace).

        Writes to:
          - XMP:dc:subject (Keywords — read by Lightroom, Bridge, etc.)
          - IPTC:Keywords (broader compatibility)
          - XMP:dc:description (Caption)
          - EXIF:ImageDescription (legacy caption)

        Returns dict with 'success', 'message', and 'keywords_written'.
        """
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return {
                "success": False,
                "message": "exiftool not found on system PATH",
                "keywords_written": 0,
            }

        cmd = [exiftool_path, "-overwrite_original"]

        # Use += to ADD to existing keywords rather than replacing them.
        # XMP-dc:Subject is the keyword list that Lightroom Classic,
        # Bridge, Capture One, and other DAMs read as "Keywords".
        # IPTC:Keywords provides legacy compatibility.
        for kw in keywords:
            cmd.append(f"-XMP-dc:Subject+={kw}")
            cmd.append(f"-IPTC:Keywords+={kw}")

        # Write description/caption to the correct fields:
        # - XMP-dc:Description → LRC "Caption" in metadata panel
        # - IPTC:Caption-Abstract → legacy caption
        # - EXIF:ImageDescription → EXIF-level caption
        if description:
            caption = description[:2000]
            cmd.append(f"-XMP-dc:Description={caption}")
            cmd.append(f"-IPTC:Caption-Abstract={caption}")
            cmd.append(f"-EXIF:ImageDescription={caption}")

        cmd.append(file_path)

        try:
            result = subprocess.run(
                cmd, capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0:
                logger.info(f"Wrote {len(keywords)} keywords to {Path(file_path).name}")
                return {
                    "success": True,
                    "message": result.stdout.strip(),
                    "keywords_written": len(keywords),
                }
            else:
                logger.warning(f"exiftool error: {result.stderr}")
                return {
                    "success": False,
                    "message": result.stderr.strip(),
                    "keywords_written": 0,
                }
        except subprocess.TimeoutExpired:
            return {"success": False, "message": "exiftool timed out", "keywords_written": 0}
        except Exception as e:
            return {"success": False, "message": str(e), "keywords_written": 0}

    @staticmethod
    def clear_exif_keywords(file_path: str) -> bool:
        """Remove all existing XMP Subject, IPTC Keywords, and description fields."""
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return False
        try:
            result = subprocess.run(
                [exiftool_path, "-overwrite_original",
                 "-XMP-dc:Subject=", "-IPTC:Keywords=",
                 "-XMP-dc:Description=", "-IPTC:Caption-Abstract=",
                 "-EXIF:ImageDescription=",
                 file_path],
                capture_output=True, text=True, timeout=10,
            )
            return result.returncode == 0
        except Exception as e:
            logger.warning(f"Clear keywords failed: {e}")
            return False

    @staticmethod
    def clear_exif_location(file_path: str) -> bool:
        """Remove existing IPTC/XMP location fields (Country, State, City)."""
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return False
        try:
            result = subprocess.run(
                [exiftool_path, "-overwrite_original",
                 "-IPTC:Country-PrimaryLocationName=",
                 "-XMP-photoshop:Country=",
                 "-IPTC:Province-State=",
                 "-XMP-photoshop:State=",
                 "-IPTC:City=",
                 "-XMP-photoshop:City=",
                 file_path],
                capture_output=True, text=True, timeout=10,
            )
            return result.returncode == 0
        except Exception as e:
            logger.warning(f"Clear location failed: {e}")
            return False

    @staticmethod
    def read_exif_location(file_path: str) -> Dict[str, str]:
        """Read existing IPTC/XMP location fields from an image."""
        import shutil

        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return {"country": "", "state": "", "city": ""}
        try:
            import json

            result = subprocess.run(
                [
                    exiftool_path,
                    "-json",
                    "-XMP-photoshop:City",
                    "-IPTC:City",
                    "-XMP-photoshop:State",
                    "-IPTC:Province-State",
                    "-XMP-photoshop:Country",
                    "-IPTC:Country-PrimaryLocationName",
                    file_path,
                ],
                capture_output=True,
                text=True,
                timeout=10,
            )
            if result.returncode != 0:
                return {"country": "", "state": "", "city": ""}
            payload = json.loads(result.stdout)
            if not payload:
                return {"country": "", "state": "", "city": ""}
            row = payload[0]
            return {
                "country": (row.get("Country") or row.get("Country-PrimaryLocationName") or "").strip(),
                "state": (row.get("State") or row.get("Province-State") or "").strip(),
                "city": (row.get("City") or "").strip(),
            }
        except Exception as e:
            logger.warning(f"Location read failed for {file_path}: {e}")
            return {"country": "", "state": "", "city": ""}

    @staticmethod
    def read_gps(file_path: str) -> Optional[Tuple[float, float]]:
        """Read GPS coordinates from image EXIF. Returns (lat, lon) or None."""
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return None
        try:
            import json
            result = subprocess.run(
                [exiftool_path, "-json", "-n", "-GPSLatitude", "-GPSLongitude",
                 file_path],
                capture_output=True, text=True, timeout=10,
            )
            if result.returncode != 0:
                return None
            data = json.loads(result.stdout)
            if not data:
                return None
            lat = data[0].get("GPSLatitude")
            lon = data[0].get("GPSLongitude")
            if lat is not None and lon is not None:
                return (float(lat), float(lon))
        except Exception as e:
            logger.warning(f"GPS read failed for {file_path}: {e}")
        return None

    @staticmethod
    def _merge_location_fields(*locations: Optional[Dict[str, str]]) -> Dict[str, str]:
        """Merge location dicts, keeping the first non-empty value for each field."""
        merged = {"country": "", "state": "", "city": ""}
        for location in locations:
            if not location:
                continue
            for field in ("country", "state", "city"):
                value = (location.get(field, "") or "").strip()
                if value and not merged[field]:
                    merged[field] = value
        return merged

    @staticmethod
    def reverse_geocode(lat: float, lon: float) -> Dict[str, str]:
        """Reverse-geocode GPS coordinates to country/state/city.

        Returns dict with 'country', 'state', 'city' keys (empty string if
        not found).
        """
        try:
            from geopy.geocoders import Nominatim
            g = Nominatim(user_agent="cortex_suite", timeout=10)
            loc = g.reverse((lat, lon), exactly_one=True, language="en")
            if loc is None:
                return {"country": "", "state": "", "city": ""}
            addr = loc.raw.get("address", {})
            country = addr.get("country", "")
            state = addr.get("state", "") or addr.get("region", "")
            city = (addr.get("city", "")
                    or addr.get("town", "")
                    or addr.get("village", "")
                    or addr.get("suburb", ""))
            return {"country": country, "state": state, "city": city}
        except Exception as e:
            logger.warning(f"Reverse geocode failed for ({lat}, {lon}): {e}")
            return {"country": "", "state": "", "city": ""}

    @staticmethod
    def geocode_location_hint(city: str = "", state: str = "", country: str = "") -> Optional[Tuple[float, float]]:
        """Resolve coordinates from a city/state/country hint."""
        city = (city or "").strip()
        state = (state or "").strip()
        country = (country or "").strip()
        if not any([city, state, country]):
            return None
        try:
            from geopy.geocoders import Nominatim

            geocoder = Nominatim(user_agent="cortex_suite", timeout=10)
            queries: List[str] = []

            if city:
                queries.append(", ".join(part for part in [city, state, country] if part))
                if state and not country:
                    queries.append(", ".join(part for part in [city, state] if part))
                if country and not state:
                    queries.append(", ".join(part for part in [city, country] if part))
                queries.append(city)
            elif state:
                queries.append(", ".join(part for part in [state, country] if part))
                queries.append(state)

            if country:
                if not city:
                    queries.append(f"capital of {country}")
                    queries.append(f"{country} capital city")
                queries.append(country)

            seen = set()
            for query in queries:
                normalized = query.strip()
                if not normalized or normalized.lower() in seen:
                    continue
                seen.add(normalized.lower())
                loc = geocoder.geocode(normalized, exactly_one=True, language="en")
                if loc is not None:
                    return (float(loc.latitude), float(loc.longitude))
        except Exception as e:
            logger.warning(
                "Forward geocode failed for city=%r state=%r country=%r: %s",
                city,
                state,
                country,
                e,
            )
        return None

    def resolve_photo_location(
        self,
        file_path: str,
        fallback_city: str = "",
        fallback_country: str = "",
    ) -> Dict[str, Any]:
        """Resolve the best available location/GPS combination for a photo."""
        existing_location = self.read_exif_location(file_path)
        existing_gps = self.read_gps(file_path)

        fallback_location = {"country": "", "state": "", "city": ""}
        if not any(existing_location.values()) and not existing_gps:
            fallback_location = {
                "country": (fallback_country or "").strip(),
                "state": "",
                "city": (fallback_city or "").strip(),
            }

        seed_location = self._merge_location_fields(existing_location, fallback_location)
        resolved_location = dict(seed_location)
        resolved_gps = existing_gps

        if resolved_gps:
            reverse = self.reverse_geocode(resolved_gps[0], resolved_gps[1])
            resolved_location = self._merge_location_fields(seed_location, reverse)
        elif any(seed_location.values()):
            derived_gps = self.geocode_location_hint(
                city=seed_location.get("city", ""),
                state=seed_location.get("state", ""),
                country=seed_location.get("country", ""),
            )
            if derived_gps:
                resolved_gps = derived_gps
                reverse = self.reverse_geocode(derived_gps[0], derived_gps[1])
                resolved_location = self._merge_location_fields(seed_location, reverse)

        return {
            "existing_location": existing_location,
            "existing_gps": existing_gps,
            "location": resolved_location,
            "gps": resolved_gps,
            "used_fallback_location": bool(any(fallback_location.values())),
        }

    @staticmethod
    def write_exif_location(file_path: str, country: str, state: str,
                            city: str) -> bool:
        """Write IPTC/XMP location fields to image using exiftool."""
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return False
        cmd = [exiftool_path, "-overwrite_original"]
        if country:
            cmd.append(f"-IPTC:Country-PrimaryLocationName={country}")
            cmd.append(f"-XMP-photoshop:Country={country}")
        if state:
            cmd.append(f"-IPTC:Province-State={state}")
            cmd.append(f"-XMP-photoshop:State={state}")
        if city:
            cmd.append(f"-IPTC:City={city}")
            cmd.append(f"-XMP-photoshop:City={city}")
        cmd.append(file_path)
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            return result.returncode == 0
        except Exception as e:
            logger.warning(f"EXIF location write failed: {e}")
            return False

    @staticmethod
    def write_gps_coordinates(file_path: str, lat: float, lon: float) -> Dict[str, Any]:
        """Write GPS coordinates to EXIF using decimal degrees."""
        import shutil

        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return {"success": False, "message": "exiftool not found on system PATH"}
        try:
            latitude = float(lat)
            longitude = float(lon)
            cmd = [
                exiftool_path,
                "-overwrite_original",
                f"-GPSLatitude={abs(latitude):.6f}",
                f"-GPSLatitudeRef={'N' if latitude >= 0 else 'S'}",
                f"-GPSLongitude={abs(longitude):.6f}",
                f"-GPSLongitudeRef={'E' if longitude >= 0 else 'W'}",
                file_path,
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                return {"success": True, "message": result.stdout.strip()}
            return {"success": False, "message": result.stderr.strip() or "GPS write failed"}
        except Exception as e:
            logger.warning(f"GPS write failed for {file_path}: {e}")
            return {"success": False, "message": str(e)}

    @staticmethod
    def write_ownership_metadata(file_path: str, ownership_notice: str) -> Dict[str, any]:
        """Write ownership/copyright metadata fields using exiftool."""
        import shutil

        notice = (ownership_notice or "").strip()
        if not notice:
            return {"success": True, "message": "No ownership notice provided", "fields_written": 0}

        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return {"success": False, "message": "exiftool not found on system PATH", "fields_written": 0}

        cmd = [
            exiftool_path,
            "-overwrite_original",
            f"-EXIF:Copyright={notice}",
            f"-IPTC:CopyrightNotice={notice}",
            f"-XMP-dc:Rights={notice}",
            "-XMP-xmpRights:Marked=True",
            file_path,
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=20)
            if result.returncode == 0:
                return {"success": True, "message": result.stdout.strip(), "fields_written": 4}
            return {"success": False, "message": result.stderr.strip(), "fields_written": 0}
        except Exception as e:
            return {"success": False, "message": str(e), "fields_written": 0}

    @staticmethod
    def _resize_image_preserving_metadata(
        file_path: str,
        max_width: int = 1920,
        max_height: int = 1080,
        convert_to_jpg: bool = False,
        jpg_quality: int = 90,
    ) -> Dict[str, any]:
        """Downsample oversized images in-place and copy metadata from original."""
        from PIL import Image
        import shutil

        suffix = Path(file_path).suffix.lower()
        format_map = {
            ".jpg": "JPEG",
            ".jpeg": "JPEG",
            ".png": "PNG",
            ".tif": "TIFF",
            ".tiff": "TIFF",
            ".webp": "WEBP",
            ".bmp": "BMP",
            ".gif": "GIF",
        }

        original_copy = None
        output_path = file_path
        try:
            with Image.open(file_path) as img:
                original_width, original_height = img.size
                image_format = img.format or format_map.get(suffix, "JPEG")

            should_convert_to_jpg = bool(convert_to_jpg and image_format != "JPEG")
            if original_width <= max_width and original_height <= max_height and not should_convert_to_jpg:
                return {
                    "resized": False,
                    "original_width": original_width,
                    "original_height": original_height,
                    "new_width": original_width,
                    "new_height": original_height,
                    "metadata_preserved": True,
                    "converted_to_jpg": False,
                    "output_path": file_path,
                }

            scale = min(max_width / float(original_width), max_height / float(original_height))
            if original_width <= max_width and original_height <= max_height:
                new_width, new_height = original_width, original_height
            else:
                new_width = max(1, int(original_width * scale))
                new_height = max(1, int(original_height * scale))

            target_format = "JPEG" if should_convert_to_jpg else image_format
            if should_convert_to_jpg:
                output_path = str(Path(file_path).with_suffix(".jpg"))

            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                original_copy = tmp.name
            shutil.copy2(file_path, original_copy)

            with Image.open(file_path) as img:
                if target_format == "JPEG" and img.mode in ("RGBA", "LA", "P"):
                    img = img.convert("RGB")
                resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                save_kwargs: Dict[str, any] = {}
                if target_format == "JPEG":
                    quality = int(max(60, min(100, jpg_quality)))
                    save_kwargs["quality"] = quality
                    save_kwargs["optimize"] = True
                resized.save(output_path, format=target_format, **save_kwargs)

            if output_path != file_path and Path(file_path).exists():
                try:
                    os.remove(file_path)
                except Exception:
                    pass

            metadata_preserved = False
            exiftool_path = shutil.which("exiftool")
            if exiftool_path:
                copy_result = subprocess.run(
                    [
                        exiftool_path,
                        "-overwrite_original",
                        "-TagsFromFile",
                        original_copy,
                        "-all:all",
                        "-unsafe",
                        output_path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=30,
                )
                metadata_preserved = copy_result.returncode == 0
                if not metadata_preserved:
                    logger.warning(
                        f"Metadata copy after resize failed for {Path(output_path).name}: {copy_result.stderr}"
                    )
            else:
                logger.warning("exiftool not found; metadata copy after resize skipped")

            return {
                "resized": (new_width != original_width or new_height != original_height),
                "original_width": original_width,
                "original_height": original_height,
                "new_width": new_width,
                "new_height": new_height,
                "metadata_preserved": metadata_preserved,
                "converted_to_jpg": should_convert_to_jpg,
                "output_path": output_path,
            }
        except Exception as e:
            logger.warning(f"Resize failed for {Path(file_path).name}: {e}")
            if original_copy and Path(original_copy).exists():
                try:
                    shutil.copy2(original_copy, file_path)
                except Exception as restore_err:
                    logger.warning(f"Could not restore original image after resize failure: {restore_err}")
            if output_path != file_path and Path(output_path).exists():
                try:
                    os.remove(output_path)
                except Exception:
                    pass
            return {
                "resized": False,
                "error": str(e),
                "metadata_preserved": False,
                "converted_to_jpg": False,
                "output_path": file_path,
            }
        finally:
            if original_copy and Path(original_copy).exists():
                try:
                    os.remove(original_copy)
                except Exception:
                    pass

    def resize_image_only(
        self,
        file_path: str,
        max_width: int = 1920,
        max_height: int = 1080,
        convert_to_jpg: bool = False,
        jpg_quality: int = 90,
    ) -> Dict[str, any]:
        """Resize photo only (no keyword generation), preserving metadata when possible."""
        self._report(0.0, "Checking image dimensions...")
        resize_info = self._resize_image_preserving_metadata(
            file_path,
            max_width=max_width,
            max_height=max_height,
            convert_to_jpg=convert_to_jpg,
            jpg_quality=jpg_quality,
        )
        output_path = resize_info.get("output_path", file_path)
        self._report(1.0, "Done")
        return {
            "file_name": Path(output_path).name,
            "output_path": output_path,
            "resize_info": resize_info,
        }

    @staticmethod
    def _filter_sensitive_keywords(
        keywords: List[str], blocked_keywords: Optional[List[str]] = None
    ) -> Tuple[List[str], List[str]]:
        """Remove blocked/sensitive keywords while preserving other tags."""
        blocked_set = {k.strip().lower() for k in (blocked_keywords or []) if k and k.strip()}
        filtered: List[str] = []
        removed: List[str] = []

        for keyword in keywords:
            kw = (keyword or "").strip()
            if not kw:
                continue
            kw_lower = kw.lower()
            is_handle = kw_lower.startswith("@") or ("_" in kw_lower)
            if kw_lower in blocked_set or is_handle:
                removed.append(kw_lower)
                continue
            filtered.append(kw_lower)

        dedup_filtered = list(dict.fromkeys(filtered))
        dedup_removed = list(dict.fromkeys(removed))
        return dedup_filtered, dedup_removed

    def _llm_filter_sensitive_keywords(self, keywords: List[str]) -> List[str]:
        """Use a text LLM to remove personal/sensitive tags while keeping neutral content tags."""
        if not keywords:
            return []
        try:
            import ollama
            client = ollama.Client(timeout=30)

            text_model = None
            for model in self.TEXT_MODELS:
                try:
                    client.show(model)
                    text_model = model
                    break
                except Exception:
                    continue

            if not text_model:
                logger.warning("No text model available for keyword anonymization; using deterministic filtering only")
                return keywords

            payload = ", ".join(keywords)
            response = client.chat(
                model=text_model,
                messages=[{
                    "role": "user",
                    "content": (
                        "You are cleaning photo keywords for privacy. "
                        "Remove personal names, usernames/handles, nicknames, and relationship words "
                        "like friends/family/partner/mum/dad. "
                        "Keep non-person subject/location/object tags unchanged. "
                        "Return ONLY a comma-separated list of kept keywords in lowercase.\n\n"
                        f"Keywords: {payload}"
                    ),
                }],
                options={"temperature": 0.0, "num_predict": 200},
            )
            raw = (response.get("message", {}) or {}).get("content", "").strip()
            if not raw:
                return keywords
            kept = [k.strip().lower() for k in raw.split(",") if k.strip()]
            return list(dict.fromkeys(kept))
        except Exception as e:
            logger.warning(f"LLM keyword anonymization failed: {e}")
            return keywords

    def _hybrid_filter_sensitive_keywords(
        self, keywords: List[str], blocked_keywords: Optional[List[str]] = None
    ) -> Tuple[List[str], List[str]]:
        """Hybrid privacy filter: LLM intent filter + deterministic blocked-list/handle filter."""
        llm_kept_keywords = self._llm_filter_sensitive_keywords(keywords)
        llm_removed = [k for k in keywords if k.lower() not in {x.lower() for x in llm_kept_keywords}]
        final_kept, blocked_removed = self._filter_sensitive_keywords(
            llm_kept_keywords, blocked_keywords=blocked_keywords
        )
        removed = list(dict.fromkeys([x.lower() for x in llm_removed] + blocked_removed))
        return final_kept, removed

    @staticmethod
    def clear_exif_keyword_tags_only(file_path: str) -> bool:
        """Remove only keyword tag fields, preserving captions/descriptions."""
        import shutil
        exiftool_path = shutil.which("exiftool")
        if not exiftool_path:
            return False
        try:
            result = subprocess.run(
                [exiftool_path, "-overwrite_original", "-XMP-dc:Subject=", "-IPTC:Keywords=", file_path],
                capture_output=True,
                text=True,
                timeout=10,
            )
            return result.returncode == 0
        except Exception as e:
            logger.warning(f"Clear keyword tags failed: {e}")
            return False

    def anonymize_existing_photo_keywords(
        self, file_path: str, blocked_keywords: Optional[List[str]] = None
    ) -> Dict[str, any]:
        """Anonymize existing keyword tags on a photo without running image description."""
        existing_keywords = self.read_exif_keywords(file_path)
        filtered_existing, removed_existing = self._hybrid_filter_sensitive_keywords(
            existing_keywords, blocked_keywords=blocked_keywords
        )
        if not removed_existing:
            return {
                "success": True,
                "existing_keywords": existing_keywords,
                "kept_keywords": filtered_existing,
                "removed_keywords": [],
                "message": "No sensitive keywords found to remove",
            }

        if not self.clear_exif_keyword_tags_only(file_path):
            return {
                "success": False,
                "existing_keywords": existing_keywords,
                "kept_keywords": existing_keywords,
                "removed_keywords": [],
                "message": "Failed to clear keyword fields before anonymization",
            }
        write_result = self.write_exif_keywords(file_path, filtered_existing, "")
        return {
            "success": bool(write_result.get("success")),
            "existing_keywords": existing_keywords,
            "kept_keywords": filtered_existing,
            "removed_keywords": removed_existing,
            "message": write_result.get("message", ""),
        }

    def keyword_image(self, file_path: str, city_radius_km: float = 5.0,
                      clear_keywords: bool = False,
                      clear_location: bool = False,
                      force_resize_max_1080: bool = False,
                      generate_description: bool = True,
                      populate_location: bool = True,
                      anonymize_keywords: bool = False,
                      blocked_keywords: Optional[List[str]] = None,
                      fallback_city: str = "",
                      fallback_country: str = "",
                      ownership_notice: str = "") -> Dict[str, any]:
        """Full pipeline: describe image, extract keywords, resolve GPS location, write to EXIF.

        Args:
            file_path: Path to image file.
            city_radius_km: Reserved for future proximity grouping.
            clear_keywords: If True, strip existing keywords/description before writing.
            clear_location: If True, strip existing location fields before writing.

        Returns dict with description, keywords, location, and write result.
        """
        self._report(0.0, "Reading image...")
        file_name = Path(file_path).name
        resize_info = {
            "resized": False,
            "metadata_preserved": True,
        }

        if force_resize_max_1080:
            self._report(0.01, "Checking image dimensions...")
            resize_info = self._resize_image_preserving_metadata(
                file_path, max_width=1920, max_height=1080
            )
            if resize_info.get("resized"):
                self._report(0.06, "Downsampled to gallery size")

        # Clear existing fields if requested
        if clear_keywords:
            self._report(0.02, "Clearing existing keywords...")
            self.clear_exif_keywords(file_path)
        if clear_location:
            self._report(0.04, "Clearing existing location fields...")
            self.clear_exif_location(file_path)

        # Read existing keywords so we can merge (after any clear)
        existing_keywords = self.read_exif_keywords(file_path)
        keywords: List[str] = []
        new_keywords: List[str] = []
        description = ""
        removed_sensitive_keywords: List[str] = []
        existing_keywords_for_merge = existing_keywords
        self._report(0.1, "Resolving GPS and location metadata...")
        location = None
        gps = None
        location_result = {
            "enabled": populate_location,
            "location_written": False,
            "gps_written": False,
            "location_write_ok": True,
            "gps_write_ok": True,
        }
        location_resolution = {
            "existing_location": {"country": "", "state": "", "city": ""},
            "existing_gps": None,
            "location": {"country": "", "state": "", "city": ""},
            "gps": None,
            "used_fallback_location": False,
        }
        if populate_location:
            location_resolution = self.resolve_photo_location(
                file_path,
                fallback_city=fallback_city,
                fallback_country=fallback_country,
            )
            location = location_resolution.get("location")
            gps = location_resolution.get("gps")
            logger.info(f"Resolved location for {file_name}: {location} / {gps}")
        else:
            location = self.read_exif_location(file_path)
            gps = self.read_gps(file_path)

        has_gps = gps is not None

        if generate_description:
            with open(file_path, "rb") as f:
                image_bytes = f.read()

            self._report(0.2, "Describing image with vision model...")
            description = self.describe_image(image_bytes)

            self._report(0.5, "Extracting keywords...")
            keywords = self.extract_keywords(description)
        else:
            self._report(0.5, "Skipping AI description generation")

        if generate_description:
            if populate_location and location and any(location.values()):
                for field in ("city", "state", "country"):
                    val = location.get(field, "")
                    if val and val.lower() not in [k.lower() for k in keywords]:
                        keywords.append(val.lower())
            elif populate_location and not has_gps:
                if "nogps" not in keywords:
                    keywords.append("nogps")

            removed_existing_sensitive: List[str] = []
            if anonymize_keywords:
                # Sanitize generated keywords and existing metadata keywords.
                keywords, removed_generated_sensitive = self._hybrid_filter_sensitive_keywords(
                    keywords, blocked_keywords=blocked_keywords
                )
                existing_keywords_for_merge, removed_existing_sensitive = self._hybrid_filter_sensitive_keywords(
                    existing_keywords, blocked_keywords=blocked_keywords
                )
                removed_sensitive_keywords = list(dict.fromkeys(
                    removed_generated_sensitive + removed_existing_sensitive
                ))

            # Deduplicate: only write keywords that aren't already in EXIF
            existing_lower = {k.lower() for k in existing_keywords_for_merge}
            new_keywords = [k for k in keywords if k.lower() not in existing_lower]

            self._report(0.8, "Writing keywords to EXIF...")
            if anonymize_keywords:
                # Replace existing keyword set with sanitized + generated keywords.
                combined_target_keywords = sorted(set(existing_keywords_for_merge + keywords))
                self.clear_exif_keyword_tags_only(file_path)
                write_result = self.write_exif_keywords(file_path, combined_target_keywords, description)
                new_keywords = [
                    k for k in combined_target_keywords
                    if k.lower() not in {x.lower() for x in existing_keywords}
                ]
            else:
                write_result = self.write_exif_keywords(file_path, new_keywords, description)
        else:
            write_result = {"success": True, "message": "Keyword generation skipped", "keywords_written": 0}

        if populate_location:
            if location and any(location.values()):
                location_write_ok = self.write_exif_location(
                    file_path,
                    location.get("country", ""),
                    location.get("state", ""),
                    location.get("city", ""),
                )
                location_result["location_written"] = True
                location_result["location_write_ok"] = location_write_ok
            original_gps = location_resolution.get("existing_gps")
            if gps and original_gps is None:
                gps_write_result = self.write_gps_coordinates(file_path, gps[0], gps[1])
                location_result["gps_written"] = True
                location_result["gps_write_ok"] = bool(gps_write_result.get("success"))
                location_result["gps_write_result"] = gps_write_result

        ownership_result = None
        if ownership_notice.strip():
            ownership_result = self.write_ownership_metadata(file_path, ownership_notice.strip())

        # Combined set for display (existing + new)
        combined_keywords = (
            sorted(set(existing_keywords_for_merge + keywords))
            if generate_description and anonymize_keywords
            else sorted(set(existing_keywords + keywords))
        )

        self._report(1.0, "Done")
        return {
            "file_name": file_name,
            "description": description,
            "keywords": combined_keywords,
            "new_keywords": new_keywords,
            "existing_keywords": existing_keywords,
            "removed_sensitive_keywords": removed_sensitive_keywords,
            "has_gps": has_gps,
            "location": location,
            "gps": gps,
            "location_result": location_result,
            "location_resolution": location_resolution,
            "exif_result": write_result,
            "ownership_result": ownership_result,
            "resize_info": resize_info,
        }

    # ------------------------------------------------------------------
    # Docling fallback
    # ------------------------------------------------------------------

    def _try_docling(self, file_path: str, timeout_seconds: Optional[float] = None) -> Optional[str]:
        """Try using Docling for higher-quality conversion. Returns None if unavailable."""
        if timeout_seconds and timeout_seconds > 0:
            result = self._try_docling_subprocess(file_path, timeout_seconds)
            if result:
                md_text, figures = result
                if str(file_path).lower().endswith(".pdf"):
                    md_text = self._enrich_docling_image_markers(md_text, file_path, figures)
                logger.info("Used Docling subprocess for PDF conversion")
                return self._normalize_markdown_output(md_text)
            return None

        # Prefer Cortex Docling reader (handles model repair / pipeline options).
        try:
            from cortex_engine.docling_reader import DoclingDocumentReader
            reader = DoclingDocumentReader()
            if reader.is_available and reader.can_process_file(file_path):
                docs = self._run_with_timeout(reader.load_data, timeout_seconds, file_path)
                if docs is None:
                    logger.warning(
                        f"Docling timed out after {timeout_seconds}s for {Path(file_path).name}; "
                        "falling back to legacy reader"
                    )
                    return None
                if docs and docs[0].text and docs[0].text.strip():
                    logger.info("Used DoclingDocumentReader for PDF conversion")
                    md_text = docs[0].text
                    figures = []
                    try:
                        figures = (docs[0].metadata or {}).get("docling_figures", []) or []
                    except Exception:
                        figures = []
                    if str(file_path).lower().endswith(".pdf"):
                        md_text = self._enrich_docling_image_markers(md_text, file_path, figures)
                    return self._normalize_markdown_output(md_text)
        except Exception as e:
            logger.info(f"DoclingDocumentReader path unavailable: {e}")

        # Fallback direct Docling converter path.
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            result = self._run_with_timeout(converter.convert, timeout_seconds, file_path)
            if result is None:
                logger.warning(
                    f"Direct Docling converter timed out after {timeout_seconds}s for {Path(file_path).name}"
                )
                return None
            md = result.document.export_to_markdown()
            if md and md.strip():
                logger.info("Used Docling for PDF conversion")
                return self._normalize_markdown_output(md)
        except Exception as e:
            logger.info(f"Direct Docling converter unavailable: {e}")
        return None
