"""
Document Textifier Module
Converts PDF, DOCX, and PPTX documents to rich Markdown with optional
vision-model descriptions of images and tables.
"""

import os
import base64
import tempfile
from pathlib import Path
from typing import Callable, List, Optional, Tuple

from cortex_engine.utils.logging_utils import get_logger

logger = get_logger(__name__)


class DocumentTextifier:
    """Converts documents to Markdown with optional VLM image descriptions."""

    # Preferred vision models in order of priority
    VISION_MODELS = ["qwen3-vl:8b", "qwen3-vl:4b", "qwen3-vl"]

    def __init__(self, use_vision: bool = True, on_progress: Optional[Callable[[float, str], None]] = None):
        self.use_vision = use_vision
        self.on_progress = on_progress
        self._vlm_client = None
        self._vlm_model = None

    def _report(self, fraction: float, message: str):
        """Send a progress update if a callback is registered."""
        if self.on_progress:
            self.on_progress(fraction, message)

    def _init_vlm(self):
        """Lazy-init the Ollama VLM client, preferring Qwen3-VL models."""
        if self._vlm_client is not None:
            return
        try:
            import ollama
            client = ollama.Client(timeout=120)

            # Try preferred Qwen3-VL models first
            for model in self.VISION_MODELS:
                try:
                    client.show(model)
                    self._vlm_client = client
                    self._vlm_model = model
                    logger.info(f"Textifier using vision model: {model}")
                    return
                except Exception:
                    continue

            # Fall back to configured VLM_MODEL (e.g. llava:7b)
            from cortex_engine.config import VLM_MODEL
            client.show(VLM_MODEL)
            self._vlm_client = client
            self._vlm_model = VLM_MODEL
            logger.info(f"Textifier falling back to VLM_MODEL: {VLM_MODEL}")
        except Exception as e:
            logger.warning(f"No vision model available: {e}")
            self._vlm_client = None
            self._vlm_model = None

    def describe_image(self, image_bytes: bytes) -> str:
        """Describe an image using the VLM. Returns placeholder on failure."""
        if not self.use_vision:
            return "[Image: vision model disabled]"
        self._init_vlm()
        if self._vlm_client is None:
            return "[Image: could not be described — vision model unavailable]"
        try:
            encoded = base64.b64encode(image_bytes).decode("utf-8")
            response = self._vlm_client.chat(
                model=self._vlm_model,
                messages=[{
                    "role": "user",
                    "content": "Describe this image concisely. Focus on key content, text, data, or diagrams visible.",
                    "images": [encoded],
                }],
                options={"temperature": 0.1, "num_predict": 300},
            )
            return response["message"]["content"].strip()
        except Exception as e:
            logger.warning(f"VLM describe failed: {e}")
            return "[Image: could not be described — vision model error]"

    @staticmethod
    def table_to_markdown(rows: List[List[str]]) -> str:
        """Convert a list-of-lists table to Markdown."""
        if not rows:
            return ""
        # Ensure all rows have same column count
        max_cols = max(len(r) for r in rows)
        normalised = [r + [""] * (max_cols - len(r)) for r in rows]
        header = normalised[0]
        lines = ["| " + " | ".join(header) + " |"]
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in normalised[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------

    def textify_file(self, file_path: str) -> str:
        """Dispatch to the correct converter based on file extension."""
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self.textify_pdf(file_path)
        elif ext == ".docx":
            return self.textify_docx(file_path)
        elif ext == ".pptx":
            return self.textify_pptx(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"):
            return self.textify_image(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------

    def textify_pdf(self, file_path: str) -> str:
        """Convert a PDF to Markdown using PyMuPDF."""
        self._report(0.0, "Checking Docling availability...")
        docling_md = self._try_docling(file_path)
        if docling_md:
            self._report(1.0, "Converted via Docling")
            return docling_md

        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        total_pages = len(doc)
        md_parts: List[str] = []

        for page_num in range(total_pages):
            self._report(page_num / total_pages, f"Page {page_num + 1}/{total_pages}")
            page = doc[page_num]
            md_parts.append(f"## Page {page_num + 1}\n")

            text = page.get_text("text")
            if text.strip():
                md_parts.append(text.strip())
                md_parts.append("")

            image_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    if base_image and base_image.get("image"):
                        self._report(
                            (page_num + (img_idx + 1) / max(len(image_list), 1)) / total_pages,
                            f"Page {page_num + 1} — describing image {img_idx + 1}/{len(image_list)}",
                        )
                        desc = self.describe_image(base_image["image"])
                        md_parts.append(f"> **[Image {img_idx + 1}]**: {desc}")
                        md_parts.append("")
                except Exception as e:
                    logger.debug(f"Could not extract image xref={xref}: {e}")

        doc.close()
        self._report(1.0, "PDF conversion complete")
        return "\n".join(md_parts)

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def textify_docx(self, file_path: str) -> str:
        """Convert a DOCX to Markdown."""
        from docx import Document as DocxDocument
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        import io

        doc = DocxDocument(file_path)
        md_parts: List[str] = []

        total_paras = len(doc.paragraphs)
        self._report(0.0, "Processing paragraphs...")

        for p_idx, para in enumerate(doc.paragraphs):
            if total_paras > 0 and p_idx % 20 == 0:
                self._report(p_idx / total_paras * 0.5, f"Paragraph {p_idx}/{total_paras}")
            style_name = (para.style.name or "").lower()
            text = para.text.strip()
            if not text:
                md_parts.append("")
                continue
            if "heading 1" in style_name:
                md_parts.append(f"# {text}")
            elif "heading 2" in style_name:
                md_parts.append(f"## {text}")
            elif "heading 3" in style_name:
                md_parts.append(f"### {text}")
            elif "heading" in style_name:
                md_parts.append(f"#### {text}")
            elif "list" in style_name:
                md_parts.append(f"- {text}")
            else:
                md_parts.append(text)
            md_parts.append("")

        self._report(0.5, "Processing tables...")
        for table_idx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                md_parts.append(f"**Table {table_idx + 1}:**")
                md_parts.append(self.table_to_markdown(rows))
                md_parts.append("")

        # Count images first for progress
        image_rels = [r for r in doc.part.rels.values() if "image" in r.reltype]
        total_images = len(image_rels)
        self._report(0.7, f"Processing {total_images} embedded images...")

        img_idx = 0
        for rel_i, rel in enumerate(image_rels):
            try:
                image_bytes = rel.target_part.blob
                self._report(0.7 + 0.3 * (rel_i / max(total_images, 1)),
                             f"Describing image {rel_i + 1}/{total_images}")
                desc = self.describe_image(image_bytes)
                img_idx += 1
                md_parts.append(f"> **[Image {img_idx}]**: {desc}")
                md_parts.append("")
            except Exception as e:
                logger.debug(f"Could not extract DOCX image: {e}")

        return "\n".join(md_parts)

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def textify_pptx(self, file_path: str) -> str:
        """Convert a PPTX to Markdown."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        md_parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            self._report((slide_num - 1) / total_slides, f"Slide {slide_num}/{total_slides}")
            md_parts.append(f"## Slide {slide_num}\n")

            img_count = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            md_parts.append(text)
                    md_parts.append("")

                if shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        md_parts.append(self.table_to_markdown(rows))
                        md_parts.append("")

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_count += 1
                        self._report(
                            (slide_num - 1) / total_slides,
                            f"Slide {slide_num}/{total_slides} — describing image {img_count}",
                        )
                        image_bytes = shape.image.blob
                        desc = self.describe_image(image_bytes)
                        md_parts.append(f"> **[Image]**: {desc}")
                        md_parts.append("")
                    except Exception as e:
                        logger.debug(f"Could not extract PPTX image: {e}")

            md_parts.append("---\n")

        self._report(1.0, "PPTX conversion complete")

        return "\n".join(md_parts)

    # ------------------------------------------------------------------
    # Image files (PNG, JPG, etc.)
    # ------------------------------------------------------------------

    def textify_image(self, file_path: str) -> str:
        """Convert an image file to Markdown with a VLM description."""
        self._report(0.0, "Reading image...")
        file_name = Path(file_path).name

        with open(file_path, "rb") as f:
            image_bytes = f.read()

        md_parts = [f"# {file_name}\n"]

        self._report(0.3, "Describing image with vision model...")
        description = self.describe_image(image_bytes)
        md_parts.append(f"> **[Image]**: {description}")
        md_parts.append("")

        # Try OCR extraction via pytesseract if available
        try:
            from PIL import Image
            import pytesseract
            self._report(0.7, "Extracting text via OCR...")
            img = Image.open(file_path)
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text:
                md_parts.append("## Extracted Text (OCR)\n")
                md_parts.append(ocr_text)
                md_parts.append("")
        except ImportError:
            logger.debug("pytesseract not available — skipping OCR")
        except Exception as e:
            logger.debug(f"OCR extraction failed: {e}")

        self._report(1.0, "Image conversion complete")
        return "\n".join(md_parts)

    # ------------------------------------------------------------------
    # Docling fallback
    # ------------------------------------------------------------------

    def _try_docling(self, file_path: str) -> Optional[str]:
        """Try using Docling for higher-quality conversion. Returns None if unavailable."""
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            result = converter.convert(file_path)
            md = result.document.export_to_markdown()
            if md and md.strip():
                logger.info("Used Docling for PDF conversion")
                return md
        except Exception:
            pass
        return None
