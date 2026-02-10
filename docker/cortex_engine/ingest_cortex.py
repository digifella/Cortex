# ## File: ingest_cortex.py
# Version: 14.1.0 (Embedding Model Safeguards)
# Date: 2025-12-27
# Purpose: Core ingestion script for Project Cortex with integrated knowledge graph extraction.
#          - CRITICAL FIX (v14.1.0): Added embedding model validation and metadata storage
#            to prevent mixed embedding corruption. System now validates model compatibility
#            before ingestion and stores model metadata for future validation.
#          - FEATURE (v13.0.0): Integrated entity extraction and knowledge graph building
#            during the ingestion process. The system now extracts people, organizations,
#            projects, and their relationships while maintaining backward compatibility.

import os
import sys
os.environ['ANONYMIZED_TELEMETRY'] = 'False'

import argparse
import json
import logging
import warnings
import subprocess

# Suppress common warnings that don't affect functionality
warnings.filterwarnings("ignore", message=".*attention_mask.*")
warnings.filterwarnings("ignore", message=".*pad_token_id.*")
warnings.filterwarnings("ignore", category=FutureWarning, module="transformers")

# Suppress Docling warnings about API compatibility (non-blocking)
warnings.filterwarnings("ignore", message=".*Docling.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*docling.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*docling_reader.*", category=UserWarning)

# Suppress Pydantic field name warnings (non-functional)
warnings.filterwarnings("ignore", message=".*field names.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*pydantic.*field.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*has conflict with protected namespace.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*model_.*", category=UserWarning)

# Suppress torchvision and model loading warnings (informational only)
warnings.filterwarnings("ignore", message=".*torchvision.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*CUDA.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*GPU.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*torch.load.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*vulnerability.*", category=UserWarning)
warnings.filterwarnings("ignore", message=".*CVE-.*", category=UserWarning)

# Suppress sentence-transformers verbosity
warnings.filterwarnings("ignore", category=FutureWarning, module="sentence_transformers")
warnings.filterwarnings("ignore", category=UserWarning, module="sentence_transformers")

# Suppress huggingface model warnings
warnings.filterwarnings("ignore", category=FutureWarning, module="huggingface_hub")

# Reduce verbosity of specific loggers that tend to be noisy
logging.getLogger("unstructured").setLevel(logging.ERROR)
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("sentence_transformers").setLevel(logging.WARNING)
logging.getLogger("huggingface_hub").setLevel(logging.WARNING)
logging.getLogger("torch").setLevel(logging.WARNING)
logging.getLogger("torchvision").setLevel(logging.ERROR)

# Suppress ChromaDB telemetry errors (harmless version mismatch)
logging.getLogger("chromadb.telemetry").setLevel(logging.CRITICAL)
import sys
# Redirect stderr temporarily to suppress telemetry print statements
class SuppressTelemetryErrors:
    def __enter__(self):
        self._original_stderr = sys.stderr
        sys.stderr = open(os.devnull, 'w')
        return self
    def __exit__(self, *args):
        sys.stderr.close()
        sys.stderr = self._original_stderr
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Literal, Tuple
import threading
import re

from pydantic import BaseModel, Field, ValidationError
from llama_index.core import Document
from llama_index.core.settings import Settings
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.core import VectorStoreIndex, StorageContext
import chromadb
from chromadb.config import Settings as ChromaSettings
import hashlib

from llama_index.readers.file import (
    DocxReader,
    PptxReader,
    PyMuPDFReader,
    FlatReader,
    UnstructuredReader
)

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from cortex_engine.config import (
    INGESTION_LOG_PATH, STAGING_INGESTION_FILE, INGESTED_FILES_LOG,
    COLLECTION_NAME, get_embed_model,
    TABLE_AWARE_CHUNKING, TABLE_SPECIFIC_EMBEDDINGS, FIGURE_ENTITY_LINKING
)
from cortex_engine.utils.file_utils import get_file_hash
from cortex_engine.utils.logging_utils import get_logger
from cortex_engine.utils.smart_ollama_llm import create_smart_ollama_llm
from cortex_engine.utils.path_utils import ensure_directory_writable
try:
    from cortex_engine.migration_to_docling import create_migration_manager
except Exception:
    def create_migration_manager():
        class _Noop:
            def migrate(self, *args, **kwargs):
                return None
        return _Noop()

logger = get_logger(__name__)
from cortex_engine.query_cortex import describe_image_with_vlm_for_ingestion, describe_image_with_vlm_async
from cortex_engine.embedding_adapters import EmbeddingServiceAdapter
from cortex_engine.entity_extractor import EntityExtractor, ExtractedEntity, ExtractedRelationship
from cortex_engine.graph_manager import EnhancedGraphManager
from cortex_engine.batch_manager import BatchState

# Phase 2 Enhancements: Table-aware chunking and figure entity linking
try:
    from cortex_engine.table_chunking_enhancer import create_table_aware_chunker
    from cortex_engine.figure_entity_linker import create_figure_entity_linker, load_knowledge_graph_for_linking
    PHASE2_AVAILABLE = True
except ImportError as e:
    logger.warning(f"Phase 2 enhancements not available: {e}")
    PHASE2_AVAILABLE = False

LOG_FILE = INGESTION_LOG_PATH
# STAGING_FILE will be set dynamically based on runtime db_path
COLLECTIONS_FILE = str(project_root / "working_collections.json")

def _ensure_directory_cross_platform(dir_path: str) -> None:
    """Ensure directory exists, with Windows PowerShell fallback for WSL permission issues."""
    is_wsl = False
    if os.name == "posix":
        if os.environ.get("WSL_DISTRO_NAME"):
            is_wsl = True
        else:
            try:
                with open("/proc/sys/kernel/osrelease") as release_file:
                    is_wsl = "microsoft" in release_file.read().lower()
            except OSError:
                is_wsl = False

    try:
        os.makedirs(dir_path, exist_ok=True)
    except PermissionError as e:
        if is_wsl and dir_path.startswith("/mnt/"):
            try:
                win_path = subprocess.run(
                    ["wslpath", "-w", dir_path],
                    capture_output=True,
                    text=True,
                    check=True,
                ).stdout.strip()
                subprocess.run(
                    [
                        "powershell.exe",
                        "-NoProfile",
                        "-Command",
                        f'New-Item -ItemType Directory -Path "{win_path}" -Force | Out-Null',
                    ],
                    check=True,
                )
                os.makedirs(dir_path, exist_ok=True)
                logging.info(f"Created directory via PowerShell fallback: {dir_path}")
            except Exception as fallback_error:
                raise PermissionError(
                    f"Permission denied creating '{dir_path}'. "
                    f"Attempted PowerShell fallback and failed: {fallback_error}"
                ) from e
        else:
            raise

def get_staging_file_path(db_path: str) -> str:
    """Get the staging file path for the given database path"""
    return os.path.join(db_path, "staging_ingestion.json")

# Define supported image extensions - Enhanced Visual Processing Support
IMAGE_EXTENSIONS = {
    # Standard formats
    ".png", ".jpg", ".jpeg",
    # Additional formats for comprehensive visual processing
    ".gif", ".bmp", ".webp", ".tiff", ".tif",
    # Vector and specialized formats
    ".svg",  # Will be converted to raster for VLM processing
    ".ico"   # Icon files
}

os.makedirs(os.path.dirname(LOG_FILE), exist_ok=True)

def initialize_script():
    """Sets up robust logging and configures LlamaIndex models."""
    log_formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    root_logger = logging.getLogger()

    if root_logger.hasHandlers():
        root_logger.handlers.clear()

    root_logger.setLevel(logging.INFO)

    file_handler = logging.FileHandler(LOG_FILE, mode='a', encoding='utf-8')
    file_handler.setFormatter(log_formatter)
    root_logger.addHandler(file_handler)

    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setFormatter(log_formatter)
    root_logger.addHandler(console_handler)

    logging.info("--- Script Initialized: Configuring models... ---")
    
    # Check if Ollama is available before configuring
    from cortex_engine.utils.ollama_utils import check_ollama_service, format_ollama_error_for_user
    
    chk = check_ollama_service()
    # Backward compatible unpack
    is_running = chk[0]
    error_msg = chk[1] if len(chk) > 1 else None
    resolved_url = chk[2] if len(chk) > 2 else None
    if not is_running:
        logging.warning(f"Ollama service not available: {error_msg}")
        logging.warning("AI-enhanced metadata extraction will be disabled. Documents will be processed with basic metadata only.")
        Settings.llm = None  # Will be handled in analysis function
    else:
        # Set reasonable timeout for metadata extraction (3 minutes)
        # Very long documents may timeout and use fallback metadata
        Settings.llm = create_smart_ollama_llm(model="mistral:latest", request_timeout=180.0)
        logging.info(f"Ollama connected successfully at {resolved_url}")
    
    # Unify embeddings with search/async via adapter around embedding_service (Docker)
    # Use adaptive model selection
    embed_model_name = get_embed_model()
    try:
        Settings.embed_model = EmbeddingServiceAdapter(model_name=embed_model_name)
        logging.info(f"Models configured via EmbeddingServiceAdapter (Embed: {embed_model_name}).")
    except Exception as e:
        logging.warning(f"EmbeddingServiceAdapter failed ({e}); falling back to HuggingFaceEmbedding")
        Settings.embed_model = HuggingFaceEmbedding(model_name=embed_model_name)
        logging.info(f"Models configured (Embed: {embed_model_name}).")

    # LLM safety: add a hard-timeout wrapper to prevent hangs during metadata extraction
    # Uses a single worker thread to enforce a wall-clock timeout per LLM call.
    # On timeout, callers should fall back to basic metadata and continue.
    global _run_with_timeout
    def _run_with_timeout(fn, timeout_s: float):
        _futures = __import__('concurrent').futures
        _executor = _futures.ThreadPoolExecutor(max_workers=1, thread_name_prefix="metadata-timeout")
        _future = _executor.submit(fn)
        try:
            return _future.result(timeout=timeout_s)
        except _futures.TimeoutError as exc:
            _future.cancel()
            raise TimeoutError(f"Operation timed out after {timeout_s:.1f}s") from exc
        finally:
            # Avoid blocking shutdown when the worker is still running after a timeout.
            _executor.shutdown(wait=_future.done())

class RichMetadata(BaseModel):
    document_type: Literal[
        "Project Plan", "Technical Documentation", "Proposal/Quote", "Case Study / Trophy",
        "Final Report", "Draft Report", "Presentation", "Contract/SOW",
        "Meeting Minutes", "Financial Report", "Research Paper", "Email Correspondence", "Image/Diagram", "Other"
    ] = Field(..., description="The primary category of the document.")
    proposal_outcome: Literal["Won", "Lost", "Pending", "N/A"] = Field(..., description="The outcome of the proposal, if applicable.")
    summary: str = Field(..., description="A concise, 1-3 sentence summary of the document's content and purpose.")
    thematic_tags: List[str] = Field(default_factory=list, description="A list of 3-5 key themes, topics, or technologies discussed.")
    credibility_tier_value: Literal[0, 1, 2, 3, 4, 5] = Field(
        0,
        description="Credibility tier score: 5=peer-reviewed, 4=institutional, 3=pre-print, 2=editorial, 1=commentary, 0=unclassified."
    )
    credibility_tier_key: Literal["peer-reviewed", "institutional", "pre-print", "editorial", "commentary", "unclassified"] = Field(
        "unclassified",
        description="Normalized credibility key for filtering and reporting."
    )
    credibility_tier_label: Literal["Peer-Reviewed", "Institutional", "Pre-Print", "Editorial", "Commentary", "Unclassified"] = Field(
        "Unclassified",
        description="Human-readable credibility tier label."
    )
    credibility: str = Field(
        "Final Unclassified Report",
        description='Human-readable credibility summary in the format "{Draft|Final} {Label} Report".'
    )


_CREDIBILITY_BY_VALUE = {
    5: ("peer-reviewed", "Peer-Reviewed"),
    4: ("institutional", "Institutional"),
    3: ("pre-print", "Pre-Print"),
    2: ("editorial", "Editorial"),
    1: ("commentary", "Commentary"),
    0: ("unclassified", "Unclassified"),
}
_CREDIBILITY_BY_KEY = {k: (v, label) for v, (k, label) in _CREDIBILITY_BY_VALUE.items()}
_CREDIBILITY_BY_LABEL = {label.lower(): (v, key) for v, (key, label) in _CREDIBILITY_BY_VALUE.items()}


def _normalize_credibility_tier(metadata_json: Dict) -> None:
    """Normalize credibility tier fields in-place so value/key/label are consistent."""
    if not isinstance(metadata_json, dict):
        return

    default_value = 0
    default_key, default_label = _CREDIBILITY_BY_VALUE[default_value]

    raw_value = metadata_json.get("credibility_tier_value")
    raw_key = str(metadata_json.get("credibility_tier_key", "")).strip().lower()
    raw_label = str(metadata_json.get("credibility_tier_label", "")).strip()

    if raw_value is not None:
        try:
            tier_value = int(raw_value)
            if tier_value in _CREDIBILITY_BY_VALUE:
                key, label = _CREDIBILITY_BY_VALUE[tier_value]
                metadata_json["credibility_tier_value"] = tier_value
                metadata_json["credibility_tier_key"] = key
                metadata_json["credibility_tier_label"] = label
                return
        except (TypeError, ValueError):
            pass

    if raw_key in _CREDIBILITY_BY_KEY:
        tier_value, label = _CREDIBILITY_BY_KEY[raw_key]
        metadata_json["credibility_tier_value"] = tier_value
        metadata_json["credibility_tier_key"] = raw_key
        metadata_json["credibility_tier_label"] = label
        return

    if raw_label.lower() in _CREDIBILITY_BY_LABEL:
        tier_value, key = _CREDIBILITY_BY_LABEL[raw_label.lower()]
        metadata_json["credibility_tier_value"] = tier_value
        metadata_json["credibility_tier_key"] = key
        metadata_json["credibility_tier_label"] = _CREDIBILITY_BY_VALUE[tier_value][1]
        return

    metadata_json["credibility_tier_value"] = default_value
    metadata_json["credibility_tier_key"] = default_key
    metadata_json["credibility_tier_label"] = default_label


def _detect_document_stage(file_path: str, text: str) -> str:
    """Classify document stage for human-readable credibility field."""
    full_text = f"{file_path}\n{text}".lower()
    return "Draft" if "draft" in full_text else "Final"


def _detect_marker_credibility_value(file_path: str, text: str) -> int:
    """Find credibility tier from marker matches (priority 5->1)."""
    haystack = f"{file_path}\n{text}".lower()
    marker_map = {
        5: ["pubmed", "nlm", "nature", "lancet", "jama", "bmj", "peer-reviewed", "peer reviewed"],
        4: ["who", "un ", "ipcc", "oecd", "world bank", "government", "department", "ministry", "university", "institute", "centre", "center"],
        3: ["arxiv", "ssrn", "biorxiv", "researchgate", "preprint", "pre-print"],
        2: ["scientific american", "the conversation", "hbr", "harvard business review", "editorial"],
        1: ["blog", "newsletter", "opinion", "consulting report", "whitepaper", "white paper"],
    }
    for tier in (5, 4, 3, 2, 1):
        if any(marker in haystack for marker in marker_map[tier]):
            return tier
    return 0


def _enforce_credibility_policy(metadata_json: Dict, file_path: str, text: str) -> None:
    """Apply canonical credibility policy and emit all canonical fields."""
    if not isinstance(metadata_json, dict):
        return

    # Legacy aliases.
    if "credibility_tier_value" not in metadata_json and "credibility_value" in metadata_json:
        metadata_json["credibility_tier_value"] = metadata_json.get("credibility_value")
    if "credibility_tier_key" not in metadata_json and "credibility_source" in metadata_json:
        metadata_json["credibility_tier_key"] = str(metadata_json.get("credibility_source", "")).strip().lower()

    _normalize_credibility_tier(metadata_json)

    normalized_value = int(metadata_json.get("credibility_tier_value", 0) or 0)
    summary = str(metadata_json.get("summary", "") or "")
    source_type = str(metadata_json.get("source_type", "") or "")
    combined_text = f"{summary}\n{text}"

    ai_markers = ["ai generated report", "generated by ai", "chatgpt", "openai", "claude", "gemini", "perplexity"]
    is_ai_generated = source_type.lower() == "ai generated report" or any(m in combined_text.lower() for m in ai_markers)
    marker_value = _detect_marker_credibility_value(file_path, combined_text)

    if is_ai_generated:
        final_value = 1
    elif marker_value > 0:
        final_value = marker_value
    else:
        final_value = normalized_value

    key, label = _CREDIBILITY_BY_VALUE.get(final_value, _CREDIBILITY_BY_VALUE[0])
    stage = _detect_document_stage(file_path, combined_text)
    metadata_json["credibility_tier_value"] = final_value
    metadata_json["credibility_tier_key"] = key
    metadata_json["credibility_tier_label"] = label
    metadata_json["credibility"] = f"{stage} {label} Report"

class DocumentMetadata(BaseModel):
    doc_id: str
    doc_posix_path: str
    file_name: str
    last_modified_date: str
    rich_metadata: Optional[RichMetadata] = None
    exclude_from_final: bool = False
    extracted_entities: List[Dict] = Field(default_factory=list)
    extracted_relationships: List[Dict] = Field(default_factory=list)
    target_collection: Optional[str] = None  # Collection assignment for finalization

def load_processed_files_log(log_path: str) -> Dict[str, str]:
    if not os.path.exists(log_path): return {}
    with open(log_path, 'r') as f:
        try:
            log_data = json.load(f)
            return {k: (v[0] if isinstance(v, list) else v) for k, v in log_data.items()}
        except json.JSONDecodeError: return {}

def write_to_processed_log(log_path: str, file_path: str, doc_id: str):
    """Write a processed file entry to the log."""
    processed_files = {}
    if os.path.exists(log_path):
        try:
            with open(log_path, 'r') as f:
                processed_files = json.load(f)
        except (json.JSONDecodeError, IOError):
            processed_files = {}
    
    processed_files[file_path] = doc_id
    
    try:
        with open(log_path, 'w') as f:
            json.dump(processed_files, f, indent=2)
    except IOError as e:
        logging.error(f"Failed to write to processed log: {e}")


def _get_existing_doc_ids_in_collection(chroma_db_path: str, candidate_doc_ids: set[str]) -> set[str]:
    """
    Return candidate doc_ids that already exist in the vector collection.
    This is a hard duplicate gate independent of ingestion logs/path history.
    """
    existing: set[str] = set()
    if not candidate_doc_ids:
        return existing
    if not os.path.isdir(chroma_db_path):
        return existing

    try:
        db_settings = ChromaSettings(anonymized_telemetry=False)
        chroma_client = chromadb.PersistentClient(path=chroma_db_path, settings=db_settings)
        collections = chroma_client.list_collections()
        if not any(c.name == COLLECTION_NAME for c in collections):
            return existing
        collection = chroma_client.get_collection(COLLECTION_NAME)
    except Exception as e:
        logging.warning(f"Could not open Chroma collection for duplicate check: {e}")
        return existing

    ids_list = list(candidate_doc_ids)
    batch_size = 200
    for i in range(0, len(ids_list), batch_size):
        batch = ids_list[i:i + batch_size]
        try:
            # Fast path: query all candidates in one where-clause.
            res = collection.get(where={"doc_id": {"$in": batch}}, include=["metadatas"])
            for meta in res.get("metadatas", []) or []:
                if isinstance(meta, dict):
                    value = meta.get("doc_id")
                    if value in candidate_doc_ids:
                        existing.add(value)
        except Exception:
            # Fallback path for backends that don't support $in.
            for doc_id in batch:
                try:
                    res = collection.get(where={"doc_id": doc_id}, limit=1, include=["metadatas"])
                    metas = res.get("metadatas", []) or []
                    if metas:
                        existing.add(doc_id)
                except Exception:
                    continue

    return existing

def get_document_content(file_path: str, skip_image_processing: bool = False) -> Dict:
    """Extract content from a single document file
    
    Returns:
        Dict with keys: 'content' (str), 'metadata' (dict), 'source_type' (str)
    """
    from cortex_engine.query_cortex import describe_image_with_vlm_for_ingestion
    
    path = Path(file_path)
    extension = path.suffix.lower()
    
    # Set up readers
    reader_map = {".pdf": PyMuPDFReader(), ".docx": DocxReader(), ".pptx": PptxReader(), 
                  ".doc": UnstructuredReader(), ".ppt": UnstructuredReader()}
    default_reader = FlatReader()
    
    try:
        # Handle image files using the VLM (if enabled)
        if extension in IMAGE_EXTENSIONS:
            if skip_image_processing:
                return {
                    'content': f"Image file: {path.name} (processing skipped)",
                    'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                    'source_type': 'image_skipped'
                }
            else:
                try:
                    description = describe_image_with_vlm_for_ingestion(file_path)
                    return {
                        'content': description,
                        'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                        'source_type': 'image'
                    }
                except Exception as e:
                    return {
                        'content': f"Image file: {path.name} (VLM processing failed: {str(e)})",
                        'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                        'source_type': 'image_error'
                    }

        # Text-based document handling
        reader = reader_map.get(extension, default_reader)
        
        # Handle different reader types with error handling
        if isinstance(reader, PyMuPDFReader):
            try:
                docs_from_file = reader.load_data(path)
            except Exception as pdf_error:
                if "cmsOpenProfileFromMem" not in str(pdf_error):
                    logger.warning(f"PDF processing warning for {path.name}: {pdf_error}")
                docs_from_file = reader.load_data(path)
        elif isinstance(reader, UnstructuredReader):
            try:
                docs_from_file = reader.load_data(path)
            except Exception as unstructured_error:
                error_str = str(unstructured_error).lower()
                if not any(x in error_str for x in ['warning', 'deprecation', 'future']):
                    logger.warning(f"UnstructuredReader processing issue for {path.name}: {unstructured_error}")
                docs_from_file = [Document(
                    text=f"Error processing old Office document: {path.name}. Reason: {unstructured_error}",
                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                )]
        else:
            try:
                docs_from_file = reader.load_data(file=path)
            except UnicodeDecodeError as encoding_error:
                logger.warning(f"Encoding error for {path.name}: {encoding_error}")
                docs_from_file = [Document(
                    text=f"File could not be processed due to encoding issues: {path.name}",
                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                )]
            except OSError as wmf_error:
                if "cannot find loader for this WMF file" in str(wmf_error):
                    logger.warning(f"WMF image error in {path.name}: {wmf_error}")
                    docs_from_file = [Document(
                        text=f"Document processing failed due to WMF image issues: {path.name}",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
                else:
                    raise wmf_error
        
        # Combine text from multiple documents if needed
        if docs_from_file:
            combined_text = "\n\n".join([doc.text for doc in docs_from_file if doc.text])
            metadata = docs_from_file[0].metadata if docs_from_file else {}
            metadata.update({'file_path': str(path.as_posix()), 'file_name': path.name})
            
            return {
                'content': combined_text,
                'metadata': metadata,
                'source_type': metadata.get('source_type', 'document')
            }
        else:
            return {
                'content': f"No content extracted from {path.name}",
                'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
                'source_type': 'document_error'
            }
            
    except Exception as e:
        logger.error(f"Failed to process {file_path}: {e}")
        return {
            'content': f"Error processing {path.name}: {str(e)}",
            'metadata': {'file_path': str(path.as_posix()), 'file_name': path.name},
            'source_type': 'document_error'
        }


def create_document_with_entities(content_result: Dict, rich_metadata: Optional[RichMetadata] = None) -> Document:
    """Create a Document object with entity extraction"""
    content = content_result.get('content', '')
    metadata = content_result.get('metadata', {})
    source_type = content_result.get('source_type', 'document')
    
    # Create the document
    doc = Document(text=content)
    doc.metadata.update(metadata)
    doc.metadata['source_type'] = source_type
    
    # Add rich metadata if provided
    if rich_metadata:
        doc.metadata['rich_metadata'] = rich_metadata
    
    return doc


def manual_load_documents(file_paths: List[str], args=None) -> List[Document]:
    """
    Enhanced document loading with Docling integration and intelligent migration.
    
    This function now uses a migration strategy to gradually transition from legacy
    LlamaIndex readers to Docling-enhanced processing while maintaining stability.
    """
    
    # Check for migration mode preference in args
    # In Docker environments, default to legacy mode to avoid Docling dependency conflicts
    import os
    default_mode = 'legacy' if os.path.exists('/.dockerenv') else 'gradual'
    migration_mode = getattr(args, 'migration_mode', default_mode) if args else default_mode
    skip_image_processing = getattr(args, 'skip_image_processing', False) if args else False
    
    # Use migration manager for intelligent processing
    try:
        migration_manager = create_migration_manager(mode=migration_mode)
        documents = migration_manager.process_documents(file_paths, skip_image_processing)
        
        # Log migration insights
        if hasattr(migration_manager, 'comparison_results') and migration_manager.comparison_results:
            logger.info("📊 Migration manager provided processing insights")
        
        return documents
        
    except Exception as e:
        logger.warning(f"⚠️ Migration processing failed, falling back to legacy: {e}")
        # Fallback to legacy processing
        return _legacy_manual_load_documents(file_paths, args)


def _process_images_batch(
    image_files: List[str],
    skip_image_processing: bool = False,
    max_workers: int = 1
) -> List[Document]:
    """
    Process images with VLM (sequential by default to prevent Ollama crashes).

    This function processes images reliably by:
    1. Processing 1 image at a time (llava:7b crashes on concurrent requests)
    2. Using 30s timeout per image
    3. Graceful fallback on timeout/error
    4. Retry logic with exponential backoff for transient errors

    NOTE: When Qwen3-VL is enabled, VLM processing is skipped because Qwen3-VL
    can directly embed images without needing text descriptions from llava.

    Args:
        image_files: List of image file paths to process
        skip_image_processing: Skip VLM processing if True (fast mode)

    Returns:
        List of Document objects with image descriptions or placeholders
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from pathlib import Path

    documents = []

    # Check if Qwen3-VL is enabled - if so, skip llava VLM processing
    # Qwen3-VL can directly embed images without needing text descriptions
    try:
        from cortex_engine.config import QWEN3_VL_ENABLED
        if QWEN3_VL_ENABLED:
            logging.info(f"🎨 Qwen3-VL enabled - skipping llava VLM for {len(image_files)} images (will embed directly)")
            for file_path in image_files:
                path = Path(file_path)
                # Create document with image path for direct Qwen3-VL embedding
                doc = Document(text=f"[Image: {path.name}]")
                doc.metadata['file_path'] = str(path.as_posix())
                doc.metadata['file_name'] = path.name
                doc.metadata['source_type'] = 'image_qwen3vl'
                doc.metadata['image_path'] = str(path.as_posix())  # For Qwen3-VL to embed
                documents.append(doc)
            return documents
    except ImportError:
        pass  # Config not available, continue with llava

    if skip_image_processing:
        logging.info(f"⚡ Skipping VLM processing for {len(image_files)} images (fast mode)")
        for file_path in image_files:
            path = Path(file_path)
            doc = Document(text=f"Image file: {path.name} (processing skipped)")
            doc.metadata['file_path'] = str(path.as_posix())
            doc.metadata['file_name'] = path.name
            doc.metadata['source_type'] = 'image_skipped'
            documents.append(doc)
        return documents

    logging.info(f"🖼️ Processing {len(image_files)} images with VLM (parallel, 30s timeout each)")

    # Be conservative with workers to reduce heat/CPU load
    workers = max(1, int(max_workers))
    with ThreadPoolExecutor(max_workers=workers) as executor:
        # Submit all image processing tasks
        future_to_path = {
            executor.submit(describe_image_with_vlm_async, img_path, 30): img_path
            for img_path in image_files
        }

        # Process results as they complete
        for idx, future in enumerate(as_completed(future_to_path), 1):
            img_path = future_to_path[future]
            path = Path(img_path)

            logging.info(f"📄 Processing image {idx}/{len(image_files)}: {path.name}")

            try:
                description = future.result()

                if description:
                    doc = Document(text=description)
                    doc.metadata['file_path'] = str(path.as_posix())
                    doc.metadata['file_name'] = path.name
                    doc.metadata['source_type'] = 'image'
                    documents.append(doc)
                    logging.info(f"✅ Successfully processed image: {path.name}")
                else:
                    # Timeout or error - create placeholder
                    doc = Document(text=f"Image file: {path.name} (VLM timeout/error)")
                    doc.metadata['file_path'] = str(path.as_posix())
                    doc.metadata['file_name'] = path.name
                    doc.metadata['source_type'] = 'image_fallback'
                    documents.append(doc)
                    logging.warning(f"⚠️ Image processing failed, using fallback: {path.name}")

            except Exception as e:
                logging.error(f"❌ Unexpected error processing {path.name}: {e}")
                # Create error placeholder
                doc = Document(text=f"Image file: {path.name} (processing error)")
                doc.metadata['file_path'] = str(path.as_posix())
                doc.metadata['file_name'] = path.name
                doc.metadata['source_type'] = 'image_error'
                documents.append(doc)

    successful = len([d for d in documents if d.metadata.get('source_type') == 'image'])
    logging.info(f"🎯 Image processing complete: {successful}/{len(image_files)} successful")

    return documents


def _legacy_manual_load_documents(file_paths: List[str], args=None) -> List[Document]:
    """
    Legacy document loading function (original implementation).
    Kept as fallback for compatibility and robustness.
    """
    documents = []
    
    # Initialize readers with graceful fallback for torch version issues
    reader_map = {".pdf": PyMuPDFReader(), ".docx": DocxReader()}
    
    # Try to initialize PptxReader, fallback to UnstructuredReader if torch issues
    try:
        reader_map[".pptx"] = PptxReader()
        reader_map[".ppt"] = PptxReader()
        logging.info("✅ PowerPoint reader initialized successfully")
    except Exception as e:
        logging.warning(f"⚠️ PowerPoint reader failed to initialize (torch version issue): {e}")
        logging.info("📋 Using UnstructuredReader fallback for PowerPoint files")
        reader_map[".pptx"] = UnstructuredReader()
        reader_map[".ppt"] = UnstructuredReader()
    
    # Add other readers
    reader_map[".doc"] = UnstructuredReader()
    default_reader = FlatReader()
    
    # Suppress PyMuPDF warnings
    import fitz
    fitz.TOOLS.mupdf_display_errors(False)
    
    # Expand directories to get actual files
    expanded_file_paths = []
    for file_path in file_paths:
        path = Path(file_path)
        if path.is_dir():
            # Recursively find all files in directory
            for file in path.rglob('*'):
                if file.is_file():
                    expanded_file_paths.append(str(file))
        elif path.is_file():
            expanded_file_paths.append(str(path))
        else:
            logging.warning(f"Path does not exist or is not a file/directory: {file_path}")

    logging.info(f"Expanded {len(file_paths)} paths to {len(expanded_file_paths)} files")

    # Separate images from other documents for batch processing
    skip_images = getattr(args, 'skip_image_processing', False) if hasattr(args, 'skip_image_processing') else False
    image_files = []
    other_files = []

    for file_path in expanded_file_paths:
        path = Path(file_path)
        extension = path.suffix.lower()
        if extension in IMAGE_EXTENSIONS:
            image_files.append(file_path)
        else:
            other_files.append(file_path)

    logging.info(f"📊 File breakdown: {len(image_files)} images, {len(other_files)} documents")

    # Process images in batch (parallel with timeout)
    if image_files:
        # Sequential processing (1 worker) - llava:7b crashes on concurrent requests
        # Use --image-workers to increase if your VLM supports concurrency
        image_workers = getattr(args, 'image_workers', 1) if hasattr(args, 'image_workers') else 1
        image_docs = _process_images_batch(image_files, skip_images, max_workers=image_workers)
        documents.extend(image_docs)

    # Process other documents normally
    total_files = len(other_files)
    for idx, file_path in enumerate(other_files, 1):
        try:
            path = Path(file_path)
            extension = path.suffix.lower()
            logging.info(f"📄 Processing file {idx}/{total_files}: {path.name}")

            # Images are already processed in batch above, this loop only handles documents
            # Original text-based document handling
            reader = reader_map.get(extension, default_reader)
            logging.info(f"Loading '{file_path}' with reader: {reader.__class__.__name__}")
            
            # For PDFs, add extra error handling
            if isinstance(reader, PyMuPDFReader):
                try:
                    docs_from_file = reader.load_data(path)
                except Exception as pdf_error:
                    # If PDF fails completely, log it but continue
                    if "cmsOpenProfileFromMem" not in str(pdf_error):
                        logging.warning(f"PDF processing warning for {path.name}: {pdf_error}")
                    docs_from_file = reader.load_data(path)
            elif isinstance(reader, UnstructuredReader):
                # UnstructuredReader also uses positional path parameter
                try:
                    docs_from_file = reader.load_data(path)
                except Exception as unstructured_error:
                    # Many UnstructuredReader warnings are harmless parsing issues - only log real errors
                    error_str = str(unstructured_error).lower()
                    if not any(x in error_str for x in ['warning', 'deprecation', 'future']):
                        logging.warning(f"UnstructuredReader processing issue for {path.name}: {unstructured_error}")
                    # Create error document if UnstructuredReader fails
                    docs_from_file = [Document(
                        text=f"Error processing old Office document: {path.name}. Reason: {unstructured_error}",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
            else:
                try:
                    docs_from_file = reader.load_data(file=path)
                except UnicodeDecodeError as encoding_error:
                    logging.warning(f"Encoding error for {path.name} - file appears to be binary or uses unsupported encoding: {encoding_error}")
                    # Create an informational document noting the file couldn't be processed
                    docs_from_file = [Document(
                        text=f"File could not be processed due to encoding issues: {path.name}. This file may be binary or use unsupported text encoding.",
                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                    )]
                except OSError as wmf_error:
                    if "cannot find loader for this WMF file" in str(wmf_error):
                        logging.warning(f"WMF image error in {path.name}, attempting text-only extraction: {wmf_error}")
                        # For PowerPoint files with WMF image issues, try alternative extraction
                        if isinstance(reader, PptxReader):
                            try:
                                # Try to extract just text content without images
                                from pptx import Presentation
                                prs = Presentation(path)
                                text_content = []
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text") and shape.text.strip():
                                            text_content.append(shape.text.strip())
                                
                                if text_content:
                                    combined_text = "\n\n".join(text_content)
                                    docs_from_file = [Document(
                                        text=f"PowerPoint content (text-only extraction due to WMF image issues):\n\n{combined_text}",
                                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_partial'}
                                    )]
                                else:
                                    docs_from_file = [Document(
                                        text=f"PowerPoint file processed but no extractable text found: {path.name}",
                                        metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_partial'}
                                    )]
                            except Exception as pptx_error:
                                logging.error(f"Alternative PowerPoint extraction also failed for {path.name}: {pptx_error}")
                                docs_from_file = [Document(
                                    text=f"PowerPoint file could not be processed: {path.name}. WMF image and text extraction both failed.",
                                    metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                                )]
                        else:
                            # For non-PowerPoint files with WMF issues, create error document
                            docs_from_file = [Document(
                                text=f"Document processing failed due to WMF image issues: {path.name}",
                                metadata={'file_path': str(path.as_posix()), 'file_name': path.name, 'source_type': 'document_error'}
                            )]
                    else:
                        raise wmf_error
                
            for doc in docs_from_file:
                doc.metadata['file_path'] = str(path.as_posix())
                doc.metadata['file_name'] = path.name
                doc.metadata['source_type'] = 'document'
            documents.extend(docs_from_file)
            
        except Exception as e:
            logging.error(f"Failed to load file {file_path}: {e}", exc_info=True)
            documents.append(Document(
                text=f"Error reading this document. Could not load content from {Path(file_path).name}. Reason: {e}",
                metadata={'file_path': str(Path(file_path).as_posix()), 'file_name': Path(file_path).name}
            ))
    return documents

def get_gpu_utilization() -> Optional[float]:
    """
    Get current GPU utilization percentage using nvidia-smi.

    Returns:
        GPU utilization as float (0-100), or None if unavailable
    """
    # 1) Try nvidia-smi (best signal on supported systems)
    try:
        import subprocess
        result = subprocess.run(
            ['nvidia-smi', '--query-gpu=utilization.gpu', '--format=csv,noheader,nounits'],
            capture_output=True,
            text=True,
            timeout=2
        )
        if result.returncode == 0:
            gpu_util = float(result.stdout.strip().split('\n')[0])
            return gpu_util
    except (FileNotFoundError, subprocess.TimeoutExpired, ValueError, IndexError):
        pass

    # 2) Fallback: approximate utilization from CUDA memory (works in WSL when nvidia-smi is unavailable)
    try:
        import torch
        if torch.cuda.is_available():
            total = torch.cuda.get_device_properties(0).total_memory
            reserved = torch.cuda.memory_reserved(0)
            if total and total > 0:
                util_pct = (reserved / total) * 100.0
                # Clamp to [0, 100]
                util_pct = max(0.0, min(100.0, util_pct))
                return float(util_pct)
    except Exception:
        pass
    return None


def get_cpu_utilization() -> Optional[float]:
    """
    Get current CPU utilization percentage using psutil.

    Returns:
        CPU utilization as float (0-100), or None if unavailable
    """
    try:
        import psutil
        # Use interval=None for non-blocking call (uses cached values from background thread)
        # This avoids adding 1 second blocking delay per document
        return psutil.cpu_percent(interval=None)
    except ImportError:
        # Fallback to system load average if psutil is unavailable
        try:
            import os
            load1, _, _ = os.getloadavg()
            cores = os.cpu_count() or 1
            # Approximate utilization percentage relative to core count
            return float(min(100.0, max(0.0, (load1 / cores) * 100.0)))
        except Exception:
            pass
    return None


# ---------------------------------------------------------------------------------
# Heartbeat: periodic status pings to UI (CPU/GPU load) to keep UI fresh
# ---------------------------------------------------------------------------------
_heartbeat_stop_event: Optional[threading.Event] = None
_heartbeat_thread: Optional[threading.Thread] = None


def _heartbeat_loop(period_s: float = 5.0):
    try:
        import time as _t
        while _heartbeat_stop_event and not _heartbeat_stop_event.is_set():
            gpu_util = get_gpu_utilization()
            cpu_util = get_cpu_utilization()
            gpu_str = f"{gpu_util:.0f}" if gpu_util is not None else "N/A"
            cpu_str = f"{cpu_util:.0f}" if cpu_util is not None else "N/A"
            print(f"CORTEX_HEARTBEAT::{gpu_str}::{cpu_str}", flush=True)
            # Sleep in small steps to allow responsive stop
            for _ in range(int(period_s * 10)):
                if _heartbeat_stop_event and _heartbeat_stop_event.is_set():
                    break
                _t.sleep(0.1)
    except Exception:
        pass


def start_heartbeat(period_s: float = 5.0):
    """Start background heartbeat pings for UI."""
    global _heartbeat_stop_event, _heartbeat_thread
    try:
        if _heartbeat_thread and _heartbeat_thread.is_alive():
            return
        _heartbeat_stop_event = threading.Event()
        _heartbeat_thread = threading.Thread(target=_heartbeat_loop, args=(period_s,), daemon=True)
        _heartbeat_thread.start()
    except Exception:
        pass


def stop_heartbeat(timeout_s: float = 1.0):
    """Stop background heartbeat pings."""
    global _heartbeat_stop_event, _heartbeat_thread
    try:
        if _heartbeat_stop_event:
            _heartbeat_stop_event.set()
        if _heartbeat_thread and _heartbeat_thread.is_alive():
            _heartbeat_thread.join(timeout=timeout_s)
        _heartbeat_thread = None
        _heartbeat_stop_event = None
    except Exception:
        pass


def calculate_adaptive_throttle(
    base_delay: float,
    gpu_util: Optional[float],
    cpu_util: Optional[float],
    gpu_threshold: float = 80.0,
    cpu_threshold: float = 85.0,
    increment: float = 0.5,
    max_delay: float = 5.0
) -> Tuple[float, str]:
    """
    Calculate adaptive throttle delay based on system load.

    STABILITY: Enforces a mandatory minimum delay of 1.0 second regardless of
    detected load. Load detection is not always reliable, and uncontrolled
    processing can drive GPU/CPU to dangerous levels causing system instability.

    Args:
        base_delay: User-configured baseline delay (minimum)
        gpu_util: Current GPU utilization (0-100)
        cpu_util: Current CPU utilization (0-100)
        gpu_threshold: GPU threshold to trigger throttling
        cpu_threshold: CPU threshold to trigger throttling
        increment: Amount to increase delay by
        max_delay: Maximum allowed delay

    Returns:
        Tuple of (delay_seconds, reason_string)
    """
    # STABILITY: Mandatory minimum delay to prevent system instability
    # Load detection is not always reliable - this ensures we never run unchecked
    MANDATORY_MINIMUM_DELAY = 1.0

    current_delay = max(base_delay, MANDATORY_MINIMUM_DELAY)
    reasons = []

    # GPU has priority (more likely to cause freezing with Ollama)
    if gpu_util is not None and gpu_util > gpu_threshold:
        overage = (gpu_util - gpu_threshold) / 10  # Scale factor
        additional = increment * max(1, int(overage))
        current_delay += additional
        reasons.append(f"GPU:{gpu_util:.0f}%")

    # CPU fallback or additional throttling
    if cpu_util is not None and cpu_util > cpu_threshold:
        overage = (cpu_util - cpu_threshold) / 10
        additional = increment * max(1, int(overage))
        current_delay += additional
        reasons.append(f"CPU:{cpu_util:.0f}%")

    # Cap at max delay
    current_delay = min(current_delay, max_delay)

    # Always report the delay reason (including mandatory baseline)
    if reasons:
        reason = f"throttle={current_delay:.1f}s ({', '.join(reasons)})"
    else:
        reason = f"stability={current_delay:.1f}s (mandatory)"

    return current_delay, reason


def extract_entities_and_relationships(doc_text: str, metadata: Dict, entity_extractor: EntityExtractor) -> Tuple[List[Dict], List[Dict]]:
    """Extract entities and relationships from document and convert to serializable format."""
    try:
        entities, relationships = entity_extractor.extract_entities_and_relationships(doc_text, metadata)

        # Convert to dictionaries for JSON serialization
        entities_dict = [entity.model_dump() for entity in entities]
        relationships_dict = [rel.model_dump() for rel in relationships]

        return entities_dict, relationships_dict
    except Exception as e:
        logging.error(f"Failed to extract entities: {e}")
        return [], []

def analyze_documents(include_paths: List[str], fresh_start: bool, args=None, target_collection: str = None):
    logging.info(f"--- Starting Stage 2: Document Analysis with Graph Extraction (Cortex v13.0.0) ---")
    print("CORTEX_STAGE::ANALYSIS_START", flush=True)
    # Start periodic heartbeat for UI responsiveness
    start_heartbeat(period_s=5.0)
    
    # Initialize batch manager if db_path is available
    batch_manager = None
    if hasattr(args, 'db_path') and args.db_path:
        batch_manager = BatchState(args.db_path)

        # Load scan configuration if available (for resume capability)
        scan_config = {}
        scan_config_path = Path(args.db_path) / "scan_config.json"
        if scan_config_path.exists():
            try:
                with open(scan_config_path, 'r') as f:
                    scan_config = json.load(f)
                logging.info(f"Loaded scan configuration from {scan_config_path}")
            except Exception as e:
                logging.warning(f"Failed to load scan configuration: {e}")

        # Handle resume logic
        if not fresh_start:
            batch_id, files_to_process, completed_count = batch_manager.resume_or_create_batch(include_paths, scan_config)
            if not files_to_process:
                logging.info("No files to process - all files already completed")
                return

            logging.info(f"Batch {batch_id}: Processing {len(files_to_process)} files ({completed_count} already completed)")
            include_paths = files_to_process
        else:
            # Fresh start - clear any existing batch
            batch_manager.clear_batch()
            batch_manager.create_batch(include_paths, scan_config)
    
    # Get staging file path based on batch manager's db_path
    staging_file = get_staging_file_path(batch_manager.db_path) if batch_manager else STAGING_INGESTION_FILE
    if fresh_start and os.path.exists(staging_file): os.remove(staging_file)
    
    if target_collection:
        logging.info(f"Collection assignment configured: documents will be assigned to '{target_collection}'")
    
    # Initialize entity extractor
    entity_extractor = EntityExtractor()
    
    # Check memory and implement chunked processing for large batches
    if len(include_paths) > 1000:
        logging.warning(f"Large batch detected ({len(include_paths)} files). Consider processing in smaller chunks to avoid memory issues.")
    
    docs = manual_load_documents(include_paths, args)
    logging.info(f"Loaded {len(docs)} document objects from {len(include_paths)} files.")
    unique_docs = list({doc.metadata['file_path']: doc for doc in docs}.values())
    logging.info(f"--- Found {len(unique_docs)} unique files to process. ---")

    # Hard duplicate guard: skip documents already present in vector DB by doc_id hash.
    chroma_db_path = os.path.join(args.db_path, "knowledge_hub_db") if args and getattr(args, "db_path", None) else ""
    processed_log_path = os.path.join(chroma_db_path, INGESTED_FILES_LOG) if chroma_db_path else ""
    path_to_doc_id = {}
    for doc in unique_docs:
        fp = doc.metadata.get('file_path', '')
        if fp:
            path_to_doc_id[fp] = get_file_hash(fp)

    existing_doc_ids = set()
    try:
        if processed_log_path:
            existing_doc_ids.update(load_processed_files_log(processed_log_path).values())
    except Exception as e:
        logging.warning(f"Could not load processed log for duplicate guard: {e}")
    try:
        existing_doc_ids.update(_get_existing_doc_ids_in_collection(chroma_db_path, set(path_to_doc_id.values())))
    except Exception as e:
        logging.warning(f"Could not perform vector-store duplicate check: {e}")

    if existing_doc_ids:
        logging.info(f"Duplicate guard active: {len(existing_doc_ids)} known doc_ids already in DB/log")
    
    # Memory check
    try:
        import psutil
        memory_percent = psutil.virtual_memory().percent
        if memory_percent > 80:
            logging.warning(f"High memory usage detected: {memory_percent}%. Consider reducing batch size.")
    except ImportError:
        pass  # psutil not available

    staged_docs = []
    schema_json_str = json.dumps(RichMetadata.model_json_schema(), indent=2)
    prompt_lines = [
        "Analyze the document content and its file path to return a single, valid JSON object that strictly conforms to this Pydantic schema:", "```json", "{schema}", "```", "---",
        "**SPECIAL INSTRUCTIONS:**",
        '- If the source is an image (`source_type` is "image"), you **MUST** set `document_type` to "Image/Diagram".',
        '- If the file is a `.pptx`, you **MUST** set `document_type` to "Presentation".',
        '- If the `file_path` contains "trophy" or the filename mentions "Case Study", you **MUST** set `document_type` to "Case Study / Trophy".',
        '- If the filename contains the word "Final", you should strongly prefer the "Final Report" type.', '- If the filename contains the word "Draft", you should strongly prefer the "Draft Report" type.',
        '- You **MUST** classify source credibility using these exact tiers:',
        '  - 5 / peer-reviewed / Peer-Reviewed: NLM/PubMed, Nature, The Lancet, JAMA, BMJ',
        '  - 4 / institutional / Institutional: WHO, UN/IPCC, OECD, World Bank, ABS, government departments',
        '  - 3 / pre-print / Pre-Print: arXiv, SSRN, bioRxiv, ResearchGate',
        '  - 2 / editorial / Editorial: Scientific American, The Conversation, HBR',
        '  - 1 / commentary / Commentary: blogs, newsletters, consulting reports, opinion',
        '  - 0 / unclassified / Unclassified: not yet assessed (default)',
        '- If source_type indicates AI-generated content, classify credibility as tier 1 (commentary).',
        '- Populate all four fields consistently: `credibility_tier_value`, `credibility_tier_key`, `credibility_tier_label`, `credibility`.',
        '- `credibility` format must be: "{Draft|Final} {Label} Report".',
        '- If no other category seems appropriate, you **MUST** use "Other" as a fallback.', "---", "File Path: {file_path}", "Source Type: {source_type}", "Document Content (first 8000 characters):",
        "-----------------", "{text}", "-----------------", "IMPORTANT: Your response must be ONLY the JSON object itself, with no extra text, explanations, or wrapper keys."
    ]
    metadata_prompt_template = "\n".join(prompt_lines)

    # Cooldown configuration (conservative defaults)
    cooldown_every = int(getattr(args, 'cooldown_every', 25) or 0) if args else 25
    cooldown_seconds = float(getattr(args, 'cooldown_seconds', 20.0)) if args else 20.0

    for i, doc in enumerate(unique_docs):
        # Check for pause request
        if batch_manager and batch_manager.is_paused():
            logging.info("Batch processing paused by user request")
            break
            
        rich_metadata = None
        file_path_str, file_name = doc.metadata.get('file_path', ''), doc.metadata.get('file_name', 'Unknown File')
        source_type = doc.metadata.get('source_type', 'document')
        existing_doc_id = path_to_doc_id.get(file_path_str) or get_file_hash(file_path_str)

        # Skip if this exact document hash already exists in DB/log.
        if existing_doc_id in existing_doc_ids:
            logging.info(f"Skipping duplicate document (already ingested): {file_name} [{existing_doc_id[:12]}...]")
            if processed_log_path and file_path_str:
                write_to_processed_log(processed_log_path, file_path_str, existing_doc_id)
            if batch_manager:
                batch_manager.update_progress(file_path_str)
            continue

        logging.info(f"--- ({i+1}/{len(unique_docs)}) Analyzing: {file_name} ({source_type.upper()}) ---")
        
        # Print machine-readable progress for the UI
        print(f"CORTEX_PROGRESS::{i+1}/{len(unique_docs)}::{file_name}", flush=True)
        
        try:
            if not doc.text.strip(): 
                raise ValueError("Document is empty or could not be read.")
            
            # Check if Ollama LLM is available
            if Settings.llm is None:
                logging.info(f"Ollama not available - using basic metadata for {file_name}")
                # Create basic metadata when Ollama is not available
                rich_metadata = RichMetadata(
                    document_type="Other",
                    proposal_outcome="N/A", 
                    summary=f"Document processed without AI analysis (Ollama unavailable). File: {file_name}",
                    thematic_tags=["basic-processing", "no-ai-analysis"]
                )
            else:
                try:
                    prompt = metadata_prompt_template.format(schema=schema_json_str, text=doc.get_content()[:8000], file_path=file_path_str, source_type=source_type)
                    logging.info("Sending prompt to LLM for metadata extraction...")

                    # Add timeout handling with progress feedback
                    import time
                    start_time = time.time()
                    # Enforce a hard timeout around the LLM call to avoid indefinite stalls
                    llm_timeout = float(getattr(args, 'llm_timeout', 120.0)) if args else 120.0
                    try:
                        response_str = str(_run_with_timeout(lambda: Settings.llm.complete(prompt), llm_timeout))
                        elapsed = time.time() - start_time
                        if elapsed > 60:
                            logging.info(f"LLM call took {elapsed:.1f}s (longer document)")
                    except TimeoutError as te:
                        logging.warning(f"LLM timeout after {time.time() - start_time:.1f}s for {file_name}, using fallback metadata")
                        # Brief backoff to let Ollama recover a bit before the next document
                        try:
                            import time as _t
                            _t.sleep(5)
                        except Exception:
                            pass
                        raise ValueError(f"LLM timeout: {te}")
                    except Exception as llm_ex:
                        logging.warning(f"LLM error after {time.time() - start_time:.1f}s for {file_name}: {llm_ex}")
                        raise

                    logging.info("Received raw response from LLM. Cleaning and parsing JSON...")
                    
                    # Check if response is empty (indicates LLM error)
                    if not response_str or response_str.strip() == "":
                        raise ValueError("LLM returned empty response")
                    
                    json_match = re.search(r'\{.*\}', response_str, re.DOTALL)
                    if not json_match:
                        error_snippet = response_str.strip().replace('\n', ' ')[:200]
                        logging.error(f"LLM did not return valid JSON for {file_name}. Response: {error_snippet}...")
                        raise ValueError("LLM did not return a valid JSON object.")
                    
                    clean_json_str = json_match.group(0)
                    metadata_json = json.loads(clean_json_str)
                    
                    # Only process if we got valid JSON data
                    if 'metadata_json' in locals():
                        # Normalize common validation fields before Pydantic
                        if 'proposal_outcome' in metadata_json and isinstance(metadata_json['proposal_outcome'], str):
                            metadata_json['proposal_outcome'] = metadata_json['proposal_outcome'].title()
                        _enforce_credibility_policy(metadata_json, file_path_str, doc.get_content()[:8000])

                        try: 
                            rich_metadata = RichMetadata.model_validate(metadata_json)
                        except ValidationError as e:
                            logging.warning(f"Initial validation failed for {file_name}. Retrying with nested key check...")
                            if isinstance(metadata_json, dict) and len(metadata_json) == 1:
                                nested_key = list(metadata_json.keys())[0]
                                if isinstance(metadata_json[nested_key], dict): 
                                    rich_metadata = RichMetadata.model_validate(metadata_json[nested_key])
                                else: 
                                    raise e
                            else:
                                raise e
                                
                except (ValueError, json.JSONDecodeError, RuntimeError, Exception) as llm_error:
                    logging.warning(f"LLM processing failed for {file_name}: {llm_error}")
                    logging.info(f"Falling back to basic metadata for {file_name}")
                    # Create fallback metadata when LLM fails
                    rich_metadata = RichMetadata(
                        document_type="Other",
                        proposal_outcome="N/A", 
                        summary=f"Document processed with fallback metadata (LLM error). File: {file_name}",
                        thematic_tags=["llm-error", "fallback-metadata"]
                    )
            logging.info(f"Successfully parsed and validated metadata for {file_name}.")
            
            # Extract entities and relationships
            logging.info(f"Extracting entities and relationships from {file_name}...")
            if rich_metadata is not None:
                entities_dict, relationships_dict = extract_entities_and_relationships(
                    doc.get_content()[:8000],
                    {
                        'document_type': rich_metadata.document_type,
                        'file_name': file_name,
                        'summary': rich_metadata.summary,
                        'thematic_tags': rich_metadata.thematic_tags
                    },
                    entity_extractor
                )
            else:
                logging.warning(f"Skipping entity extraction for {file_name} - metadata is None")
                entities_dict, relationships_dict = {}, {}
            logging.info(f"Extracted {len(entities_dict)} entities and {len(relationships_dict)} relationships.")
            
        except Exception as e:
            logging.error(f"CRITICAL ERROR analyzing {file_name}: {e}", exc_info=True)
            # Record error in batch manager
            if batch_manager:
                batch_manager.record_error(file_path_str, str(e))
            
            default_doc_type = "Image/Diagram" if source_type == 'image' else 'Other'
            rich_metadata = RichMetadata(
                document_type=default_doc_type, 
                proposal_outcome="N/A", 
                summary=f"ERROR: Could not analyze document. Reason: {e}", 
                thematic_tags=["error", "analysis-failed"]
            )
            entities_dict = []
            relationships_dict = []

        doc_meta = DocumentMetadata(
            doc_id=existing_doc_id,
            doc_posix_path=Path(file_path_str).as_posix(),
            file_name=doc.metadata.get('file_name'), 
            last_modified_date=str(datetime.fromtimestamp(os.path.getmtime(file_path_str))),
            rich_metadata=rich_metadata,
            extracted_entities=entities_dict,
            extracted_relationships=relationships_dict
        )
        staged_docs.append(doc_meta.model_dump())
        
        # Update batch progress for successfully processed file
        if batch_manager:
            batch_manager.update_progress(file_path_str)

        # Smart adaptive throttling to reduce CPU/GPU load
        base_throttle_delay = getattr(args, 'throttle_delay', 1.0) if args else 1.0

        # Always check system load, even if base delay is 0 (enables auto-throttling)
        if base_throttle_delay >= 0:  # -1 would disable entirely (not exposed in UI)
            import time

            # Always sample system load to provide responsive feedback
            gpu_util = get_gpu_utilization()
            cpu_util = get_cpu_utilization()

            # Calculate adaptive delay (more conservative defaults)
            actual_delay, throttle_reason = calculate_adaptive_throttle(
                base_throttle_delay,
                gpu_util,
                cpu_util,
                gpu_threshold=float(getattr(args, 'gpu_threshold', 60.0)) if args else 60.0,
                cpu_threshold=float(getattr(args, 'cpu_threshold', 70.0)) if args else 70.0,
                increment=0.5,
                max_delay=float(getattr(args, 'max_throttle_delay', 8.0)) if args else 8.0
            )

            # Always print machine-readable throttle status for UI (even when not throttling)
            gpu_str = f"{gpu_util:.0f}" if gpu_util is not None else "N/A"
            cpu_str = f"{cpu_util:.0f}" if cpu_util is not None else "N/A"

            if throttle_reason:
                # High load - throttling active
                logging.info(f"⏱️ {throttle_reason}")
                print(f"CORTEX_THROTTLE::{actual_delay:.1f}::{gpu_str}::{cpu_str}", flush=True)
            else:
                # Normal load - baseline delay for responsiveness
                print(f"CORTEX_THROTTLE::{actual_delay:.1f}::{gpu_str}::{cpu_str}", flush=True)

            # Apply delay (ensures system stays responsive)
            if actual_delay > 0:
                time.sleep(actual_delay)

        # Periodic cooldown to prevent thermal throttling on laptops
        if cooldown_every and (i + 1) % cooldown_every == 0 and cooldown_seconds > 0:
            logging.info(f"🌡️ Cooldown: sleeping {cooldown_seconds:.1f}s after {(i + 1)} documents")
            print(f"CORTEX_COOLDOWN::{cooldown_seconds:.1f}::{i+1}", flush=True)
            import time as _t
            _t.sleep(cooldown_seconds)

    # Create staging data structure with collection assignment
    staging_data = {
        "documents": staged_docs,  # staged_docs already contains dicts from doc_meta.model_dump()
        "target_collection": target_collection,
        "created_at": str(datetime.now()),
        "version": "2.0"  # New format with collection assignment
    }
    
    with open(staging_file, 'w') as f: 
        json.dump(staging_data, f, indent=2)
    
    collection_info = f" -> '{target_collection}'" if target_collection else " -> 'default'"
    logging.info(f"--- Analysis complete. {len(staged_docs)} documents staged{collection_info} at {staging_file} ---")
    print(f"CORTEX_STAGE::ANALYSIS_DONE::{len(staged_docs)}::{staging_file}", flush=True)
    # Stop heartbeat
    stop_heartbeat()

def finalize_ingestion(db_path: str, args=None):
    logging.info(f"--- Starting Stage 3: Finalize from Staging with Graph Building (Cortex v13.0.0) ---")
    print("CORTEX_STAGE::FINALIZE_START", flush=True)
    staging_file = get_staging_file_path(db_path)
    if not os.path.exists(staging_file): 
        logging.error(f"Staging file not found at: {staging_file}")
        return
    
    batch_manager = BatchState(db_path)
    chroma_db_path = os.path.join(db_path, "knowledge_hub_db")
    _ensure_directory_cross_platform(chroma_db_path)
    
    # Initialize graph manager
    graph_file_path = os.path.join(db_path, "knowledge_cortex.gpickle")
    graph_manager = EnhancedGraphManager(graph_file_path)
    
    with open(staging_file, 'r') as f: 
        staging_data = json.load(f)
    
    # Handle both old format (list) and new format (dict with metadata)
    if isinstance(staging_data, list):
        docs_to_process = [DocumentMetadata(**data) for data in staging_data]
        target_collection = None  # No collection assignment for old format
    else:
        docs_to_process = [DocumentMetadata(**data) for data in staging_data.get('documents', [])]
        target_collection = staging_data.get('target_collection', None)
    
    logging.info(f"Target collection for finalization: {target_collection or 'default'}")

    # Finalize-stage duplicate guard in case staged data predates checks above.
    staged_doc_ids = {d.doc_id for d in docs_to_process if d and d.doc_id}
    already_in_db_doc_ids = _get_existing_doc_ids_in_collection(chroma_db_path, staged_doc_ids)
    if already_in_db_doc_ids:
        logging.info(f"Finalize duplicate guard: {len(already_in_db_doc_ids)} staged documents already exist in DB")

    docs_to_index_paths, metadata_map, doc_ids_to_add_to_default = [], {}, []
    processed_log_path = os.path.join(chroma_db_path, INGESTED_FILES_LOG)
    
    for doc_meta in docs_to_process:
        if doc_meta.doc_id in already_in_db_doc_ids:
            logging.info(f"Skipping duplicate during finalize: {doc_meta.file_name} [{doc_meta.doc_id[:12]}...]")
            write_to_processed_log(processed_log_path, doc_meta.doc_posix_path, doc_meta.doc_id)
            continue
        if doc_meta.exclude_from_final:
            logging.info(f"User excluded {doc_meta.file_name}. Skipping.")
            write_to_processed_log(processed_log_path, doc_meta.doc_posix_path, doc_meta.doc_id)
            continue
        if doc_meta.rich_metadata and "ERROR:" in doc_meta.rich_metadata.summary:
            logging.warning(f"Skipping finalization for {doc_meta.file_name} due to prior analysis error.")
            continue
        
        docs_to_index_paths.append(doc_meta.doc_posix_path)
        metadata_map[doc_meta.doc_posix_path] = doc_meta
        doc_ids_to_add_to_default.append(doc_meta.doc_id)
        
        # Add document to graph
        logging.info(f"Adding {doc_meta.file_name} to knowledge graph...")
        graph_manager.add_entity(
            doc_meta.doc_id,
            'Document',
            file_name=doc_meta.file_name,
            document_type=doc_meta.rich_metadata.document_type if doc_meta.rich_metadata else 'Unknown',
            summary=doc_meta.rich_metadata.summary if doc_meta.rich_metadata else '',
            last_modified=doc_meta.last_modified_date,
            posix_path=doc_meta.doc_posix_path
        )
        
        # Add entities and relationships to graph
        entity_map = {}  # Track entity names to node IDs
        
        for entity_dict in doc_meta.extracted_entities:
            entity_id = f"{entity_dict['entity_type']}:{entity_dict['name']}"
            entity_map[entity_dict['name']] = entity_id
            
            # Add entity if it doesn't exist
            if entity_id not in graph_manager.graph:
                graph_manager.add_entity(
                    entity_id,
                    entity_dict['entity_type'],
                    name=entity_dict['name'],
                    aliases=entity_dict.get('aliases', [])
                )
            
            # Link entity to document
            if entity_dict['entity_type'] == 'person':
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'authored'
                )
            elif entity_dict['entity_type'] == 'organization':
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'client_of'
                )
            else:
                graph_manager.add_relationship(
                    entity_id,
                    doc_meta.doc_id,
                    'mentioned_in'
                )
        
        # Add extracted relationships
        for rel_dict in doc_meta.extracted_relationships:
            source_id = entity_map.get(rel_dict['source'], rel_dict['source'])
            target_id = entity_map.get(rel_dict['target'], rel_dict['target'])

            # Handle special case where target might be the document
            if rel_dict['target'] == doc_meta.file_name:
                target_id = doc_meta.doc_id

            graph_manager.add_relationship(
                source_id,
                target_id,
                rel_dict['relationship_type'],
                context=rel_dict.get('context', '')
            )

        # STABILITY: Mandatory delay between document graph processing
        # Prevents runaway resource consumption during intensive operations
        import time
        import gc
        gc.collect()
        time.sleep(0.5)  # Half-second delay for graph operations (lighter than embedding)

    if not docs_to_index_paths:
        logging.warning("No new, valid documents to ingest. Finalization complete.")
        if os.path.exists(staging_file): 
            os.remove(staging_file)
        try:
            batch_manager.clear_batch()
        except Exception as e:
            logging.warning(f"Failed to clear batch state after empty finalization: {e}")
        return

    # Save the graph
    graph_manager.save_graph()
    logging.info(f"Knowledge graph saved to {graph_file_path}")

    # Clear query cache since new data was added
    try:
        from cortex_engine.graph_query import clear_query_cache
        clear_query_cache()
        logging.info("🔄 Query cache cleared due to new ingestion")
    except Exception as e:
        logging.warning(f"Failed to clear query cache: {e}")

    # Continue with regular vector indexing
    documents_for_indexing = manual_load_documents(docs_to_index_paths, args)
    for doc in documents_for_indexing:
        path_key = doc.metadata['file_path']
        if path_key in metadata_map:
            doc_meta = metadata_map[path_key]
            # IMPORTANT: LlamaIndex uses `id_` as the document identity that flows
            # into node metadata. Using `doc_id` here does not persist.
            doc.id_ = doc_meta.doc_id
            flat_metadata = {
                "doc_id": doc_meta.doc_id,
                "file_name": doc_meta.file_name,
                "doc_posix_path": doc_meta.doc_posix_path,
                "last_modified_date": doc_meta.last_modified_date
            }
            if doc_meta.rich_metadata:
                flat_metadata.update(doc_meta.rich_metadata.model_dump())
                flat_metadata['thematic_tags'] = ', '.join(flat_metadata.get('thematic_tags', []))

            # Phase 1 Enhancement: Preserve complex Docling metadata
            # ChromaDB requires flat metadata, so serialize complex structures as JSON
            if 'docling_provenance' in doc.metadata:
                flat_metadata['docling_provenance'] = json.dumps(doc.metadata['docling_provenance'])
            if 'docling_figures' in doc.metadata:
                flat_metadata['docling_figures'] = json.dumps(doc.metadata['docling_figures'])
            if 'docling_structure' in doc.metadata:
                flat_metadata['docling_structure'] = json.dumps(doc.metadata['docling_structure'])

            doc.metadata = flat_metadata

    # Phase 2 Enhancement: Apply table-aware chunking and figure entity linking
    if PHASE2_AVAILABLE and (TABLE_AWARE_CHUNKING or FIGURE_ENTITY_LINKING):
        logging.info("🚀 Applying Phase 2 enhancements to documents...")

        enhanced_documents = []

        # Apply figure entity linking first (operates on full documents)
        if FIGURE_ENTITY_LINKING:
            try:
                # Load knowledge graph for entity linking
                graph_path = args.graph_path if args and hasattr(args, 'graph_path') else None
                if not graph_path:
                    from cortex_engine.config import GRAPH_FILE_PATH
                    graph_path = GRAPH_FILE_PATH

                knowledge_graph = load_knowledge_graph_for_linking(graph_path)
                entity_linker = create_figure_entity_linker(knowledge_graph)

                logging.info(f"📎 Linking figures to knowledge graph entities...")
                documents_for_indexing = entity_linker.link_batch_documents(
                    documents_for_indexing,
                    knowledge_graph
                )

                # Serialize figure_entities metadata for ChromaDB
                for doc in documents_for_indexing:
                    if 'figure_entities' in doc.metadata:
                        doc.metadata['figure_entities'] = json.dumps(doc.metadata['figure_entities'])

                logging.info(f"✅ Figure entity linking complete")

            except Exception as e:
                logging.warning(f"Figure entity linking failed (non-critical): {e}")

        # Apply table-aware chunking (expands documents into chunks)
        if TABLE_AWARE_CHUNKING:
            try:
                chunker = create_table_aware_chunker(
                    chunk_size=1024,
                    chunk_overlap=200,
                    table_context_sentences=2
                )

                logging.info(f"📊 Applying table-aware chunking to {len(documents_for_indexing)} documents...")

                for doc in documents_for_indexing:
                    # Chunk document while preserving tables
                    chunks = chunker.process_document(doc)
                    enhanced_documents.extend(chunks)

                original_count = len(documents_for_indexing)
                documents_for_indexing = enhanced_documents

                logging.info(
                    f"✅ Table-aware chunking complete: "
                    f"{original_count} docs → {len(documents_for_indexing)} chunks"
                )

            except Exception as e:
                logging.warning(f"Table-aware chunking failed (non-critical), using original documents: {e}")
                # Keep original documents if chunking fails

    # Validate embedding model compatibility before ingestion
    # CRITICAL: This prevents database corruption from dimension mismatches
    from cortex_engine.utils.embedding_validator import (
        validate_embedding_compatibility,
        get_embedding_dimension,
        get_database_embedding_dimension,
        EmbeddingModelMismatchError
    )
    from cortex_engine.collection_manager import WorkingCollectionManager

    # Get adaptive embedding model
    current_embed_model = get_embed_model()
    current_model_dimension = get_embedding_dimension(current_embed_model)
    logging.info(f"Current embedding model: {current_embed_model} ({current_model_dimension}D)")

    # CRITICAL: Check actual database dimension directly from ChromaDB
    # This catches dimension mismatches even for imported databases without metadata
    db_dimension = get_database_embedding_dimension(db_path)

    if db_dimension is not None:
        # Existing database - validate dimension match
        if db_dimension != current_model_dimension:
            error_msg = (
                f"DIMENSION MISMATCH DETECTED!\n"
                f"   Database contains: {db_dimension}-dimensional embeddings\n"
                f"   Current model produces: {current_model_dimension}-dimensional embeddings\n"
                f"   Model: {current_embed_model}\n"
                f"\n"
                f"   Using this model would CORRUPT the database!\n"
                f"\n"
                f"   SOLUTIONS:\n"
                f"   1. Select a model with {db_dimension}D embeddings (e.g., "
                f"{'Qwen3-VL-2B' if db_dimension == 2048 else 'Qwen3-VL-8B' if db_dimension == 4096 else 'matching model'})\n"
                f"   2. Or delete the database and start fresh with the current model\n"
            )
            logging.error(f"❌ CRITICAL: {error_msg}")
            print(f"CORTEX_ERROR::DIMENSION_MISMATCH::{db_dimension}::{current_model_dimension}", flush=True)
            raise RuntimeError(f"Embedding dimension mismatch prevents ingestion: {db_dimension}D database vs {current_model_dimension}D model")
        else:
            logging.info(f"✅ Embedding dimension validated: {db_dimension}D (database matches model)")
    else:
        # New database - any dimension is fine
        logging.info(f"Creating new database with {current_model_dimension}D embeddings")

    # Check if collection already exists and validate compatibility
    db_settings = ChromaSettings(anonymized_telemetry=False)
    chroma_client = chromadb.PersistentClient(path=chroma_db_path, settings=db_settings)

    # Get or create collection with embedding model metadata
    try:
        existing_collections = chroma_client.list_collections()
        collection_exists = any(c.name == COLLECTION_NAME for c in existing_collections)

        if collection_exists:
            # Validate existing collection metadata (secondary check)
            collection_mgr_temp = WorkingCollectionManager()
            collection_metadata = collection_mgr_temp.get_embedding_model_metadata("default")

            # Only validate metadata if it exists (imported DBs may not have it)
            if collection_metadata.get("embedding_model") or collection_metadata.get("embedding_dimension"):
                try:
                    validation_result = validate_embedding_compatibility(
                        collection_metadata,
                        current_model=current_embed_model,
                        strict=True
                    )
                    logging.info(f"✅ Embedding model metadata validation passed")
                except EmbeddingModelMismatchError as e:
                    logging.error(f"❌ CRITICAL: {e}")
                    logging.error(f"   Current model: {e.current_model}")
                    logging.error(f"   Expected model: {e.expected_model}")
                    raise RuntimeError(f"Embedding model mismatch prevents ingestion: {e}")
            else:
                logging.info("No embedding metadata found - relying on dimension check (passed above)")
        else:
            # New collection - store embedding model metadata
            logging.info(f"Creating new collection with embedding model: {current_embed_model}")
            logging.info(f"Embedding dimension: {current_model_dimension}")

    except EmbeddingModelMismatchError:
        raise  # Re-raise validation errors
    except RuntimeError:
        raise  # Re-raise dimension mismatch errors
    except Exception as e:
        logging.warning(f"Could not validate embedding model metadata (dimension check passed): {e}")

    chroma_collection = chroma_client.get_or_create_collection(COLLECTION_NAME)
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    storage_context = StorageContext.from_defaults(vector_store=vector_store)

    # Emit coarse-grained progress during finalization by indexing in small batches
    total_docs = len(documents_for_indexing)
    batch_size = 25 if total_docs > 50 else max(5, total_docs // 4 or 1)
    indexed = 0
    logging.info(f"Indexing {total_docs} documents in batches of {batch_size}...")
    # Optional cooldown between index batches
    index_batch_cooldown = float(getattr(args, 'index_batch_cooldown', 1.0)) if args else 1.0

    for start in range(0, total_docs, batch_size):
        batch_docs = documents_for_indexing[start:start+batch_size]
        try:
            VectorStoreIndex.from_documents(batch_docs, storage_context=storage_context, show_progress=True)
        except Exception as e:
            logging.error(f"Batch indexing failed at {start}-{start+len(batch_docs)}: {e}", exc_info=True)
            # Continue with next batch to be resilient
        indexed += len(batch_docs)
        # Emit machine-readable progress update for UI
        last_name = batch_docs[-1].metadata.get('file_name', 'batch') if batch_docs else 'batch'
        print(f"CORTEX_PROGRESS::{indexed}/{total_docs}::{last_name}", flush=True)

        # STABILITY: Mandatory cleanup and pause between index batches
        # Prevents memory buildup and GPU exhaustion during embedding
        import gc
        import torch
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()

        # Mandatory delay (minimum 1 second for stability)
        import time as _t
        actual_cooldown = max(1.0, index_batch_cooldown)
        _t.sleep(actual_cooldown)

    logging.info(f"Persisting index to disk at {chroma_db_path}...")
    storage_context.persist(persist_dir=chroma_db_path)
    
    for doc in documents_for_indexing: 
        write_to_processed_log(processed_log_path, doc.metadata['doc_posix_path'], doc.metadata['doc_id'])
    os.remove(staging_file)

    if doc_ids_to_add_to_default:
        # Use target collection if specified, otherwise default
        collection_name = target_collection or "default"
        logging.info(f"Adding {len(doc_ids_to_add_to_default)} new documents to the '{collection_name}' collection.")
        try:
            from cortex_engine.collection_manager import WorkingCollectionManager
            from cortex_engine.utils.embedding_validator import get_embedding_dimension

            # Get adaptive embedding model
            current_embed_model = get_embed_model()

            collection_mgr = WorkingCollectionManager()
            collection_mgr.add_docs_by_id_to_collection(collection_name, doc_ids_to_add_to_default)

            # Store/update embedding model metadata for this collection
            try:
                embedding_dimension = get_embedding_dimension(current_embed_model)
                collection_mgr.set_embedding_model_metadata(
                    collection_name,
                    current_embed_model,
                    embedding_dimension
                )
                logging.info(f"Stored embedding model metadata: {current_embed_model} ({embedding_dimension}D)")
            except Exception as meta_error:
                logging.warning(f"Could not store embedding metadata (non-critical): {meta_error}")

            logging.info("Collections updated via WorkingCollectionManager (Docker)")
        except Exception as e:
            logging.error(f"Could not automatically add documents to '{collection_name}' collection: {e}")
    
    logging.info("--- Finalization complete. Knowledge base and graph are up to date. ---")
    print("CORTEX_STAGE::FINALIZE_DONE", flush=True)
    try:
        batch_manager.clear_batch()
    except Exception as e:
        logging.warning(f"Failed to clear batch state after finalization: {e}")

def main():
    parser = argparse.ArgumentParser(description="Cortex Ingestion Engine with GraphRAG")
    parser.add_argument("--db-path", type=str, required=True)
    parser.add_argument("--include", type=str, nargs='*')
    parser.add_argument("--analyze-only", action="store_true")
    parser.add_argument("--finalize-from-staging", action="store_true")
    parser.add_argument("--fresh", action="store_true")
    parser.add_argument("--skip-image-processing", action="store_true", 
                       help="Skip VLM image processing for faster ingestion")
    parser.add_argument("--resume", action="store_true",
                       help="Resume from existing batch state (opposite of --fresh)")
    parser.add_argument("--target-collection", type=str, default=None,
                       help="Target collection for document assignment")
    parser.add_argument("--throttle-delay", type=float, default=1.0,
                       help="Baseline delay (s) between documents. Auto-adjusts with CPU/GPU load.")
    parser.add_argument("--cpu-threshold", type=float, default=70.0,
                       help="CPU utilization (%) to start increasing delay (default 70)")
    parser.add_argument("--gpu-threshold", type=float, default=60.0,
                       help="GPU utilization (%) to start increasing delay (default 60)")
    parser.add_argument("--max-throttle-delay", type=float, default=8.0,
                       help="Maximum adaptive delay (seconds) when system is under load")
    parser.add_argument("--cooldown-every", type=int, default=25,
                       help="After this many documents, sleep for cooldown period")
    parser.add_argument("--cooldown-seconds", type=float, default=20.0,
                       help="Cooldown duration (seconds) to prevent thermal throttling")
    parser.add_argument("--image-workers", type=int, default=1,
                       help="Max concurrent image VLM workers (default 1, llava crashes on concurrent)")
    parser.add_argument("--index-batch-cooldown", type=float, default=1.0,
                       help="Sleep (s) between indexing batches to reduce sustained load")
    parser.add_argument("--llm-timeout", type=float, default=120.0,
                       help="Hard timeout (seconds) per LLM metadata call; on timeout, fallback metadata is used")
    parser.add_argument("--gpu-intensity", type=int, default=75,
                       help="GPU intensity 25-100%%. Lower = smaller batches + longer delays. Use 50-75%% if multitasking.")
    args = parser.parse_args()

    # Store GPU intensity in environment for embedding services to use
    import os
    os.environ["CORTEX_GPU_INTENSITY"] = str(args.gpu_intensity)

    writable, reason = ensure_directory_writable(args.db_path)
    if not writable:
        logging.error(f"Database path '{args.db_path}' is not writable: {reason}")
        sys.exit(1)
    logging.info(f"Knowledge base path resolved to: {args.db_path}")

    initialize_script()

    try:
        if args.analyze_only:
            if not args.include: 
                logging.error("--include paths are required.")
                sys.exit(1)
            # Resume takes precedence over fresh
            fresh_start = args.fresh and not args.resume
            analyze_documents(args.include, fresh_start, args, getattr(args, 'target_collection', None))
        elif args.finalize_from_staging:
            finalize_ingestion(args.db_path, args)
        else:
            logging.error("Specify either --analyze-only or --finalize-from-staging.")
            sys.exit(1)
    except KeyboardInterrupt:
        logging.info("Process interrupted by user (Ctrl+C)")
        sys.exit(0)
    except Exception as e:
        logging.error(f"CRITICAL ERROR: Batch processing failed with exception: {e}", exc_info=True)
        # Try to save any progress made so far
        try:
            if hasattr(args, 'db_path') and args.db_path:
                batch_manager = BatchState(args.db_path)
                batch_manager.record_error("SYSTEM_CRASH", f"Critical system error: {e}")
        except Exception as save_error:
            logging.error(f"Failed to save error state: {save_error}")  # Log but don't re-raise
        sys.exit(1)

if __name__ == "__main__":
    main()
